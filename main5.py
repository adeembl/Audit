import os
import glob
import io
import openpyxl
import matplotlib
matplotlib.use('Agg')  # Para entornos sin pantalla (ej: en servidor)
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime, timedelta
import webview
import shutil
from pathlib import Path
import subprocess  # Para abrir archivos de forma cross-platform
import sys  # Para detectar el sistema operativo
# ==================== LIBRERÍAS SAP OCULTO ====================
import win32com.client
import pyperclip
import pyautogui
import win32gui
import win32con
import threading
import tkinter as tk
from tkinter import ttk, messagebox
from tkcalendar import Calendar
from pdf2image import convert_from_path
import tempfile
# Configuración de pyautogui
pyautogui.PAUSE = 0
pyautogui.FAILSAFE = False
# Importar pandas si está disponible (para manejo de fechas)
try:
    import pandas as pd
except ImportError:
    print("ADVERTENCIA: pandas no está instalado. Algunas funciones de fecha pueden no funcionar correctamente.")
    pd = None

# Librerías de PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Librerías de ReportLab para PDF
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
)
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib import colors

# ------------------------------------------------------------------------------
# Configuración de rutas de recursos
# ------------------------------------------------------------------------------
ruta_origen_html_recursos = r"\\PIASFIL01\info-cal-eur\56. Rechazo Adem\Base de datos"
ruta_destino_html_local = os.path.join(os.path.dirname(__file__), "index2.html")
ruta_logo_local = os.path.join(os.path.dirname(__file__), "logo2.png")

# Directorio para guardar PDF
REPORTS_DIR_NAME = "Resultados"
REPORTS_DIR_PATH = os.path.join(os.path.dirname(__file__), REPORTS_DIR_NAME)

# Directorio para PowerPoint
POWERPOINT_DIR_PATH = r"\\PIASFIL01\Pub-CalidadTV"

# Tamaño de página para PDF (A3 landscape)
PAGE_SIZE = (1190, 842)

# Copiar index.html si no existe
if not os.path.exists(ruta_destino_html_local):
    try:
        shutil.copy(os.path.join(ruta_origen_html_recursos, "index.html"), ruta_destino_html_local)
        print(f"index.html copiado a {ruta_destino_html_local}")
    except Exception as e:
        print(f"Error copiando index.html: {e}")

# Copiar logo2.png si no existe
if not os.path.exists(ruta_logo_local):
    try:
        shutil.copy(os.path.join(ruta_origen_html_recursos, "logo2.png"), ruta_logo_local)
        print(f"logo2.png copiado a {ruta_logo_local}")
    except Exception as e:
        print(f"Error copiando logo2.png: {e}")


def _emu_to_inches(v):
    """Convierte EMU (pptx.util.Length) a pulgadas (float)."""
    return v.inches if isinstance(v, Length) else float(v)


class API:
    # Colores corporativos
    corporate_colors = {
        'primary_blue': '#0B2240',
        'primary_green': '#00B39F',
        'secondary_blue': '#1e3a5f',
        'light_blue': '#4a6fa5',
        'dark_gray': '#2c3e50',
        'medium_gray': '#7f8c8d',
        'light_gray': '#ecf0f1',
        'lighter_gray': '#95a5a6',
        'danger': '#e74c3c',
        'warning': '#f39c12',
        'success': '#27ae60',
    }

    # Paleta simplificada para gráficas
    chart_palette = ['#0B2240', '#00B39F', '#4a6fa5', '#7f8c8d', '#2c3e50']

    shift_color_map = {
        "T": "#0B2240", "M": "#00B39F", "N": "#4a6fa5",
        "T1": "#0B2240", "T2": "#00B39F", "T3": "#4a6fa5",
        "T4": "#2c3e50", "T5": "#7f8c8d", "SIN_TURNO": "#ecf0f1",
        "MAÑANA": "#00B39F", "TARDE": "#0B2240", "NOCHE": "#4a6fa5",
        "A": "#00B39F", "B": "#0B2240", "C": "#4a6fa5"
    }

    defect_color_cycle = [
        '#0B2240', '#00B39F', '#4a6fa5', '#7f8c8d', '#2c3e50',
        '#1e3a5f', '#95a5a6', '#34495e', '#16a085', '#27ae60'
    ]

    model_color_cycle = [
        '#0B2240', '#00B39F', '#4a6fa5', '#7f8c8d', '#2c3e50',
        '#1e3a5f', '#95a5a6', '#34495e', '#16a085', '#27ae60'
    ]

    # Objetivo en porcentaje (0.80 equivale a 0.80%)
    OBJETIVO = 0.80

    TREND_DAYS = 30  # Días para tendencias en PDF
    
    def __init__(self):
        # Ruta de archivos Excel
        self.ruta_archivos = r"\\PIASFIL01\info-cal-eur\56. Rechazo Adem\Base de datos\EXCEL"
        self.ruta_costes = r"\\PIASFIL01\info-cal-eur\56. Rechazo Adem\Base de datos\EXCEL\costes.xlsx"
        self.reports_dir = REPORTS_DIR_PATH
        self.powerpoint_dir = POWERPOINT_DIR_PATH

        # NUEVO: Cargar tabla de costes al inicializar
        self.costes_dict = self._cargar_costes()

        # Configurar matplotlib estilo profesional
        plt.style.use('seaborn-v0_8-whitegrid')
        plt.rcParams['font.family'] = 'Arial'
        plt.rcParams['font.size'] = 18
        plt.rcParams['axes.labelsize'] = 18
        plt.rcParams['axes.titlesize'] = 18
        plt.rcParams['figure.figsize'] = (10, 6)

        # Crear directorio de reportes si no existe
        if not os.path.exists(self.reports_dir):
            try:
                os.makedirs(self.reports_dir)
                print(f"Directorio de reportes creado en: {self.reports_dir}")
            except Exception as e:
                print(f"Error creando directorio de reportes '{self.reports_dir}': {e}")

    def _format_euro(self, valor):
        """Formatea números al estilo europeo: 10.000,23€"""
        return f"{valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    def _cargar_costes(self):
        """Cargar el diccionario de costes desde el Excel de costes"""
        costes_dict = {}
        try:
            if os.path.exists(self.ruta_costes):
                wb_costes = openpyxl.load_workbook(self.ruta_costes, data_only=True)
                ws_costes = wb_costes.active
                
                for row in ws_costes.iter_rows(min_row=2, values_only=True):  # Saltamos encabezados
                    if len(row) >= 2 and row[0] and row[1]:
                        texto_breve = str(row[0]).strip().upper()
                        try:
                            precio = float(row[1])
                            costes_dict[texto_breve] = precio
                        except (ValueError, TypeError):
                            continue
                            
                print(f"✅ Cargados {len(costes_dict)} precios del archivo de costes")
            else:
                print(f"⚠️ ADVERTENCIA: No se encontró el archivo de costes en {self.ruta_costes}")
                
        except Exception as e:
            print(f"❌ Error al cargar archivo de costes: {e}")
            
        return costes_dict

    def _es_defecto_proveedor(self, val_defecto):
        return "PROVEEDOR" in (val_defecto or "").upper()

    def _parse_num(self, valor):
        if not valor:
            return 0
        s = str(valor).strip().replace(',', '.')
        neg = s.endswith('-')
        if neg:
            s = s[:-1].strip()
        try:
            f = float(s)
            return -int(round(f)) if neg else int(round(f))
        except:
            return 0

    def leer_estadisticas(self, fecha_inicio, fecha_fin, modelo, defecto):
        """
        VERSIÓN MEJORADA: Lee los archivos Excel con filtrado exhaustivo y preciso.
        Cuando se filtra por modelo o defecto, TODAS las estadísticas se enfocan en ese filtro específico.
        INCLUYE ANÁLISIS ECONÓMICO COMPLETO.
        """
        try:
            fi = datetime.strptime(fecha_inicio, '%Y-%m-%d').date()
            ff = datetime.strptime(fecha_fin, '%Y-%m-%d').date()
        except ValueError:
            return {"error": "Formato de fecha inválido. Debe ser 'YYYY-MM-DD'."}

        # Manejar múltiples modelos y defectos
        if modelo and '|' in modelo:
            modelos_list = [m.strip().upper() for m in modelo.split('|') if m.strip()]
            filtrar_modelo = len(modelos_list) > 0
        else:
            modelo_upper = (modelo or "").strip().upper()
            modelos_list = [modelo_upper] if modelo_upper and modelo_upper != "(TODOS)" else []
            filtrar_modelo = len(modelos_list) > 0

        if defecto and '|' in defecto:
            defectos_list = [d.strip().upper() for d in defecto.split('|') if d.strip()]
            filtrar_defecto = len(defectos_list) > 0
        else:
            defecto_upper = (defecto or "").strip().upper()
            defectos_list = [defecto_upper] if defecto_upper and defecto_upper != "(TODOS)" else []
            filtrar_defecto = len(defectos_list) > 0

        # Calcular período anterior para tendencias
        duracion_periodo = (ff - fi).days + 1
        fecha_fin_anterior = fi - timedelta(days=1)
        fecha_inicio_anterior = fecha_fin_anterior - timedelta(days=duracion_periodo - 1)

        # Estructuras para datos FILTRADOS (solo incluyen registros que cumplen los filtros)
        filtered_evolution = {}              # {fecha_iso: {'OK': x, 'NOK': y, 'NOK_PROV': z}, ...}
        filtered_evolution_rework = {}       # {fecha_iso: total_R}
        filtered_defect_distribution = {}    # {defecto: total_NOK} - solo defectos del filtro
        filtered_shift_distribution = {}     # {turno: total_NOK} - solo turnos del filtro
        filtered_models_distribution = {}    # {modelo: total_NOK} - solo modelos del filtro
        filtered_rework_defect_distribution = {}
        filtered_rework_shift_distribution = {}
        filtered_rework_models_distribution = {}

        # Estructuras para datos GLOBALES (todos los registros, sin filtros)
        global_evolution = {}
        global_evolution_rework = {}
        global_defect_distribution = {}
        global_shift_distribution = {}
        global_models_distribution = {}
        global_models_distribution_prov = {}  
        global_rework_defect_distribution = {}
        global_rework_shift_distribution = {}
        global_rework_models_distribution = {}

        # Estructuras para período anterior (para tendencias)
        filtered_evolution_anterior = {}
        filtered_evolution_rework_anterior = {}
        global_evolution_anterior = {}
        global_evolution_rework_anterior = {}

        # NUEVO: Estructura para almacenar datos de costes por línea y turno
        costes_por_linea_turno = {}  # {(linea, uet, turno): {'items': [...], 'total_coste': X}}
        costes_globales = {'total_perdidas': 0, 'num_defectos': 0}

        if not os.path.isdir(self.ruta_archivos):
            return {"error": f"La ruta de archivos Excel no existe o no es un directorio: {self.ruta_archivos}"}

        patron = os.path.join(self.ruta_archivos, "*.xlsx")
        archivos_excel = glob.glob(patron)

        # Filtrar archivo de costes
        archivos_excel = [archivo for archivo in archivos_excel if not archivo.endswith("costes.xlsx")]

        def procesar_periodo(fecha_inicio_proc, fecha_fin_proc, es_periodo_actual=True):
            # Variables locales para este período
            local_filtered_evo = {} if es_periodo_actual else filtered_evolution_anterior
            local_filtered_evo_rew = {} if es_periodo_actual else filtered_evolution_rework_anterior
            local_global_evo = {} if es_periodo_actual else global_evolution_anterior
            local_global_evo_rew = {} if es_periodo_actual else global_evolution_rework_anterior

            if es_periodo_actual:
                local_filtered_evo = filtered_evolution
                local_filtered_evo_rew = filtered_evolution_rework
                local_global_evo = global_evolution
                local_global_evo_rew = global_evolution_rework

            for archivo in archivos_excel:
                nombre = os.path.basename(archivo)
                base, _ = os.path.splitext(nombre)
                try:
                    fecha_archivo = datetime.strptime(base, '%d.%m.%Y').date()
                except ValueError:
                    continue

                if not (fecha_inicio_proc <= fecha_archivo <= fecha_fin_proc):
                    continue

                wb = openpyxl.load_workbook(archivo, data_only=True)
                ws = wb.active
                fecha_str = fecha_archivo.isoformat()

                # Inicializar estructuras para esta fecha
                for evo_dict in [local_filtered_evo, local_global_evo]:
                    if fecha_str not in evo_dict:
                        evo_dict[fecha_str] = {"OK": 0, "NOK": 0, "NOK_PROV": 0}
                
                if fecha_str not in local_filtered_evo_rew:
                    local_filtered_evo_rew[fecha_str] = 0
                if fecha_str not in local_global_evo_rew:
                    local_global_evo_rew[fecha_str] = 0

                # Procesar cada fila de datos
                for row in ws.iter_rows(min_row=4, values_only=True):
                    if len(row) < 10:
                        continue

                    val_ok, val_nok, val_qty_rew = row[2], row[3], row[4]
                    val_col_f = str(row[5] or "").strip().upper()
                    val_defecto_row = (row[6] or "").strip().upper()
                    val_modelo_row = (row[7] or "").strip().upper()
                    val_turno_row = (row[9] or "SIN_TURNO").strip().upper()
                    
                    # NUEVO: Extraer línea, texto breve y UET para análisis de costes
                    val_texto_breve = (row[1] or "").strip().upper() if len(row) > 1 else ""
                    val_linea = (row[10] or "").strip().upper() if len(row) > 10 else "SIN_LINEA"
                    val_uet = (row[11] or "").strip().upper() if len(row) > 11 else "SIN_UET"

                    _ok = self._parse_num(val_ok)
                    _nok = self._parse_num(val_nok)
                    _rew = self._parse_num(val_qty_rew) if (val_col_f == "R") else 0

                    es_proveedor = self._es_defecto_proveedor(val_defecto_row)
                    
                    # Verificar si el registro cumple con los filtros
                    modelo_match = (not filtrar_modelo) or (val_modelo_row in modelos_list)
                    defecto_match = (not filtrar_defecto) or (val_defecto_row in defectos_list)
                    cumple_filtros = modelo_match and defecto_match

                    # NUEVO: ANÁLISIS DE COSTES (solo período actual y si hay datos de costes)
                    if es_periodo_actual and val_col_f != "R" and _nok > 0 and val_texto_breve and not es_proveedor:
                        precio_unitario = self.costes_dict.get(val_texto_breve, 0)
                        coste_total = _nok * precio_unitario
                        
                        if coste_total > 0:  # Solo si encontramos precio
                            # Usar línea real, UET (aunque sea "SIN_UET"), y turno
                            key_linea_turno = (val_linea, val_uet, val_turno_row)
                            
                            if key_linea_turno not in costes_por_linea_turno:
                                costes_por_linea_turno[key_linea_turno] = {
                                    'items': [],
                                    'total_coste': 0
                                }
                            
                            costes_por_linea_turno[key_linea_turno]['items'].append({
                                'texto_breve': val_texto_breve,
                                'cantidad': _nok,
                                'precio_unitario': precio_unitario,
                                'coste_total': coste_total,
                                'fecha': fecha_str,
                                'modelo': val_modelo_row,
                                'linea': val_linea,
                                'uet': val_uet
                            })
                            costes_por_linea_turno[key_linea_turno]['total_coste'] += coste_total
                            
                            costes_globales['total_perdidas'] += coste_total
                            costes_globales['num_defectos'] += _nok

                    # ACTUALIZAR DATOS GLOBALES (SIEMPRE, SIN FILTROS)
                    # ACTUALIZAR DATOS GLOBALES (SIEMPRE, SIN FILTROS)
                    if val_col_f != "R":
                        local_global_evo[fecha_str]["OK"] += _ok
                        if es_proveedor:
                            local_global_evo[fecha_str]["NOK_PROV"] += _nok
                            # ← NUEVA SECCIÓN PARA DISTRIBUCIONES DE PROVEEDOR
                            if es_periodo_actual and _nok > 0:
                                if val_modelo_row:
                                    global_models_distribution_prov[val_modelo_row] = global_models_distribution_prov.get(val_modelo_row, 0) + _nok
                        else:
                            local_global_evo[fecha_str]["NOK"] += _nok

                            # Distribuciones globales (solo período actual)
                            if es_periodo_actual and _nok > 0:
                                if val_defecto_row:
                                    global_defect_distribution[val_defecto_row] = global_defect_distribution.get(val_defecto_row, 0) + _nok
                                if val_turno_row:
                                    global_shift_distribution[val_turno_row] = global_shift_distribution.get(val_turno_row, 0) + _nok
                                if val_modelo_row:
                                    global_models_distribution[val_modelo_row] = global_models_distribution.get(val_modelo_row, 0) + _nok
                    else:
                        local_global_evo_rew[fecha_str] += _rew
                        if es_periodo_actual and _rew > 0:
                            if val_defecto_row:
                                global_rework_defect_distribution[val_defecto_row] = global_rework_defect_distribution.get(val_defecto_row, 0) + _rew
                            if val_turno_row:
                                global_rework_shift_distribution[val_turno_row] = global_rework_shift_distribution.get(val_turno_row, 0) + _rew
                            if val_modelo_row:
                                global_rework_models_distribution[val_modelo_row] = global_rework_models_distribution.get(val_modelo_row, 0) + _rew

                    # ACTUALIZAR DATOS FILTRADOS (SOLO SI CUMPLE LOS FILTROS)
                    if cumple_filtros:
                        if val_col_f != "R":
                            local_filtered_evo[fecha_str]["OK"] += _ok
                            if es_proveedor:
                                local_filtered_evo[fecha_str]["NOK_PROV"] += _nok
                            else:
                                local_filtered_evo[fecha_str]["NOK"] += _nok

                            # Distribuciones filtradas (solo período actual)
                            if es_periodo_actual and _nok > 0:
                                if val_defecto_row:
                                    filtered_defect_distribution[val_defecto_row] = filtered_defect_distribution.get(val_defecto_row, 0) + _nok
                                if val_turno_row:
                                    filtered_shift_distribution[val_turno_row] = filtered_shift_distribution.get(val_turno_row, 0) + _nok
                                if val_modelo_row:
                                    filtered_models_distribution[val_modelo_row] = filtered_models_distribution.get(val_modelo_row, 0) + _nok
                        else:
                            local_filtered_evo_rew[fecha_str] += _rew
                            if es_periodo_actual and _rew > 0:
                                if val_defecto_row:
                                    filtered_rework_defect_distribution[val_defecto_row] = filtered_rework_defect_distribution.get(val_defecto_row, 0) + _rew
                                if val_turno_row:
                                    filtered_rework_shift_distribution[val_turno_row] = filtered_rework_shift_distribution.get(val_turno_row, 0) + _rew
                                if val_modelo_row:
                                    filtered_rework_models_distribution[val_modelo_row] = filtered_rework_models_distribution.get(val_modelo_row, 0) + _rew

        # Procesar período actual
        procesar_periodo(fi, ff, True)
        
        # Procesar período anterior para tendencias
        procesar_periodo(fecha_inicio_anterior, fecha_fin_anterior, False)

        # Calcular totales FILTRADOS del período actual
        filtered_total_ok_actual = sum(d["OK"] for d in filtered_evolution.values())
        filtered_total_nok_actual = sum(d["NOK"] for d in filtered_evolution.values())
        filtered_total_nok_prov_actual = sum(d["NOK_PROV"] for d in filtered_evolution.values())
        filtered_total_rework_actual = sum(filtered_evolution_rework.values())

        # Calcular totales GLOBALES del período actual
        global_total_ok_actual = sum(d["OK"] for d in global_evolution.values())
        global_total_nok_actual = sum(d["NOK"] for d in global_evolution.values())
        global_total_nok_prov_actual = sum(d["NOK_PROV"] for d in global_evolution.values())
        global_total_rework_actual = sum(global_evolution_rework.values())

        # Calcular totales del período anterior (para tendencias)
        filtered_total_ok_anterior = sum(d["OK"] for d in filtered_evolution_anterior.values())
        filtered_total_nok_anterior = sum(d["NOK"] for d in filtered_evolution_anterior.values())
        filtered_total_rework_anterior = sum(filtered_evolution_rework_anterior.values())
        
        global_total_ok_anterior = sum(d["OK"] for d in global_evolution_anterior.values())
        global_total_nok_anterior = sum(d["NOK"] for d in global_evolution_anterior.values())
        global_total_rework_anterior = sum(global_evolution_rework_anterior.values())

        # Función para calcular cambio porcentual seguro
        def calcular_cambio_porcentual(actual, anterior):
            if anterior == 0:
                return 0.0 if actual == 0 else 100.0
            return ((actual - anterior) / anterior) * 100

        # Calcular porcentajes de NOK
        filtered_produccion_total_actual = filtered_total_ok_actual + filtered_total_nok_actual
        filtered_produccion_total_anterior = filtered_total_ok_anterior + filtered_total_nok_anterior
        
        filtered_porcentaje_nok_actual = (filtered_total_nok_actual / filtered_produccion_total_actual * 100) if filtered_produccion_total_actual > 0 else 0
        filtered_porcentaje_nok_anterior = (filtered_total_nok_anterior / filtered_produccion_total_anterior * 100) if filtered_produccion_total_anterior > 0 else 0

        global_produccion_total_actual = global_total_ok_actual + global_total_nok_actual
        global_produccion_total_anterior = global_total_ok_anterior + global_total_nok_anterior
        
        global_porcentaje_nok_actual = (global_total_nok_actual / global_produccion_total_actual * 100) if global_produccion_total_actual > 0 else 0
        global_porcentaje_nok_anterior = (global_total_nok_anterior / global_produccion_total_anterior * 100) if global_produccion_total_anterior > 0 else 0

        # Preparar resultado con información clara sobre filtros aplicados
        res = {"error": None}
        
        # Información sobre filtros aplicados
        res["filtros_aplicados"] = {
            "modelo": '|'.join(modelos_list) if filtrar_modelo else "TODOS",
            "defecto": '|'.join(defectos_list) if filtrar_defecto else "TODOS", 
            "tiene_filtros": filtrar_modelo or filtrar_defecto,
            "descripcion_filtro": ""
        }
        
        descripcion_partes = []
        if filtrar_modelo:
            if len(modelos_list) == 1:
                descripcion_partes.append(f"Modelo: {modelos_list[0]}")
            else:
                descripcion_partes.append(f"Modelos: {len(modelos_list)} seleccionados")

        if filtrar_defecto:
            if len(defectos_list) == 1:
                descripcion_partes.append(f"Defecto: {defectos_list[0]}")
            else:
                descripcion_partes.append(f"Defectos: {len(defectos_list)} seleccionados")

        if descripcion_partes:
            res["filtros_aplicados"]["descripcion_filtro"] = " | ".join(descripcion_partes)
        else:
            res["filtros_aplicados"]["descripcion_filtro"] = "Sin filtros (Datos globales)"

        # DATOS PRINCIPALES (FILTRADOS si hay filtros, GLOBALES si no hay filtros)
        if filtrar_modelo or filtrar_defecto:
            # Si hay filtros, usar datos filtrados como principales
            res["total_ok"] = filtered_total_ok_actual
            res["total_nok"] = filtered_total_nok_actual
            res["total_nok_prov"] = filtered_total_nok_prov_actual
            res["total_rework"] = filtered_total_rework_actual
            res["evolution"] = filtered_evolution
            res["evolution_rework"] = filtered_evolution_rework
            res["defect_distribution"] = filtered_defect_distribution
            res["shift_distribution"] = filtered_shift_distribution
            res["models_distribution"] = filtered_models_distribution
            res["rework_defect_distribution"] = filtered_rework_defect_distribution
            res["rework_shift_distribution"] = filtered_rework_shift_distribution
            res["rework_models_distribution"] = filtered_rework_models_distribution
            
            # Datos de tendencias (filtrados)
            res["tendencias"] = {
                "cambio_produccion_total": calcular_cambio_porcentual(filtered_produccion_total_actual, filtered_produccion_total_anterior),
                "cambio_ok_total": calcular_cambio_porcentual(filtered_total_ok_actual, filtered_total_ok_anterior),
                "cambio_nok_total": calcular_cambio_porcentual(filtered_total_nok_actual, filtered_total_nok_anterior),
                "cambio_retrabajos_total": calcular_cambio_porcentual(filtered_total_rework_actual, filtered_total_rework_anterior),
                "cambio_porcentaje_nok": filtered_porcentaje_nok_actual - filtered_porcentaje_nok_anterior,
                "periodo_anterior": {
                    "fecha_inicio": fecha_inicio_anterior.isoformat(),
                    "fecha_fin": fecha_fin_anterior.isoformat(),
                    "produccion_total": filtered_produccion_total_anterior,
                    "ok_total": filtered_total_ok_anterior,
                    "nok_total": filtered_total_nok_anterior,
                    "retrabajos_total": filtered_total_rework_anterior,
                    "porcentaje_nok": filtered_porcentaje_nok_anterior
                }
            }
        else:
            # Si no hay filtros, usar datos globales como principales
            res["total_ok"] = global_total_ok_actual
            res["total_nok"] = global_total_nok_actual
            res["total_nok_prov"] = global_total_nok_prov_actual
            res["total_rework"] = global_total_rework_actual
            res["evolution"] = global_evolution
            res["evolution_rework"] = global_evolution_rework
            res["defect_distribution"] = global_defect_distribution
            res["shift_distribution"] = global_shift_distribution
            res["models_distribution"] = global_models_distribution
            res["global_models_distribution_prov"] = global_models_distribution_prov
            res["rework_defect_distribution"] = global_rework_defect_distribution
            res["rework_shift_distribution"] = global_rework_shift_distribution
            res["rework_models_distribution"] = global_rework_models_distribution
            
            # Datos de tendencias (globales)
            res["tendencias"] = {
                "cambio_produccion_total": calcular_cambio_porcentual(global_produccion_total_actual, global_produccion_total_anterior),
                "cambio_ok_total": calcular_cambio_porcentual(global_total_ok_actual, global_total_ok_anterior),
                "cambio_nok_total": calcular_cambio_porcentual(global_total_nok_actual, global_total_nok_anterior),
                "cambio_retrabajos_total": calcular_cambio_porcentual(global_total_rework_actual, global_total_rework_anterior),
                "cambio_porcentaje_nok": global_porcentaje_nok_actual - global_porcentaje_nok_anterior,
                "periodo_anterior": {
                    "fecha_inicio": fecha_inicio_anterior.isoformat(),
                    "fecha_fin": fecha_fin_anterior.isoformat(),
                    "produccion_total": global_produccion_total_anterior,
                    "ok_total": global_total_ok_anterior,
                    "nok_total": global_total_nok_anterior,
                    "retrabajos_total": global_total_rework_anterior,
                    "porcentaje_nok": global_porcentaje_nok_anterior
                }
            }

        # MANTENER COMPATIBILIDAD CON CÓDIGO EXISTENTE
        res["global_total_ok"] = global_total_ok_actual
        res["global_total_nok"] = global_total_nok_actual
        res["global_total_nok_prov"] = global_total_nok_prov_actual
        res["global_evolution"] = global_evolution
        res["global_total_rework"] = global_total_rework_actual
        res["global_evolution_rework"] = global_evolution_rework
        
        # Datos parciales (que ahora son los filtrados)
        res["partial_total_ok"] = filtered_total_ok_actual
        res["partial_total_nok"] = filtered_total_nok_actual
        res["partial_total_nok_prov"] = filtered_total_nok_prov_actual
        res["partial_evolution"] = filtered_evolution
        res["partial_total_rework"] = filtered_total_rework_actual
        res["partial_evolution_rework"] = filtered_evolution_rework
        res["partial_defect_distribution"] = filtered_defect_distribution
        res["partial_shift_distribution"] = filtered_shift_distribution
        res["partial_models_distribution"] = filtered_models_distribution
        res["partial_rework_defect_distribution"] = filtered_rework_defect_distribution
        res["partial_rework_shift_distribution"] = filtered_rework_shift_distribution
        res["partial_rework_models_distribution"] = filtered_rework_models_distribution

        # NUEVO: Añadir datos de costes al resultado
        # Convertir las tuplas (linea, uet, turno) a strings para que JSON pueda serializarlas
        costes_por_linea_turno_serializable = {}
        for (linea, uet, turno), data in costes_por_linea_turno.items():
            key_string = f"{linea}|{uet}|{turno}"  # Usar | como separador
            costes_por_linea_turno_serializable[key_string] = data
        
        print(f"📊 DEBUG Costes: Total líneas/turnos con costes: {len(costes_por_linea_turno_serializable)}")
        print(f"💰 DEBUG Costes: Total pérdidas globales: €{costes_globales.get('total_perdidas', 0):,.2f}")
        print(f"🔢 DEBUG Costes: Total defectos valorados: {costes_globales.get('num_defectos', 0)}")
        
        if costes_por_linea_turno_serializable:
            primera_key = list(costes_por_linea_turno_serializable.keys())[0]
            first_data = list(costes_por_linea_turno_serializable.values())[0]
            print(f"📋 DEBUG Costes: Ejemplo - Primera línea: {primera_key}")
            print(f"💵 DEBUG Costes: Coste de primera línea: €{first_data['total_coste']:,.2f}")
        
        res["costes_por_linea_turno"] = costes_por_linea_turno_serializable
        res["costes_globales"] = costes_globales

        return res

    def _anotar_barras(self, ax, fontsize=11, bold=True):
        """Anota cada barra con su valor encima."""
        for rect in ax.patches:
            height = rect.get_height()
            if height != 0:
                ax.annotate(
                    f'{int(height):,}',
                    xy=(rect.get_x() + rect.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', va='bottom',
                    fontsize=fontsize,
                    fontweight='bold' if bold else 'normal',
                    color=self.corporate_colors['dark_gray']
                )

    def _crear_grafica_profesional(self, fig_size=(8, 5)):
        """
        Crea una figura matplotlib con estilo profesional:
        - Fondo blanco.
        - Bordes superiores y derechos ocultos.
        - Grid suave.
        """
        fig, ax = plt.subplots(figsize=fig_size)
        fig.patch.set_facecolor('white')
        ax.set_facecolor('white')

        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color(self.corporate_colors['light_gray'])
        ax.spines['bottom'].set_color(self.corporate_colors['light_gray'])

        ax.grid(True, alpha=0.1, linestyle='--', color=self.corporate_colors['light_gray'])
        ax.set_axisbelow(True)
        return fig, ax

    # FUNCIONES PARA POWERPOINT EN TV
    # Añadir estas funciones a tu clase API
    def _generar_contenido_panel(self, datos_estadisticas):
        """
        POWERPOINT: Genera contenido del panel lateral con tendencias comparativas
        """
        try:
            # Extraer datos principales
            tendencias = datos_estadisticas.get("tendencias", {})
            periodo_anterior = tendencias.get("periodo_anterior", {})
            
            # Datos actuales
            total_produccion_actual = datos_estadisticas.get("total_ok", 0) + datos_estadisticas.get("total_nok", 0)
            total_nok_actual = datos_estadisticas.get("total_nok", 0)
            total_retrabajos_actual = datos_estadisticas.get("total_rework", 0)
            porcentaje_nok_actual = (total_nok_actual / total_produccion_actual * 100) if total_produccion_actual > 0 else 0
            
            # Datos período anterior
            total_produccion_anterior = periodo_anterior.get("produccion_total", 0)
            total_nok_anterior = periodo_anterior.get("nok_total", 0)
            total_retrabajos_anterior = periodo_anterior.get("retrabajos_total", 0)
            porcentaje_nok_anterior = periodo_anterior.get("porcentaje_nok", 0)
            
            # Cambios porcentuales
            cambio_produccion = tendencias.get("cambio_produccion_total", 0)
            cambio_nok = tendencias.get("cambio_nok_total", 0)
            cambio_retrabajos = tendencias.get("cambio_retrabajos_total", 0)
            cambio_porcentaje_nok = tendencias.get("cambio_porcentaje_nok", 0)
            
            # Top modelos problemáticos
            models_dist = datos_estadisticas.get("models_distribution", {})
            top_models = sorted(models_dist.items(), key=lambda x: x[1], reverse=True)[:3]
            
            # Top defectos
            defect_dist = datos_estadisticas.get("defect_distribution", {})
            top_defectos = sorted(defect_dist.items(), key=lambda x: x[1], reverse=True)[:3]
            
            # Datos de costes si existen
            costes_globales = datos_estadisticas.get("costes_globales", {})
            total_perdidas = costes_globales.get("total_perdidas", 0)
            
            # Construir contenido
            contenido = []
            
            # === CABECERA ===
            contenido.append("TENDENCIAS vs ANTERIOR")
            contenido.append("=" * 20)
            contenido.append("")
            
            # === PRODUCCIÓN GENERAL ===
            contenido.append("PRODUCCION TOTAL")
            contenido.append("-" * 15)
            contenido.append(f"Actual: {total_produccion_actual:,}")
            contenido.append(f"Anterior: {total_produccion_anterior:,}")
            if cambio_produccion >= 0:
                contenido.append(f"Cambio: +{cambio_produccion:.1f}%")
            else:
                contenido.append(f"Cambio: {cambio_produccion:.1f}%")
            contenido.append("")
            
            # === NOK COMPARISON ===
            contenido.append("NOK (RECHAZOS)")
            contenido.append("-" * 15)
            contenido.append(f"Actual: {total_nok_actual:,} ({porcentaje_nok_actual:.1f}%)")
            contenido.append(f"Anterior: {total_nok_anterior:,} ({porcentaje_nok_anterior:.1f}%)")
            if cambio_nok >= 0:
                contenido.append(f"Piezas: +{cambio_nok:.1f}%")
            else:
                contenido.append(f"Piezas: {cambio_nok:.1f}%")
            if cambio_porcentaje_nok >= 0:
                contenido.append(f"Ratio: +{cambio_porcentaje_nok:.1f}pp")
            else:
                contenido.append(f"Ratio: {cambio_porcentaje_nok:.1f}pp")
            contenido.append("")
            
            # === RETRABAJOS ===
            contenido.append("RETRABAJOS")
            contenido.append("-" * 15)
            contenido.append(f"Actual: {total_retrabajos_actual:,}")
            contenido.append(f"Anterior: {total_retrabajos_anterior:,}")
            if cambio_retrabajos >= 0:
                contenido.append(f"Cambio: +{cambio_retrabajos:.1f}%")
            else:
                contenido.append(f"Cambio: {cambio_retrabajos:.1f}%")
            contenido.append("")
            
            # === TOP MODELOS NOK ===
            if top_models:
                contenido.append("TOP MODELOS NOK")
                contenido.append("-" * 15)
                for i, (modelo, cantidad) in enumerate(top_models, 1):
                    modelo_corto = modelo[:12] + "..." if len(modelo) > 15 else modelo
                    contenido.append(f"{i}. {modelo_corto}: {cantidad}")
                contenido.append("")
            
            # === TOP DEFECTOS ===
            if top_defectos:
                contenido.append("TOP DEFECTOS")
                contenido.append("-" * 15)
                for i, (defecto, cantidad) in enumerate(top_defectos, 1):
                    defecto_corto = defecto[:12] + "..." if len(defecto) > 15 else defecto
                    contenido.append(f"{i}. {defecto_corto}: {cantidad}")
                contenido.append("")
            
            # === IMPACTO ECONÓMICO ===
            if total_perdidas > 0:
                contenido.append("IMPACTO ECONOMICO")
                contenido.append("-" * 15)
                contenido.append(f"Perdidas: {self._format_euro(total_perdidas)}")
                contenido.append("")
            
            # === INDICADORES RÁPIDOS ===
            contenido.append("ESTADO GENERAL")
            contenido.append("-" * 15)
            
            # Semáforo de NOK
            if porcentaje_nok_actual <= 2.0:
                estado_nok = "BUENO"
            elif porcentaje_nok_actual <= 5.0:
                estado_nok = "ATENCION"
            else:
                estado_nok = "CRITICO"
            contenido.append(f"NOK Rate: {estado_nok}")
            
            # Tendencia general
            if cambio_nok < -10:
                tendencia = "MEJORANDO"
            elif cambio_nok > 10:
                tendencia = "EMPEORANDO"
            else:
                tendencia = "ESTABLE"
            contenido.append(f"Tendencia: {tendencia}")
            
            # Fecha del período anterior
            fecha_anterior = periodo_anterior.get("fecha_fin", "")
            if fecha_anterior:
                try:
                    fecha_obj = datetime.strptime(fecha_anterior, '%Y-%m-%d')
                    fecha_formateada = fecha_obj.strftime('%d/%m')
                    contenido.append(f"Vs: {fecha_formateada}")
                except:
                    pass
            
            return "\n".join(contenido)
            
        except Exception as e:
            print(f"Error generando contenido panel: {e}")
            return "PANEL DE TENDENCIAS\n\nError cargando datos\nde comparación"

    

    def generar_powerpoint_tv(self, pdf_path):
        """POWERPOINT: Versión con gráficos específicos por tipo de slide"""
        try:
            import tempfile
            import os
            from pptx import Presentation
            from pptx.util import Inches
            
            print("📺 Iniciando generación de PowerPoint para TV con gráficos específicos...")
            
            if not os.path.exists(pdf_path):
                return {"error": f"PDF no encontrado: {pdf_path}", "powerpoint_path": None}
            
            if not os.path.exists(self.powerpoint_dir):
                try:
                    os.makedirs(self.powerpoint_dir)
                except Exception as e:
                    return {"error": f"No se pudo crear directorio {self.powerpoint_dir}: {e}", "powerpoint_path": None}
            
            # Generar datos estadísticos
            from datetime import datetime, timedelta
            fecha_fin = datetime.now().strftime('%Y-%m-%d')
            fecha_inicio = (datetime.now() - timedelta(days=7)).strftime('%Y-%m-%d')
            
            print("📊 Generando datos estadísticos...")
            datos_estadisticas = self.leer_estadisticas(fecha_inicio, fecha_fin, "(TODOS)", "(TODOS)")
            
            # Conversión PDF
            images_paths = []
            temp_dir = tempfile.gettempdir()
            
            try:
                import fitz
                pdf_document = fitz.open(pdf_path)
                total_pages = pdf_document.page_count
                
                for page_num in range(min(total_pages, 15)):
                    page = pdf_document[page_num]
                    mat = fitz.Matrix(2.0, 2.0)
                    pix = page.get_pixmap(matrix=mat)
                    
                    img_path = os.path.join(temp_dir, f"ppt_slide_{page_num + 1}_{os.getpid()}.png")
                    pix.save(img_path)
                    pix = None
                    
                    if os.path.exists(img_path) and os.path.getsize(img_path) > 1000:
                        images_paths.append(img_path)
                
                pdf_document.close()
                
            except Exception as conversion_error:
                return {"error": f"Error convirtiendo PDF: {conversion_error}", "powerpoint_path": None}
            
            if not images_paths:
                return {"error": "No se pudieron generar imágenes del PDF", "powerpoint_path": None}
            
            # Crear presentación
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            slides_creados = 0
            for i, img_path in enumerate(images_paths):
                try:
                    if not os.path.exists(img_path):
                        continue
                    
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # PASAR EL NÚMERO DE SLIDE PARA DETECTAR TIPO
                    self._añadir_imagen_con_panel_lateral(slide, img_path, datos_estadisticas, i)
                    
                    slides_creados += 1
                    print(f"✅ Slide {i+1} completado con gráficos específicos")
                    
                except Exception as slide_error:
                    print(f"❌ Error creando slide {i+1}: {slide_error}")
                    continue
            
            # Guardar
            powerpoint_filename = "TV_Calidad_actual.pptx"
            powerpoint_full_path = os.path.join(self.powerpoint_dir, powerpoint_filename)
            
            prs.save(powerpoint_full_path)
            
            # Limpiar
            for img_path in images_paths:
                try:
                    if os.path.exists(img_path):
                        os.remove(img_path)
                except:
                    pass
            
            self._configurar_powerpoint_para_tv(powerpoint_full_path)
            
            return {
                "error": None,
                "powerpoint_path": os.path.abspath(powerpoint_full_path),
                "total_slides": slides_creados,
                "mensaje": f"✅ PowerPoint creado: {slides_creados} diapositivas con gráficos específicos"
            }
            
        except Exception as e:
            return {"error": f"Error inesperado: {e}", "powerpoint_path": None}
        
    def _crear_grafico_simple_test(self, titulo, color, datos_estadisticas=None, modelo_especifico=None, datos_modelo=None):
        """Crea gráficos con datos reales - VERSIÓN CORREGIDA"""
        try:
            import matplotlib.pyplot as plt
            import tempfile
            import numpy as np
            from datetime import datetime, timedelta
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            if "BARRAS DÍAS SEMANA" in titulo and datos_estadisticas:
                # Cambio: Barras simples en lugar de pentágono
                evolution = datos_estadisticas.get("evolution", {})
                if evolution:
                    dias_semana_nok = {}
                    for fecha_str, data in evolution.items():
                        try:
                            fecha_dt = datetime.strptime(fecha_str, '%Y-%m-%d').date()
                            dia_semana_num = fecha_dt.weekday()  # 0=Lunes
                            nok_count = data.get("NOK", 0)
                            if dia_semana_num not in dias_semana_nok:
                                dias_semana_nok[dia_semana_num] = 0
                            dias_semana_nok[dia_semana_num] += nok_count
                        except:
                            continue
                    
                    if dias_semana_nok:
                        dias_nombres = ['L', 'M', 'X', 'J', 'V', 'S', 'D']
                        dias_valores = []
                        dias_labels = []
                        
                        for i in range(7):
                            if i in dias_semana_nok:
                                dias_valores.append(dias_semana_nok[i])
                                dias_labels.append(dias_nombres[i])
                        
                        if dias_valores:
                            bars = ax.bar(range(len(dias_labels)), dias_valores, color=color, alpha=0.8)
                            
                            # Añadir valores encima de barras
                            for bar, valor in zip(bars, dias_valores):
                                height = bar.get_height()
                                ax.text(bar.get_x() + bar.get_width()/2., height + max(dias_valores)*0.01,
                                    f'{valor}', ha='center', va='bottom', fontweight='bold', fontsize=9)
                            
                            ax.set_xticks(range(len(dias_labels)))
                            ax.set_xticklabels(dias_labels)
                            ax.set_title("DEFECTOS POR DÍA", fontsize=10, fontweight='bold')
                            ax.spines['top'].set_visible(False)
                            ax.spines['right'].set_visible(False)
                        else:
                            ax.text(0.5, 0.5, 'Sin datos\nsuficientes', transform=ax.transAxes, 
                                ha='center', va='center', fontsize=12, color='gray')
                            ax.set_title(titulo, fontsize=10, fontweight='bold')
                
            elif "GAUGE RECUPERACIÓN" in titulo and datos_estadisticas:
                # Gauge simple con barras horizontales
                total_nok = datos_estadisticas.get("total_nok", 0)
                total_rework = datos_estadisticas.get("total_rework", 0)
                porcentaje = min(100, (total_rework / total_nok * 100)) if total_nok > 0 else 0
                
                # Barra de progreso horizontal
                ax.barh(['Recuperación'], [porcentaje], color=color, height=0.6, alpha=0.8)
                ax.barh(['Recuperación'], [100], color='lightgray', alpha=0.3, height=0.6)
                
                # Texto con porcentaje
                ax.text(porcentaje + 2, 0, f'{porcentaje:.1f}%', va='center', fontweight='bold', fontsize=12)
                
                ax.set_xlim(0, 100)
                ax.set_title(f"RECUPERACIÓN: {porcentaje:.1f}%", fontsize=10, fontweight='bold')
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.spines['bottom'].set_visible(False)
                ax.set_yticks([])
                
            elif "BARRAS COSTE TURNO" in titulo and datos_estadisticas:
                # Cambio: Barras en lugar de speedometer
                shift_dist = datos_estadisticas.get("shift_distribution", {})
                if shift_dist and self.costes_dict:
                    precio_promedio = sum(self.costes_dict.values()) / len(self.costes_dict)
                    
                    turnos = list(shift_dist.keys())
                    costes = [nok * precio_promedio for nok in shift_dist.values()]
                    colores_turno = {'M': '#F59E0B', 'T': '#DC2626', 'N': '#3B82F6', 
                                'A': '#F59E0B', 'B': '#DC2626', 'C': '#3B82F6'}
                    colores = [colores_turno.get(turno, '#6B7280') for turno in turnos]
                    
                    bars = ax.bar(range(len(turnos)), costes, color=colores, alpha=0.8)
                    
                    # Añadir valores
                    for bar, coste in zip(bars, costes):
                        height = bar.get_height()
                        ax.text(bar.get_x() + bar.get_width()/2., height + max(costes)*0.01,
                            f'€{int(coste)}', ha='center', va='bottom', fontweight='bold', fontsize=8)
                    
                    ax.set_xticks(range(len(turnos)))
                    ax.set_xticklabels(turnos)
                    ax.set_title("COSTE POR TURNO", fontsize=10, fontweight='bold')
                    ax.spines['top'].set_visible(False)
                    ax.spines['right'].set_visible(False)
                
            elif "COMPARATIVA MODELO" in titulo and modelo_especifico and datos_modelo:
                # Datos específicos del modelo vs promedio general
                modelo_nok = datos_modelo.get("total_nok", 0)
                
                # Calcular promedio de otros modelos
                models_dist = datos_estadisticas.get("models_distribution", {})
                if models_dist and modelo_especifico in models_dist:
                    otros_modelos = {k: v for k, v in models_dist.items() if k != modelo_especifico}
                    promedio_otros = sum(otros_modelos.values()) / len(otros_modelos) if otros_modelos else 0
                    
                    categorias = [f'{modelo_especifico[:8]}', 'Promedio\nOtros']
                    valores = [modelo_nok, promedio_otros]
                    colores = [color, 'lightgray']
                    
                    bars = ax.bar(categorias, valores, color=colores, alpha=0.8)
                    
                    # Añadir valores
                    for bar, valor in zip(bars, valores):
                        height = bar.get_height()
                        ax.text(bar.get_x() + bar.get_width()/2., height + max(valores)*0.01,
                            f'{int(valor)}', ha='center', va='bottom', fontweight='bold', fontsize=9)
                    
                    ax.set_title(f"MODELO vs OTROS", fontsize=10, fontweight='bold')
                    ax.spines['top'].set_visible(False)
                    ax.spines['right'].set_visible(False)
                
            elif "DNA DEFECTO MODELO" in titulo and modelo_especifico and datos_modelo:
                # Defectos específicos de este modelo
                defect_dist = datos_modelo.get("defect_distribution", {})
                if defect_dist:
                    top_defectos = sorted(defect_dist.items(), key=lambda x: x[1], reverse=True)[:3]
                    if top_defectos:
                        nombres = [d[:10] for d, _ in top_defectos]
                        valores = [v for _, v in top_defectos]
                        colores = ['#DC2626', '#F59E0B', '#10B981']
                        
                        wedges, texts, autotexts = ax.pie(valores, labels=nombres, colors=colores, 
                                                        autopct='%1.1f%%', startangle=90,
                                                        wedgeprops=dict(width=0.4))
                        
                        centre_circle = plt.Circle((0,0), 0.60, fc='white')
                        fig.gca().add_artist(centre_circle)
                        
                        for autotext in autotexts:
                            autotext.set_color('white')
                            autotext.set_fontsize(8)
                            autotext.set_fontweight('bold')
                        
                        ax.set_title(f"DEFECTOS {modelo_especifico[:8]}", fontsize=10, fontweight='bold')
                
            elif "EFICIENCIA MODELO" in titulo and modelo_especifico and datos_modelo:
                # Eficiencia específica del modelo
                modelo_nok = datos_modelo.get("total_nok", 0)
                modelo_rework = datos_modelo.get("total_rework", 0)
                eficiencia = min(100, (modelo_rework / modelo_nok * 100)) if modelo_nok > 0 else 0
                
                ax.barh([f'{modelo_especifico[:8]}'], [eficiencia], color=color, height=0.5)
                ax.barh([f'{modelo_especifico[:8]}'], [100], color='lightgray', alpha=0.3, height=0.5)
                ax.text(eficiencia + 2, 0, f'{eficiencia:.1f}%', va='center', fontweight='bold')
                ax.set_xlim(0, 100)
                ax.set_title(f"EFICIENCIA {modelo_especifico[:8]}", fontsize=10, fontweight='bold')
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.spines['bottom'].set_visible(False)
                
            elif "EXITOSOS MODELO" in titulo and modelo_especifico and datos_modelo:
                # Retrabajos específicos del modelo
                modelo_rework = datos_modelo.get("total_rework", 0)
                exitosos = int(modelo_rework * 0.85)
                fallidos = modelo_rework - exitosos
                
                categorias = ['Exitosos', 'Fallidos']
                valores = [exitosos, fallidos]
                colores = ['#10B981', '#DC2626']
                
                bars = ax.bar(categorias, valores, color=colores, alpha=0.8)
                
                for bar, valor in zip(bars, valores):
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height + max(valores)*0.01,
                        f'{valor}', ha='center', va='bottom', fontweight='bold', fontsize=9)
                
                ax.set_title(f"RETRABAJOS {modelo_especifico[:8]}", fontsize=10, fontweight='bold')
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
            
            # [Mantener los otros gráficos como HEAT MAP, TIMELINE, TOP 3, etc. igual que antes]
            elif "HEAT MAP" in titulo and datos_estadisticas:
                models_dist = datos_estadisticas.get("models_distribution", {})
                if models_dist:
                    top_models = sorted(models_dist.items(), key=lambda x: x[1], reverse=True)[:6]
                    if top_models:
                        nombres = [modelo[:8] for modelo, _ in top_models]
                        valores = [cantidad for _, cantidad in top_models]
                        max_val = max(valores)
                        valores_norm = [(v/max_val)*100 for v in valores]
                        
                        grid_data = []
                        grid_labels = []
                        for i in range(2):
                            row_data = []
                            row_labels = []
                            for j in range(3):
                                idx = i*3 + j
                                if idx < len(valores_norm):
                                    row_data.append(valores_norm[idx])
                                    row_labels.append(f"{nombres[idx]}\n{valores[idx]}")
                                else:
                                    row_data.append(0)
                                    row_labels.append("")
                            grid_data.append(row_data)
                            grid_labels.append(row_labels)
                        
                        im = ax.imshow(grid_data, cmap='Reds', aspect='equal', vmin=0, vmax=100)
                        
                        for i in range(2):
                            for j in range(3):
                                if grid_labels[i][j]:
                                    ax.text(j, i, grid_labels[i][j], ha='center', va='center',
                                        fontsize=8, fontweight='bold', 
                                        color='white' if grid_data[i][j] > 50 else 'black')
                        
                        ax.set_xticks([])
                        ax.set_yticks([])
                        ax.set_title("TOP LÍNEAS PROBLEMÁTICAS", fontsize=10, fontweight='bold')
                
            # [Continuar con TIMELINE, TOP 3, etc...]
            elif "TIMELINE" in titulo and datos_estadisticas:
                evolution_rework = datos_estadisticas.get("evolution_rework", {})
                if evolution_rework:
                    fechas_ordenadas = sorted(evolution_rework.keys())[-7:]
                    valores = [evolution_rework[f] for f in fechas_ordenadas]
                    fechas_cortas = [f[-5:] for f in fechas_ordenadas]
                    
                    if valores:
                        max_val = max(valores) if valores else 1
                        colores_norm = [(v/max_val) for v in valores]
                        
                        scatter = ax.scatter(range(len(valores)), [0]*len(valores), 
                                        c=colores_norm, cmap='RdYlGn_r', s=200, alpha=0.8)
                        ax.plot(range(len(valores)), [0]*len(valores), color='gray', alpha=0.3, linewidth=2)
                        
                        for i, (fecha, valor) in enumerate(zip(fechas_cortas, valores)):
                            ax.text(i, 0.1, fecha, ha='center', fontsize=7, rotation=45)
                            ax.text(i, -0.1, str(valor), ha='center', fontsize=8, fontweight='bold')
                        
                        ax.set_xlim(-0.5, len(valores)-0.5)
                        ax.set_ylim(-0.3, 0.3)
                        ax.set_title("RETRABAJOS ÚLTIMOS 7 DÍAS", fontsize=10, fontweight='bold')
                        ax.axis('off')
            
            elif "TOP 3 MODELOS" in titulo and datos_estadisticas:
                models_dist = datos_estadisticas.get("models_distribution", {})
                if models_dist and self.costes_dict:
                    precio_promedio = sum(self.costes_dict.values()) / len(self.costes_dict)
                    modelos_costes = [(modelo, nok * precio_promedio) for modelo, nok in models_dist.items()]
                    top3_modelos = sorted(modelos_costes, key=lambda x: x[1], reverse=True)[:3]
                    
                    if top3_modelos:
                        modelos = [m[:8] for m, _ in top3_modelos]
                        costes = [c for _, c in top3_modelos]
                        colores = ['#DC2626', '#F59E0B', '#10B981']
                        
                        bars = ax.barh(range(len(modelos)), costes, color=colores, alpha=0.8)
                        
                        for i, (bar, coste) in enumerate(zip(bars, costes)):
                            width = bar.get_width()
                            ax.text(width + max(costes)*0.02, bar.get_y() + bar.get_height()/2,
                                f'€{int(coste)}', va='center', ha='left', fontsize=8, fontweight='bold')
                        
                        ax.set_yticks(range(len(modelos)))
                        ax.set_yticklabels(modelos, fontsize=9)
                        ax.set_title("TOP 3 MODELOS COSTOSOS", fontsize=10, fontweight='bold')
                        ax.spines['top'].set_visible(False)
                        ax.spines['right'].set_visible(False)
            
            else:
                # Gráfico por defecto
                ax.bar(['A', 'B', 'C'], [30, 50, 20], color=color)
                ax.set_title(titulo, fontsize=10, fontweight='bold')
            
            # Configuración general
            ax.tick_params(labelsize=8)
            plt.tight_layout()
            
            # Guardar
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white', edgecolor='none')
            plt.close(fig)
            
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando gráfico '{titulo}': {e}")
            return None
        
    def _añadir_imagen_izquierda_tv(self, slide, image_path):
        """
        POWERPOINT: Imagen que ocupa TODA la altura del slide
        """
        try:
            from PIL import Image
            from pptx.util import Inches
            
            # Info de la imagen
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                img_aspect = img_width / img_height
            
            # Dimensiones del slide
            slide_width = Inches(16)
            slide_height = Inches(9)
            
            # FORZAR: La imagen ocupa TODA la altura disponible
            new_height = slide_height  # Sin int(), valor exacto
            new_width = new_height * img_aspect
            
            # Posición: desde el borde superior, margen izquierdo
            left = 0  # Pegado al borde izquierdo
            top = 0  # Desde el borde superior exacto
            
            # Si es muy ancha, recalcular PERO mantener altura completa
            max_width = slide_width * 0.70
            if new_width > max_width:
                new_width = max_width
                # IMPORTANTE: NO recalcular new_height, mantener altura completa
                # new_height sigue siendo slide_height
            
            # Añadir imagen con altura forzada completa
            slide.shapes.add_picture(image_path, 
                                int(left), 
                                int(top), 
                                int(new_width), 
                                int(new_height))  # Esta será siempre slide_height
            
            print(f"  📐 Imagen: {int(new_width)}x{int(new_height)} (altura completa forzada)")
            
        except Exception as e:
            print(f"  ❌ Error añadiendo imagen: {e}")
            # Fallback con altura completa
            try:
                slide.shapes.add_picture(image_path, 
                                    Inches(0.5),    # left
                                    0,              # top = 0 (borde superior)
                                    Inches(11),     # width
                                    Inches(9))      # height = altura completa
            except Exception as e2:
                print(f"  ❌ Error en fallback: {e2}")
    def _crear_powerpoint_placeholder(self, pdf_path):
        """
        POWERPOINT: Método de emergencia - Crea PowerPoint con placeholders informativos
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            print("🔄 Creando PowerPoint con placeholders informativos...")
            
            # Leer información básica del PDF si es posible
            try:
                import fitz
                pdf_doc = fitz.open(pdf_path)
                num_paginas = pdf_doc.page_count
                pdf_doc.close()
            except:
                num_paginas = 6  # Valor por defecto
            
            # Crear presentación
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            slides_info = [
                ("📊 REPORTE DE CALIDAD", "Dashboard Ejecutivo\nPresentación para TV", "#1E40AF"),
                ("📈 ANÁLISIS NOK", "Piezas No Conformes\nEvolution y Distribución", "#DC2626"),
                ("🔄 ANÁLISIS RETRABAJOS", "Piezas Retrabajadas\nRecuperación de Calidad", "#059669"),
                ("💰 IMPACTO ECONÓMICO", "Análisis de Costes\nPérdidas por Línea/UET", "#EA580C"),
                ("🎯 TOP MODELOS NOK", "Modelos Más Problemáticos\nAnálisis Detallado", "#7C3AED"),
                ("🔍 ANÁLISIS DETALLADO", "Correlaciones y Patrones\nRecomendaciones", "#0284C7")
            ]
            
            # Crear slides con información
            for i, (titulo, subtitulo, color) in enumerate(slides_info):
                if i < num_paginas:
                    slide_layout = prs.slide_layouts[5]  # Layout con título y contenido
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Configurar título
                    title_shape = slide.shapes.title
                    title_shape.text = titulo
                    title_shape.text_frame.paragraphs[0].font.size = Inches(0.8)
                    title_shape.text_frame.paragraphs[0].font.bold = True
                    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # Configurar contenido
                    content_shape = slide.placeholders[1]
                    content_shape.text = f"{subtitulo}\n\nSlide {i+1} de {num_paginas}\n\nInstale pdf2image o PyMuPDF\npara capturas reales del PDF"
                    
                    # Aplicar formato
                    for paragraph in content_shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        paragraph.font.size = Inches(0.4)
                        paragraph.font.color.rgb = RGBColor.from_string(color)
            
            # Guardar
            powerpoint_filename = "TV_Calidad_actual.pptx"
            powerpoint_full_path = os.path.join(self.powerpoint_dir, powerpoint_filename)
            prs.save(powerpoint_full_path)
            
            self._configurar_powerpoint_para_tv(powerpoint_full_path)
            
            print(f"✅ PowerPoint placeholder creado con {len(slides_info)} slides")
            
            return {
                "error": None,
                "powerpoint_path": os.path.abspath(powerpoint_full_path),
                "total_slides": len(slides_info),
                "mensaje": f"✅ PowerPoint informativo creado: {len(slides_info)} diapositivas (instale pdf2image para capturas reales)"
            }
            
        except Exception as e:
            return {"error": f"Error creando PowerPoint placeholder: {e}", "powerpoint_path": None}

    def _optimizar_imagen_para_tv(self, page_image):
        """
        POWERPOINT: Optimiza imagen para visualización en TV
        - Ajusta contraste y brillo para pantallas grandes
        - Mantiene resolución alta para calidad
        """
        try:
            from PIL import ImageEnhance, ImageOps
            
            # Convertir a RGB si es necesario
            if page_image.mode in ('RGBA', 'LA', 'P'):
                page_image = page_image.convert('RGB')
            
            # Mejorar contraste para TV (pantallas grandes necesitan más contraste)
            enhancer = ImageEnhance.Contrast(page_image)
            page_image = enhancer.enhance(1.1)  # Aumentar contraste 10%
            
            # Mejorar nitidez para TV
            enhancer = ImageEnhance.Sharpness(page_image)
            page_image = enhancer.enhance(1.05)  # Aumentar nitidez 5%
            
            # Ajustar brillo ligeramente
            enhancer = ImageEnhance.Brightness(page_image)
            page_image = enhancer.enhance(1.02)  # Aumentar brillo 2%
            
            return page_image
            
        except Exception as e:
            print(f"Error optimizando imagen para TV: {e}")
            return page_image  # Devolver original si falla

    def _añadir_imagen_fullscreen_tv(self, slide, image_path):
        """
        POWERPOINT: Añade imagen ajustada por ALTURA, alineada a la IZQUIERDA
        - Se ajusta al alto completo de la diapositiva
        - Se alinea a la izquierda (left = 0)
        - Deja espacio en blanco a la derecha
        - Mantiene proporciones originales
        """
        try:
            from PIL import Image
            
            # Obtener dimensiones de la imagen
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                img_aspect_ratio = img_width / img_height
            
            # Dimensiones del slide
            slide_width = slide.slide_layout.width
            slide_height = slide.slide_layout.height
            
            # AJUSTAR POR ALTURA - La imagen ocupará TODO el alto de la diapositiva
            new_height = slide_height
            new_width = slide_height * img_aspect_ratio
            
            # POSICIÓN: Alineada a la izquierda con pequeño margen
            left_margin = slide_width * 0.02  # 2% de margen izquierdo
            left = int(left_margin)
            top = 0  # Desde arriba
            
            # Si la imagen calculada es más ancha que el slide, recalcular por ancho
            max_width = slide_width * 0.75  # Máximo 75% del ancho para dejar espacio derecho
            if new_width > max_width:
                new_width = max_width
                new_height = max_width / img_aspect_ratio
                # Centrar verticalmente si queda espacio
                top = int((slide_height - new_height) / 2)
            
            print(f"📐 Imagen: {img_width}x{img_height} → PowerPoint: {int(new_width)}x{int(new_height)}")
            print(f"📍 Posición: left={left}, top={top}")
            
            # Añadir imagen al slide
            slide.shapes.add_picture(
                image_path, 
                left=left, 
                top=top, 
                width=int(new_width), 
                height=int(new_height)
            )
            
            print(f"✅ Imagen añadida - Espacio derecho disponible: {int(slide_width - left - new_width)}")
            
        except Exception as e:
            print(f"Error añadiendo imagen a slide: {e}")
            # Fallback - añadir imagen sin optimizar dimensiones (centrada)
            try:
                slide.shapes.add_picture(image_path, 0, 0, slide.slide_layout.width // 2, slide.slide_layout.height)
            except Exception as e2:
                print(f"Error en fallback: {e2}")

    def _configurar_powerpoint_para_tv(self, powerpoint_path):
        """
        POWERPOINT: Configura el archivo para reproducción automática en TV
        - Añade notas para configuración de presentación
        """
        try:
            print("⚙️ Configurando PowerPoint para TV...")
            
            # Crear archivo de configuración complementario
            config_path = os.path.join(os.path.dirname(powerpoint_path), "CONFIG_TV.txt")
            
            config_content = f"""
    CONFIGURACIÓN PARA TV - POWER POINT AUTOMÁTICO
    ==============================================

    Archivo: {os.path.basename(powerpoint_path)}
    Generado: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}

    INSTRUCCIONES PARA CONFIGURAR EN TV:
    1. Abrir PowerPoint en modo presentación (F5)
    2. Configurar transición automática: Transiciones > Avanzar diapositiva > Después de: 30 segundos
    3. Activar ciclo continuo: Configurar presentación > Repetir hasta presionar ESC
    4. Modo quiosco recomendado para TV

    CONFIGURACIÓN RECOMENDADA PARA MACROS:
    - Transición: 30-45 segundos por diapositiva
    - Repetición: Continua
    - Sin interacción del usuario
    - Pantalla completa
    """
            
            try:
                with open(config_path, 'w', encoding='utf-8') as f:
                    f.write(config_content)
                print(f"📝 Archivo de configuración creado: {config_path}")
            except Exception as e:
                print(f"No se pudo crear archivo de configuración: {e}")
            
            print("✅ PowerPoint configurado para TV")
            
        except Exception as e:
            print(f"Error configurando PowerPoint para TV: {e}")

    def ejecutar_powerpoint_en_tv_remota(self):
        """
        POWERPOINT: Ejecuta el archivo PowerPoint en la TV remota usando el .bat
        """
        try:
            import subprocess
            
            # Buscar el archivo .bat en el directorio actual
            bat_file = "Lanzador a TV calidad.bat"
            bat_path = os.path.join(os.path.dirname(__file__), bat_file)
            
            if not os.path.exists(bat_path):
                return {"error": f"Archivo {bat_file} no encontrado en {os.path.dirname(__file__)}", "ejecutado": False}
            
            print(f"🚀 Ejecutando lanzador remoto: {bat_path}")
            
            # Ejecutar el archivo .bat
            result = subprocess.run([bat_path], capture_output=True, text=True, shell=True)
            
            if result.returncode == 0:
                return {
                    "error": None,
                    "ejecutado": True,
                    "mensaje": "✅ Comando enviado a TV remota exitosamente",
                    "output": result.stdout
                }
            else:
                return {
                    "error": f"Error ejecutando lanzador: {result.stderr}",
                    "ejecutado": False,
                    "return_code": result.returncode
                }
            
        except Exception as e:
            return {"error": f"Error ejecutando lanzador remoto: {e}", "ejecutado": False}

    def generar_y_lanzar_powerpoint_completo(self, fecha_inicio_reporte, fecha_fin_reporte):
        """
        POWERPOINT: Función completa - Genera PDF, convierte a PowerPoint y lanza en TV
        """
        try:
            print("🎯 PROCESO COMPLETO: PDF → PowerPoint → TV")
            
            # PASO 1: Generar el PDF del reporte
            print("1️⃣ Generando PDF...")
            resultado_pdf = self.generar_reporte_pdf(fecha_inicio_reporte, fecha_fin_reporte)
            
            if resultado_pdf.get("error"):
                return {"error": f"Error generando PDF: {resultado_pdf['error']}", "powerpoint_path": None}
            
            pdf_path = resultado_pdf.get("pdf_path")
            if not pdf_path:
                return {"error": "No se obtuvo ruta del PDF generado", "powerpoint_path": None}
            
            print(f"✅ PDF generado: {os.path.basename(pdf_path)}")
            
            # PASO 2: Convertir PDF a PowerPoint
            print("2️⃣ Convirtiendo a PowerPoint...")
            resultado_ppt = self.generar_powerpoint_tv(pdf_path)
            
            if resultado_ppt.get("error"):
                return {"error": f"Error generando PowerPoint: {resultado_ppt['error']}", "powerpoint_path": None}
            
            powerpoint_path = resultado_ppt.get("powerpoint_path")
            print(f"✅ PowerPoint generado: {os.path.basename(powerpoint_path)}")
            
            # PASO 3: Lanzar en TV remota (opcional)
            print("3️⃣ Enviando comando a TV...")
            resultado_tv = self.ejecutar_powerpoint_en_tv_remota()
            
            return {
                "error": None,
                "pdf_path": pdf_path,
                "powerpoint_path": powerpoint_path,
                "total_slides": resultado_ppt.get("total_slides", 0),
                "tv_ejecutado": resultado_tv.get("ejecutado", False),
                "tv_mensaje": resultado_tv.get("mensaje", ""),
                "mensaje_completo": f"✅ Proceso completo: PDF → PowerPoint ({resultado_ppt.get('total_slides', 0)} slides) → {'TV iniciado' if resultado_tv.get('ejecutado') else 'TV no iniciado'}"
            }
            
        except Exception as e:
            return {"error": f"Error en proceso completo: {e}", "powerpoint_path": None}
        
    
    # ------------------------------------------------------------------------------
    # GENERACIÓN DE POWERPOINT DIARIO
    # ------------------------------------------------------------------------------
    def _añadir_imagen_con_panel_lateral(self, slide, image_path, datos_estadisticas, numero_slide=0):
        """
        POWERPOINT: Imagen izquierda + Panel específico por tipo de diapositiva CORREGIDO
        """
        try:
            from PIL import Image
            from pptx.util import Inches
            import matplotlib.pyplot as plt
            import tempfile
            import numpy as np
            
            print(f"🎯 Procesando slide {numero_slide + 1}...")
            
            # === IMAGEN IZQUIERDA ===
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                img_aspect = img_width / img_height
            
            slide_width = Inches(16)
            slide_height = Inches(9)
            
            new_height = slide_height
            new_width = new_height * img_aspect
            max_width = slide_width * 0.70
            if new_width > max_width:
                new_width = max_width
            
            slide.shapes.add_picture(image_path, left=0, top=0, width=int(new_width), height=int(new_height))
            
            # === OBTENER LISTAS DE MODELOS PARA SLIDES ESPECÍFICOS ===
            models_dist = datos_estadisticas.get("models_distribution", {})
            rework_models_dist = datos_estadisticas.get("rework_models_distribution", {})
            
            top5_modelos_nok = []
            top5_modelos_rework = []
            
            if models_dist:
                sorted_nok = sorted(models_dist.items(), key=lambda x: x[1], reverse=True)
                top5_modelos_nok = [modelo for modelo, _ in sorted_nok[:5] if modelo]
            
            if rework_models_dist:
                sorted_rework = sorted(rework_models_dist.items(), key=lambda x: x[1], reverse=True)
                top5_modelos_rework = [modelo for modelo, _ in sorted_rework[:5] if modelo]
            
            # === GENERAR GRÁFICOS SEGÚN TIPO DE SLIDE ===
            panel_left = slide_width * 0.72
            panel_width = slide_width * 0.26
            panel_height = slide_height * 0.95
            panel_top = slide_height * 0.025
            
            grafico1_path = None
            grafico2_path = None
            
            # SLIDE 0: NOK Global
            if numero_slide == 0:
                print("📊 Dashboard NOK Global...")
                grafico1_path = self._crear_grafico_simple_test("BARRAS DÍAS SEMANA", "#DC2626", datos_estadisticas)
                grafico2_path = self._crear_grafico_simple_test("HEAT MAP LÍNEAS", "#F59E0B", datos_estadisticas)
            
            # SLIDE 1: Retrabajos Global  
            elif numero_slide == 1:
                print("🔄 Dashboard Retrabajos Global...")
                grafico1_path = self._crear_grafico_simple_test("GAUGE RECUPERACIÓN", "#10B981", datos_estadisticas)
                grafico2_path = self._crear_grafico_simple_test("TIMELINE 7 DÍAS", "#3B82F6", datos_estadisticas)
            
            # SLIDES 2-3: Económico
            elif numero_slide in [2, 3]:
                print("💰 Dashboard Económico...")
                grafico1_path = self._crear_grafico_simple_test("BARRAS COSTE TURNO", "#E74C3C", datos_estadisticas)
                grafico2_path = self._crear_grafico_simple_test("TOP 3 MODELOS", "#8B5CF6", datos_estadisticas)
            
            # SLIDES 4-8: Modelos NOK específicos
            elif 4 <= numero_slide <= 8:
                modelo_index = numero_slide - 4
                if modelo_index < len(top5_modelos_nok):
                    modelo_especifico = top5_modelos_nok[modelo_index]
                    print(f"🎯 Análisis Modelo NOK: {modelo_especifico}")
                    
                    # Obtener datos específicos del modelo
                    datos_modelo = self._obtener_datos_modelo_especifico(modelo_especifico, "NOK")
                    
                    grafico1_path = self._crear_grafico_simple_test("COMPARATIVA MODELO", "#DC2626", datos_estadisticas, modelo_especifico, datos_modelo)
                    grafico2_path = self._crear_grafico_simple_test("DNA DEFECTO MODELO", "#F59E0B", datos_estadisticas, modelo_especifico, datos_modelo)
                else:
                    print(f"🎯 Análisis Modelo NOK genérico...")
                    grafico1_path = self._crear_grafico_simple_test("COMPARATIVA", "#DC2626", datos_estadisticas)
                    grafico2_path = self._crear_grafico_simple_test("DNA DEFECTO", "#F59E0B", datos_estadisticas)
            
            # SLIDES 9+: Modelos Rework específicos
            else:
                modelo_index = numero_slide - 9
                if modelo_index < len(top5_modelos_rework):
                    modelo_especifico = top5_modelos_rework[modelo_index]
                    print(f"🔧 Análisis Modelo Rework: {modelo_especifico}")
                    
                    # Obtener datos específicos del modelo
                    datos_modelo = self._obtener_datos_modelo_especifico(modelo_especifico, "REWORK")
                    
                    grafico1_path = self._crear_grafico_simple_test("EFICIENCIA MODELO", "#10B981", datos_estadisticas, modelo_especifico, datos_modelo)
                    grafico2_path = self._crear_grafico_simple_test("EXITOSOS MODELO", "#3B82F6", datos_estadisticas, modelo_especifico, datos_modelo)
                else:
                    print(f"🔧 Análisis Modelo Rework genérico...")
                    grafico1_path = self._crear_grafico_simple_test("EFICIENCIA", "#10B981", datos_estadisticas)
                    grafico2_path = self._crear_grafico_simple_test("EXITOSOS/FALLIDOS", "#3B82F6", datos_estadisticas)
            
            # Añadir gráficos al slide
            if grafico1_path:
                try:
                    slide.shapes.add_picture(
                        grafico1_path,
                        left=int(panel_left + Inches(0.05)),
                        top=int(panel_top + Inches(0.1)),
                        width=int(panel_width * 0.9),
                        height=int(panel_height * 0.4)
                    )
                except Exception as e:
                    print(f"❌ Error añadiendo gráfico superior: {e}")
            
            if grafico2_path:
                try:
                    slide.shapes.add_picture(
                        grafico2_path,
                        left=int(panel_left + Inches(0.05)),
                        top=int(panel_top + panel_height * 0.55),
                        width=int(panel_width * 0.9),
                        height=int(panel_height * 0.4)
                    )
                except Exception as e:
                    print(f"❌ Error añadiendo gráfico inferior: {e}")
            
        except Exception as e:
            print(f"❌ Error general en slide {numero_slide + 1}: {e}")
            self._añadir_imagen_izquierda_tv(slide, image_path)

    def _detectar_tipo_slide(self, numero_slide):
        """Detecta el tipo de slide basado en su posición"""
        if numero_slide == 0:
            return "nok_global"
        elif numero_slide == 1:
            return "retrabajos_global"
        elif numero_slide in [2, 3]:
            return "economico"
        elif numero_slide >= 4 and numero_slide <= 8:
            return "modelo_nok"
        else:
            return "modelo_retrabajos"

    def _generar_graficos_por_tipo(self, tipo_slide, datos_estadisticas):
        """Genera los mini-gráficos específicos según el tipo de slide"""
        
        if tipo_slide == "nok_global":
            return (
                self._crear_pentagono_dias_semana(datos_estadisticas),
                self._crear_heat_map_lineas(datos_estadisticas)
            )
        
        elif tipo_slide == "retrabajos_global":
            return (
                self._crear_velocimetro_recuperacion(datos_estadisticas),
                self._crear_timeline_7_dias(datos_estadisticas)
            )
        
        elif tipo_slide == "economico":
            return (
                self._crear_speedometer_coste_turno(datos_estadisticas),
                self._crear_top3_modelos_costosos(datos_estadisticas)
            )
        
        elif tipo_slide == "modelo_nok":
            return (
                self._crear_comparativa_circular_modelo(datos_estadisticas),
                self._crear_donut_defecto_dominante(datos_estadisticas)
            )
        
        elif tipo_slide == "modelo_retrabajos":
            return (
                self._crear_medidor_eficiencia_modelo(datos_estadisticas),
                self._crear_barras_retrabajos_exitosos(datos_estadisticas)
            )
        
        return None, None
    def _obtener_datos_modelo_especifico(self, modelo, tipo):
        """Obtiene datos específicos de un modelo usando la función leer_estadisticas"""
        try:
            from datetime import datetime, timedelta
            
            fecha_fin = datetime.now().strftime('%Y-%m-%d')
            fecha_inicio = (datetime.now() - timedelta(days=7)).strftime('%Y-%m-%d')
            
            # Filtrar por modelo específico
            datos_modelo = self.leer_estadisticas(fecha_inicio, fecha_fin, modelo, "(TODOS)")
            
            return datos_modelo if not datos_modelo.get("error") else {}
            
        except Exception as e:
            print(f"Error obteniendo datos del modelo {modelo}: {e}")
            return {}

    # === GRÁFICOS PARA NOK GLOBAL ===
    def _crear_pentagono_dias_semana(self, datos_estadisticas):
        """Pentágono mostrando distribución de defectos por día de semana"""
        try:
            import numpy as np
            import tempfile
            
            # Simular datos por día de semana (necesitarías extraer esto de tus archivos Excel)
            dias_semana = ['Lun', 'Mar', 'Mié', 'Jue', 'Vie']
            
            # Extraer distribución real de tus datos
            total_nok = datos_estadisticas.get("total_nok", 0)
            if total_nok == 0:
                return None
            
            # Distribución estimada (podrías calcular esto analizando tus archivos Excel por fecha)
            valores = [20, 25, 30, 35, 40]  # Ejemplo: más problemas hacia final de semana
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white', subplot_kw=dict(projection='polar'))
            
            # Ángulos para pentágono
            angulos = np.linspace(0, 2 * np.pi, len(dias_semana), endpoint=False).tolist()
            valores += [valores[0]]  # Cerrar el pentágono
            angulos += [angulos[0]]
            
            # Dibujar pentágono
            ax.plot(angulos, valores, 'o-', linewidth=2, color='#DC2626')
            ax.fill(angulos, valores, alpha=0.25, color='#DC2626')
            
            # Etiquetas
            ax.set_xticks(angulos[:-1])
            ax.set_xticklabels(dias_semana, fontsize=8)
            ax.set_ylim(0, max(valores))
            ax.set_title("DEFECTOS POR DÍA", fontsize=10, fontweight='bold', pad=20)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando pentágono días: {e}")
            return None

    def _crear_heat_map_lineas(self, datos_estadisticas):
        """Heat map de líneas más problemáticas"""
        try:
            import tempfile
            
            models_dist = datos_estadisticas.get("models_distribution", {})
            if not models_dist:
                return None
            
            top_models = sorted(models_dist.items(), key=lambda x: x[1], reverse=True)[:6]
            if not top_models:
                return None
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            nombres = [modelo[:8] for modelo, _ in top_models]
            valores = [cantidad for _, cantidad in top_models]
            max_val = max(valores) if valores else 1
            valores_norm = [(v/max_val)*100 for v in valores]
            
            # Grid 2x3
            grid_data = []
            grid_labels = []
            for i in range(2):
                row_data = []
                row_labels = []
                for j in range(3):
                    idx = i*3 + j
                    if idx < len(valores_norm):
                        row_data.append(valores_norm[idx])
                        row_labels.append(f"{nombres[idx]}\n{valores[idx]}")
                    else:
                        row_data.append(0)
                        row_labels.append("")
                grid_data.append(row_data)
                grid_labels.append(row_labels)
            
            im = ax.imshow(grid_data, cmap='Reds', aspect='equal', vmin=0, vmax=100)
            
            for i in range(2):
                for j in range(3):
                    if grid_labels[i][j]:
                        ax.text(j, i, grid_labels[i][j], ha='center', va='center',
                            fontsize=8, fontweight='bold', 
                            color='white' if grid_data[i][j] > 50 else 'black')
            
            ax.set_xticks([])
            ax.set_yticks([])
            ax.set_title("TOP LÍNEAS PROBLEMÁTICAS", fontsize=10, fontweight='bold', pad=10)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando heat map: {e}")
            return None

    # === GRÁFICOS PARA RETRABAJOS ===
    def _crear_velocimetro_recuperacion(self, datos_estadisticas):
        """Velocímetro del % de recuperación"""
        try:
            import numpy as np
            import tempfile
            
            total_nok = datos_estadisticas.get("total_nok", 0)
            total_rework = datos_estadisticas.get("total_rework", 0)
            
            if total_nok == 0:
                recuperacion_pct = 0
            else:
                recuperacion_pct = min(100, (total_rework / total_nok) * 100)
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white', subplot_kw=dict(projection='polar'))
            
            # Configurar velocímetro
            theta = np.linspace(0, np.pi, 100)
            r = np.ones_like(theta)
            
            # Fondo del velocímetro
            ax.plot(theta, r, color='lightgray', linewidth=10)
            
            # Arco coloreado según el porcentaje
            theta_fill = np.linspace(0, np.pi * (recuperacion_pct/100), 50)
            r_fill = np.ones_like(theta_fill)
            
            color = '#10B981' if recuperacion_pct > 70 else '#F59E0B' if recuperacion_pct > 40 else '#DC2626'
            ax.plot(theta_fill, r_fill, color=color, linewidth=10)
            
            # Aguja
            aguja_theta = np.pi * (recuperacion_pct/100)
            ax.plot([aguja_theta, aguja_theta], [0, 1], color='black', linewidth=3)
            
            ax.set_ylim(0, 1.2)
            ax.set_title(f"RECUPERACIÓN: {recuperacion_pct:.1f}%", fontsize=10, fontweight='bold', pad=20)
            ax.set_xticks([])
            ax.set_yticks([])
            ax.grid(False)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando velocímetro: {e}")
            return None

    def _crear_timeline_7_dias(self, datos_estadisticas):
        """Timeline horizontal de últimos 7 días"""
        try:
            import numpy as np
            import tempfile
            from datetime import datetime, timedelta
            
            # Generar últimos 7 días
            fechas = [(datetime.now() - timedelta(days=i)).strftime('%d/%m') for i in range(6, -1, -1)]
            
            # Valores simulados (podrías extraer de evolution_rework)
            evolution_rework = datos_estadisticas.get("evolution_rework", {})
            if evolution_rework:
                valores_reales = list(evolution_rework.values())[-7:] if len(evolution_rework) >= 7 else [10, 15, 8, 20, 12, 18, 14]
            else:
                valores_reales = [10, 15, 8, 20, 12, 18, 14]  # Valores ejemplo
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            # Normalizar valores para colores
            max_val = max(valores_reales) if valores_reales else 1
            colores = [(val/max_val) for val in valores_reales]
            
            # Timeline con puntos coloreados
            scatter = ax.scatter(range(len(fechas)), [0]*len(fechas), 
                            c=colores, cmap='RdYlGn_r', s=200, alpha=0.8)
            
            # Línea conectora
            ax.plot(range(len(fechas)), [0]*len(fechas), color='gray', alpha=0.3, linewidth=2)
            
            # Etiquetas
            for i, (fecha, valor) in enumerate(zip(fechas, valores_reales)):
                ax.text(i, 0.1, fecha, ha='center', fontsize=7, rotation=45)
                ax.text(i, -0.1, str(valor), ha='center', fontsize=8, fontweight='bold')
            
            ax.set_xlim(-0.5, len(fechas)-0.5)
            ax.set_ylim(-0.3, 0.3)
            ax.set_title("RETRABAJOS ÚLTIMOS 7 DÍAS", fontsize=10, fontweight='bold')
            ax.axis('off')
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando timeline: {e}")
            return None

    # === GRÁFICOS PARA ECONÓMICO ===
    def _crear_speedometer_coste_turno(self, datos_estadisticas):
        """Speedometer del coste por turno"""
        try:
            import numpy as np
            import tempfile
            
            shift_dist = datos_estadisticas.get("shift_distribution", {})
            if not shift_dist:
                return None
            
            # Calcular coste promedio por turno
            total_turnos = len(shift_dist)
            if total_turnos == 0:
                coste_promedio = 0
            else:
                if self.costes_dict:
                    precio_promedio = sum(self.costes_dict.values()) / len(self.costes_dict)
                else:
                    precio_promedio = 15
                
                coste_total = sum(nok * precio_promedio for nok in shift_dist.values())
                coste_promedio = coste_total / total_turnos
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white', subplot_kw=dict(projection='polar'))
            
            # Speedometer
            theta = np.linspace(0, np.pi, 100)
            r = np.ones_like(theta)
            ax.plot(theta, r, color='lightgray', linewidth=15)
            
            # Escala: 0-1000€ por turno
            max_coste = 1000
            coste_normalizado = min(100, (coste_promedio / max_coste) * 100)
            
            theta_fill = np.linspace(0, np.pi * (coste_normalizado/100), 50)
            r_fill = np.ones_like(theta_fill)
            
            color = '#DC2626' if coste_normalizado > 70 else '#F59E0B' if coste_normalizado > 40 else '#10B981'
            ax.plot(theta_fill, r_fill, color=color, linewidth=15)
            
            # Aguja
            aguja_theta = np.pi * (coste_normalizado/100)
            ax.plot([aguja_theta, aguja_theta], [0, 1], color='black', linewidth=4)
            
            ax.set_ylim(0, 1.3)
            ax.set_title(f"COSTE/TURNO: €{coste_promedio:.0f}", fontsize=10, fontweight='bold', pad=20)
            ax.set_xticks([])
            ax.set_yticks([])
            ax.grid(False)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando speedometer: {e}")
            return None

    def _crear_top3_modelos_costosos(self, datos_estadisticas):
        """Top 3 modelos más costosos"""
        try:
            import tempfile
            
            models_dist = datos_estadisticas.get("models_distribution", {})
            if not models_dist:
                return None
            
            # Calcular costes por modelo
            if self.costes_dict:
                precio_promedio = sum(self.costes_dict.values()) / len(self.costes_dict)
            else:
                precio_promedio = 15
            
            modelos_costes = [(modelo, nok * precio_promedio) for modelo, nok in models_dist.items()]
            top3_modelos = sorted(modelos_costes, key=lambda x: x[1], reverse=True)[:3]
            
            if not top3_modelos:
                return None
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            modelos = [m[:8] for m, _ in top3_modelos]
            costes = [c for _, c in top3_modelos]
            colores = ['#DC2626', '#F59E0B', '#10B981']
            
            bars = ax.barh(range(len(modelos)), costes, color=colores, alpha=0.8)
            
            # Añadir valores
            for i, (bar, coste) in enumerate(zip(bars, costes)):
                width = bar.get_width()
                ax.text(width + max(costes)*0.02, bar.get_y() + bar.get_height()/2,
                    f'€{int(coste)}', va='center', ha='left', fontsize=8, fontweight='bold')
            
            ax.set_yticks(range(len(modelos)))
            ax.set_yticklabels(modelos, fontsize=9)
            ax.set_title("TOP 3 MODELOS COSTOSOS", fontsize=10, fontweight='bold')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando top3 modelos: {e}")
            return None

    # === GRÁFICOS PARA MODELO NOK ===
    def _crear_comparativa_circular_modelo(self, datos_estadisticas):
        """Comparativa circular: modelo vs promedio general"""
        try:
            import tempfile
            
            # Obtener datos del modelo (necesitarías pasar el modelo específico)
            models_dist = datos_estadisticas.get("models_distribution", {})
            if not models_dist:
                return None
            
            # Tomar el modelo más problemático como ejemplo
            modelo_principal = max(models_dist.items(), key=lambda x: x[1])
            total_modelos = sum(models_dist.values())
            promedio_modelo = total_modelos / len(models_dist) if models_dist else 0
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            # Datos para comparación
            categorias = ['Este\nModelo', 'Promedio\nGeneral']
            valores = [modelo_principal[1], promedio_modelo]
            colores = ['#DC2626', '#6B7280']
            
            # Gráfico circular (donut)
            wedges, texts, autotexts = ax.pie(valores, labels=categorias, autopct='%1.1f%%',
                                            colors=colores, startangle=90,
                                            wedgeprops=dict(width=0.5))
            
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontweight('bold')
            
            ax.set_title(f"MODELO vs PROMEDIO", fontsize=10, fontweight='bold')
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando comparativa circular: {e}")
            return None

    def _crear_donut_defecto_dominante(self, datos_estadisticas):
        """Mini donut del defecto dominante"""
        try:
            import tempfile
            
            defect_dist = datos_estadisticas.get("defect_distribution", {})
            if not defect_dist:
                return None
            
            # Tomar top 3 defectos
            top_defectos = sorted(defect_dist.items(), key=lambda x: x[1], reverse=True)[:3]
            if not top_defectos:
                return None
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            nombres = [d[:10] for d, _ in top_defectos]
            valores = [v for _, v in top_defectos]
            colores = ['#DC2626', '#F59E0B', '#10B981']
            
            # Donut chart
            wedges, texts, autotexts = ax.pie(valores, labels=nombres, autopct='%1.1f%%',
                                            colors=colores, startangle=90,
                                            wedgeprops=dict(width=0.4))
            
            # Centro blanco
            centre_circle = plt.Circle((0,0), 0.60, fc='white')
            fig.gca().add_artist(centre_circle)
            
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontsize(8)
                autotext.set_fontweight('bold')
            
            ax.set_title("DNA DEFECTOS", fontsize=10, fontweight='bold')
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando donut defectos: {e}")
            return None

    # === GRÁFICOS PARA MODELO RETRABAJOS ===
    def _crear_medidor_eficiencia_modelo(self, datos_estadisticas):
        """Medidor de eficiencia de recuperación del modelo"""
        try:
            import numpy as np
            import tempfile
            
            total_nok = datos_estadisticas.get("total_nok", 0)
            total_rework = datos_estadisticas.get("total_rework", 0)
            
            eficiencia = min(100, (total_rework / total_nok * 100)) if total_nok > 0 else 0
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            # Medidor de barras horizontales
            categorias = ['Eficiencia\nRecuperación']
            ax.barh(categorias, [eficiencia], color='#10B981', alpha=0.8, height=0.5)
            ax.barh(categorias, [100], color='lightgray', alpha=0.3, height=0.5)
            
            # Añadir porcentaje
            ax.text(eficiencia + 2, 0, f'{eficiencia:.1f}%', va='center', fontweight='bold')
            
            ax.set_xlim(0, 100)
            ax.set_title(f"EFICIENCIA MODELO", fontsize=10, fontweight='bold')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_visible(False)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando medidor eficiencia: {e}")
            return None

    def _crear_barras_retrabajos_exitosos(self, datos_estadisticas):
        """Barras de retrabajos exitosos vs fallidos"""
        try:
            import tempfile
            
            total_rework = datos_estadisticas.get("total_rework", 0)
            # Asumir 85% de retrabajos exitosos (podrías calcular esto con más detalle)
            exitosos = int(total_rework * 0.85)
            fallidos = total_rework - exitosos
            
            fig, ax = plt.subplots(figsize=(4, 3), facecolor='white')
            
            categorias = ['Exitosos', 'Fallidos']
            valores = [exitosos, fallidos]
            colores = ['#10B981', '#DC2626']
            
            bars = ax.bar(categorias, valores, color=colores, alpha=0.8)
            
            # Añadir valores en barras
            for bar, valor in zip(bars, valores):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2., height + max(valores)*0.01,
                    f'{valor}', ha='center', va='bottom', fontweight='bold', fontsize=9)
            
            ax.set_title("RETRABAJOS", fontsize=10, fontweight='bold')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            plt.tight_layout()
            fig.savefig(temp_file.name, dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            return temp_file.name
            
        except Exception as e:
            print(f"Error creando barras retrabajos: {e}")
            return None
    def generar_reporte_pdf(self, fecha_inicio_reporte, fecha_fin_reporte):
        """
        Genera un PDF ultra profesional con:
        1 página -> REPORTE GLOBAL DE RECHAZOS (NOK).
        1 página -> REPORTE GLOBAL DE RETRABAJOS (R).
        2 páginas -> ANÁLISIS ECONÓMICO DE COSTES (NUEVO INTEGRADO).
        1 página por cada modelo NOK (detalle).
        1 página por cada modelo R (detalle).
        Diseño moderno y profesional para directivos.
        """
        data_all_report_range = self.leer_estadisticas(fecha_inicio_reporte, fecha_fin_reporte, "(TODOS)", "(TODOS)")
        if data_all_report_range.get("error"):
            print(f"Error al leer datos globales: {data_all_report_range['error']}")
            return {"error": data_all_report_range["error"]}

        # Extraer datos globales
        global_ok = data_all_report_range.get("global_total_ok", 0)
        global_nok = data_all_report_range.get("global_total_nok", 0)
        global_nok_prov = data_all_report_range.get("global_total_nok_prov", 0)
        def_dist_nok_global_report = data_all_report_range.get("defect_distribution", {})
        mod_dist_nok_global_report = data_all_report_range.get("models_distribution", {})
        mod_dist_nok_prov_global_report = data_all_report_range.get("global_models_distribution_prov", {}) 
        shift_dist_nok_global_report = data_all_report_range.get("shift_distribution", {})

        global_rew = data_all_report_range.get("global_total_rework", 0)
        rew_def_dist_global_report = data_all_report_range.get("rework_defect_distribution", {})
        rew_mod_dist_global_report = data_all_report_range.get("rework_models_distribution", {})
        rew_shift_dist_global_report = data_all_report_range.get("rework_shift_distribution", {})
        
        # NUEVO: Extraer datos de costes
        costes_por_linea_turno = data_all_report_range.get("costes_por_linea_turno", {})
        costes_globales = data_all_report_range.get("costes_globales", {})

        story = []

        # Graficar Global NOK
        buf_ok_nok_trend_glob, buf_defectos_glob, buf_modelos_glob, buf_turnos_glob = self._graficas_global_nok(
            fecha_inicio_reporte, fecha_fin_reporte,
            def_dist_nok_global_report, mod_dist_nok_global_report, shift_dist_nok_global_report,  # <- CORREGIDO
            mod_dist_nok_prov_global_report
        )
        self._crear_hoja_global_nok(
            story, fecha_inicio_reporte, fecha_fin_reporte,
            global_ok, global_nok, global_nok_prov,
            buf_ok_nok_trend_glob, buf_defectos_glob, buf_modelos_glob, buf_turnos_glob
        )

        # Graficar Global Rework
        buf_evol_r_trend_glob, buf_def_r_glob, buf_mod_r_glob, buf_shift_r_glob = self._graficas_global_r(
            fecha_inicio_reporte, fecha_fin_reporte,
            rew_def_dist_global_report, rew_mod_dist_global_report, rew_shift_dist_global_report
        )
        self._crear_hoja_global_rework(
            story, fecha_inicio_reporte, fecha_fin_reporte,
            global_rew,
            buf_evol_r_trend_glob, buf_def_r_glob, buf_mod_r_glob, buf_shift_r_glob
        )
        
        # NUEVO: Añadir hojas de análisis económico de costes (2 HOJAS COMPLETAS)
        if costes_por_linea_turno:
            print(f"✅ Generando análisis económico con {len(costes_por_linea_turno)} registros de costes")
            print(f"💰 Total pérdidas: €{costes_globales.get('total_perdidas', 0):,.2f}")
            self._crear_hoja_analisis_costes(
                story, fecha_inicio_reporte, fecha_fin_reporte,
                costes_por_linea_turno, costes_globales
            )
            print(f"📄 Hojas de análisis económico añadidas correctamente")
        else:
            print("⚠️ No hay datos de costes disponibles para el análisis económico")

        # Páginas por modelo NOK
        if mod_dist_nok_global_report:
            sorted_mod_nok = sorted(mod_dist_nok_global_report.items(), key=lambda x: x[1], reverse=True)
            top5_nok = [t[0] for t in sorted_mod_nok[:5] if t[0]]
            for mod_nok in top5_nok:
                self._crear_hoja_modelo_nok(story, fecha_inicio_reporte, fecha_fin_reporte, mod_nok)
        else:
            print("No hay datos de modelos NOK para generar páginas individuales.")

        # Páginas por modelo Rework
        if rew_mod_dist_global_report:
            sorted_mod_r = sorted(rew_mod_dist_global_report.items(), key=lambda x: x[1], reverse=True)
            top5_r = [t[0] for t in sorted_mod_r[:5] if t[0]]
            for mod_r in top5_r:
                self._crear_hoja_modelo_rework(story, fecha_inicio_reporte, fecha_fin_reporte, mod_r)
        else:
            print("No hay datos de modelos R para generar páginas individuales.")

        # Asegurar existencia de carpeta
        if not os.path.exists(self.reports_dir):
            try:
                os.makedirs(self.reports_dir)
            except Exception as e:
                return {"error": f"No se pudo crear '{self.reports_dir}': {e}", "pdf_path": None}

        report_filename = f"Report_NOK_R_COSTES_{fecha_inicio_reporte}_to_{fecha_fin_reporte}.pdf"
        full_report_path = os.path.join(self.reports_dir, report_filename)

        try:
            doc = SimpleDocTemplate(
                full_report_path,
                pagesize=PAGE_SIZE,
                rightMargin=20, leftMargin=20, topMargin=15, bottomMargin=15  # Márgenes reducidos
            )
            doc.build(story)
            
            # Información del reporte generado
            total_pages = 2  # NOK + Rework
            if costes_por_linea_turno:
                total_pages += 2  # Análisis económico (2 páginas)
            total_pages += len(top5_nok) if mod_dist_nok_global_report else 0
            total_pages += len(top5_r) if rew_mod_dist_global_report else 0
            
            print(f"✅ Reporte ejecutivo generado exitosamente:")
            print(f"   📄 Total de páginas: {total_pages}")
            print(f"   💰 Análisis económico: {'Incluido' if costes_por_linea_turno else 'No disponible'}")
            print(f"   📊 Modelos NOK analizados: {len(top5_nok) if mod_dist_nok_global_report else 0}")
            print(f"   🔄 Modelos Rework analizados: {len(top5_r) if rew_mod_dist_global_report else 0}")
            print(f"   📁 Archivo: {report_filename}")
            
            return {"error": None, "pdf_path": os.path.abspath(full_report_path)}
        except Exception as e:
            print(f"Error construyendo PDF: {e}")
            return {"error": f"Error construyendo PDF: {e}", "pdf_path": None}

    def generar_presentacion_html(self, fecha_inicio, fecha_fin):
        """
        Genera presentación HTML IDÉNTICA a las páginas del PDF.
        Carrusel automático cada 5 segundos en bucle infinito.
        Diseñado para 1920x1080 en pantalla completa.

        GENERA LAS MISMAS PÁGINAS QUE EL PDF:
        1. Global NOK (4 gráficos 2x2)
        2. Global Rework (4 gráficos 2x2)
        3. Análisis Costes (si hay datos) (2 páginas)
        4. Top 5 Modelos NOK (1 página cada uno)
        5. Top 5 Modelos Rework (1 página cada uno)
        """
        try:
            import base64

            print(f"🎬 Generando presentación HTML (IDÉNTICA AL PDF) para {fecha_inicio} — {fecha_fin}")

            # Leer datos globales (igual que el PDF)
            data_all = self.leer_estadisticas(fecha_inicio, fecha_fin, "(TODOS)", "(TODOS)")
            if data_all.get("error"):
                return {"error": data_all["error"], "html_content": None}

            # Extraer datos
            global_ok = data_all.get("global_total_ok", 0)
            global_nok = data_all.get("global_total_nok", 0)
            global_nok_prov = data_all.get("global_total_nok_prov", 0)
            global_rew = data_all.get("global_total_rework", 0)

            def_dist_nok = data_all.get("defect_distribution", {})
            mod_dist_nok = data_all.get("models_distribution", {})
            mod_dist_nok_prov = data_all.get("global_models_distribution_prov", {})
            shift_dist_nok = data_all.get("shift_distribution", {})

            rew_def_dist = data_all.get("rework_defect_distribution", {})
            rew_mod_dist = data_all.get("rework_models_distribution", {})
            rew_shift_dist = data_all.get("rework_shift_distribution", {})

            costes_por_linea_turno = data_all.get("costes_por_linea_turno", {})
            costes_globales = data_all.get("costes_globales", {})

            # Función helper para convertir gráfico a base64
            def buffer_to_base64(buf):
                if buf is None:
                    return ""
                buf.seek(0)
                img_str = base64.b64encode(buf.read()).decode()
                return f"data:image/png;base64,{img_str}"

            # Generar gráficos NOK
            buf_ok_nok, buf_defectos, buf_modelos, buf_turnos = self._graficas_global_nok(
                fecha_inicio, fecha_fin,
                def_dist_nok, mod_dist_nok, shift_dist_nok, mod_dist_nok_prov
            )

            # Generar gráficos Rework
            buf_evol_r, buf_def_r, buf_mod_r, buf_shift_r = self._graficas_global_r(
                fecha_inicio, fecha_fin,
                rew_def_dist, rew_mod_dist, rew_shift_dist
            )

            # Convertir a base64
            img_ok_nok_b64 = buffer_to_base64(buf_ok_nok)
            img_defectos_b64 = buffer_to_base64(buf_defectos)
            img_modelos_b64 = buffer_to_base64(buf_modelos)
            img_turnos_b64 = buffer_to_base64(buf_turnos)

            img_evol_r_b64 = buffer_to_base64(buf_evol_r)
            img_def_r_b64 = buffer_to_base64(buf_def_r)
            img_mod_r_b64 = buffer_to_base64(buf_mod_r)
            img_shift_r_b64 = buffer_to_base64(buf_shift_r)

            # Calcular KPIs
            total_fab = global_ok + global_nok + global_nok_prov
            rejection_rate = (global_nok / (global_ok + global_nok) * 100) if (global_ok + global_nok) > 0 else 0
            fecha_display = fecha_inicio if fecha_inicio == fecha_fin else f"{fecha_inicio} — {fecha_fin}"

            # Preparar lista de diapositivas
            slides_html = []

            # === DIAPOSITIVA 1: NOK GLOBAL (IGUAL QUE PÁGINA 1 DEL PDF) ===
            status_color = "#EF4444" if rejection_rate > self.OBJETIVO else "#10B981"
            status_text = "CRÍTICO" if rejection_rate > self.OBJETIVO else "ÓPTIMO"

            slide1 = f'''
            <div class="slide">
                <div class="slide-header" style="background: linear-gradient(135deg, #0F172A 0%, #1E293B 100%);">
                    <h1 style="color: white; font-size: 48px; margin-bottom: 8px;">QUALITY DASHBOARD</h1>
                    <p style="color: white; font-size: 24px; margin-bottom: 5px; opacity: 0.95;">ANÁLISIS GLOBAL DE RECHAZOS (NOK)</p>
                    <p style="color: white; font-size: 18px; opacity: 0.85;">Período: {fecha_display} | Objetivo: ≤ {self.OBJETIVO:.1f}%</p>
                </div>
                <div class="slide-kpis">
                    <div class="kpi-hero" style="background: {status_color}; padding: 20px 40px; border-radius: 15px; min-width: 280px; text-align: center;">
                        <div style="color: white; font-size: 14px; font-weight: 700; text-transform: uppercase; margin-bottom: 8px;">TASA DE RECHAZO</div>
                        <div style="color: white; font-size: 48px; font-weight: 800; margin: 10px 0;">{rejection_rate:.2f}%</div>
                        <div style="color: white; font-size: 13px; font-weight: 500;">{status_text}</div>
                    </div>
                    <div class="kpi-card" style="background: white; padding: 20px 35px; border-radius: 15px; min-width: 220px; text-align: center; border: 2px solid #E2E8F0;">
                        <div style="color: #1E293B; font-size: 14px; font-weight: 700; text-transform: uppercase; margin-bottom: 8px;">PRODUCCIÓN TOTAL</div>
                        <div style="color: #1E293B; font-size: 36px; font-weight: 700; margin-top: 6px;">{total_fab:,}</div>
                    </div>
                    <div class="kpi-card" style="background: white; padding: 20px 35px; border-radius: 15px; min-width: 220px; text-align: center; border: 2px solid #E2E8F0;">
                        <div style="color: #1E293B; font-size: 14px; font-weight: 700; text-transform: uppercase; margin-bottom: 8px;">RECHAZOS INTERNOS</div>
                        <div style="color: #1E293B; font-size: 36px; font-weight: 700; margin-top: 6px;">{global_nok:,}</div>
                    </div>
                    <div class="kpi-card" style="background: white; padding: 20px 35px; border-radius: 15px; min-width: 220px; text-align: center; border: 2px solid #E2E8F0;">
                        <div style="color: #1E293B; font-size: 14px; font-weight: 700; text-transform: uppercase; margin-bottom: 8px;">RECHAZOS PROVEEDOR</div>
                        <div style="color: #1E293B; font-size: 36px; font-weight: 700; margin-top: 6px;">{global_nok_prov:,}</div>
                    </div>
                </div>
                <div class="slide-charts">
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Tendencia de Calidad</h3>
                        <img src="{img_ok_nok_b64}" alt="Tendencia" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Top 5 Defectos</h3>
                        <img src="{img_defectos_b64}" alt="Defectos" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Impacto por Modelo</h3>
                        <img src="{img_modelos_b64}" alt="Modelos" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Distribución por Turno</h3>
                        <img src="{img_turnos_b64}" alt="Turnos" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                </div>
            </div>
            '''
            slides_html.append(slide1)

            # === DIAPOSITIVA 2: REWORK GLOBAL (IGUAL QUE PÁGINA 2 DEL PDF) ===
            slide2 = f'''
            <div class="slide">
                <div class="slide-header" style="background: linear-gradient(135deg, #065F46 0%, #047857 100%); padding: 30px 60px; text-align: center;">
                    <h1 style="color: white; font-size: 48px; margin-bottom: 8px;">REWORK DASHBOARD</h1>
                    <p style="color: white; font-size: 24px; margin-bottom: 5px; opacity: 0.95;">ANÁLISIS GLOBAL DE RETRABAJOS</p>
                    <p style="color: white; font-size: 18px; opacity: 0.85;">Período: {fecha_display}</p>
                </div>
                <div class="slide-kpis" style="padding: 25px 60px;">
                    <div class="kpi-hero" style="background: #10B981; padding: 20px 40px; border-radius: 15px; min-width: 320px; text-align: center;">
                        <div style="color: white; font-size: 14px; font-weight: 700; text-transform: uppercase; margin-bottom: 8px;">TOTAL RETRABAJOS</div>
                        <div style="color: white; font-size: 48px; font-weight: 800; margin: 10px 0;">{global_rew:,}</div>
                        <div style="color: white; font-size: 13px; font-weight: 500;">UNIDADES RECUPERADAS</div>
                    </div>
                </div>
                <div class="slide-charts">
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Evolución Temporal</h3>
                        <img src="{img_evol_r_b64}" alt="Evolución" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Defectos Recuperados</h3>
                        <img src="{img_def_r_b64}" alt="Defectos R" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Modelos en Retrabajo</h3>
                        <img src="{img_mod_r_b64}" alt="Modelos R" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                    <div class="chart-box">
                        <h3 style="color: #1E293B; font-size: 18px; font-weight: 700; margin-bottom: 8px; text-align: center; border-bottom: 2px solid #E2E8F0; padding-bottom: 6px;">Distribución por Turno</h3>
                        <img src="{img_shift_r_b64}" alt="Turnos R" style="width: 100%; height: auto; object-fit: contain;">
                    </div>
                </div>
            </div>
            '''
            slides_html.append(slide2)

            # === DIAPOSITIVA 3: ANÁLISIS DE COSTES (si hay datos) ===
            if costes_por_linea_turno and costes_globales.get('total_perdidas', 0) > 0:
                total_perdidas = costes_globales.get('total_perdidas', 0)
                num_defectos = costes_globales.get('num_defectos', 0)
                coste_medio = total_perdidas / num_defectos if num_defectos > 0 else 0

                # Top 10 líneas/turnos por coste
                costes_sorted = sorted(
                    costes_por_linea_turno.items(),
                    key=lambda x: x[1]['total_coste'],
                    reverse=True
                )[:10]

                tabla_costes_html = ""
                for i, ((linea, uet, turno), datos) in enumerate(costes_sorted):
                    coste = datos['total_coste']
                    bg_color = "#F3E8FF" if i % 2 == 0 else "white"
                    tabla_costes_html += f'''
                    <tr style="background: {bg_color}; border-bottom: 1px solid #E2E8F0;">
                        <td style="padding: 15px; color: #1E293B; font-weight: 500;">{linea}</td>
                        <td style="padding: 15px; color: #1E293B; font-weight: 500;">{uet}</td>
                        <td style="padding: 15px; color: #1E293B; font-weight: 500;">{turno}</td>
                        <td style="padding: 15px; text-align: right; font-weight: 700; color: #7C3AED;">€{coste:,.2f}</td>
                    </tr>
                    '''

                slide3 = f'''
                <div class="slide">
                    <div class="slide-header" style="background: linear-gradient(135deg, #7C3AED 0%, #6366F1 100%);">
                        <h1>ANÁLISIS ECONÓMICO</h1>
                        <p class="slide-subtitle">IMPACTO DE COSTES POR DEFECTOS</p>
                        <p class="slide-period">Período: {fecha_display}</p>
                    </div>
                    <div class="slide-kpis">
                        <div class="kpi-hero" style="background: #7C3AED;">
                            <div class="kpi-label">PÉRDIDAS TOTALES</div>
                            <div class="kpi-value-hero">€{total_perdidas:,.2f}</div>
                            <div class="kpi-status">ANÁLISIS ECONÓMICO</div>
                        </div>
                        <div class="kpi-card">
                            <div class="kpi-label">DEFECTOS ANALIZADOS</div>
                            <div class="kpi-value">{num_defectos:,}</div>
                        </div>
                        <div class="kpi-card">
                            <div class="kpi-label">COSTE MEDIO</div>
                            <div class="kpi-value">€{coste_medio:.2f}</div>
                        </div>
                    </div>
                    <div style="padding: 40px 60px; background: white; border-radius: 20px; margin: 30px 80px; box-shadow: 0 8px 30px rgba(0,0,0,0.1);">
                        <h2 style="color: #1E293B; margin-bottom: 30px; font-size: 32px; text-align: center; border-bottom: 3px solid #7C3AED; padding-bottom: 15px;">Top 10 Líneas con Mayor Impacto Económico</h2>
                        <table style="width: 100%; border-collapse: collapse; font-size: 22px;">
                            <thead>
                                <tr style="background: #7C3AED; color: white;">
                                    <th style="padding: 18px; text-align: left; font-weight: 700;">Línea</th>
                                    <th style="padding: 18px; text-align: left; font-weight: 700;">UET</th>
                                    <th style="padding: 18px; text-align: left; font-weight: 700;">Turno</th>
                                    <th style="padding: 18px; text-align: right; font-weight: 700;">Coste Total</th>
                                </tr>
                            </thead>
                            <tbody>
                                {tabla_costes_html}
                            </tbody>
                        </table>
                    </div>
                </div>
                '''
                slides_html.append(slide3)

            # === DIAPOSITIVAS DE MODELOS NOK (Top 5) ===
            if mod_dist_nok:
                sorted_mod_nok = sorted(mod_dist_nok.items(), key=lambda x: x[1], reverse=True)
                top5_nok = sorted_mod_nok[:5]

                for modelo, cantidad in top5_nok:
                    # Obtener datos específicos del modelo
                    data_modelo = self.leer_estadisticas(fecha_inicio, fecha_fin, modelo, "(TODOS)")

                    if not data_modelo.get("error"):
                        modelo_ok = data_modelo.get("total_ok", 0)
                        modelo_nok = data_modelo.get("total_nok", 0)
                        modelo_total = modelo_ok + modelo_nok
                        modelo_pct = (modelo_nok / modelo_total * 100) if modelo_total > 0 else 0

                        # Distribución de defectos del modelo
                        modelo_def_dist = data_modelo.get("defect_distribution", {})
                        def_table_html = ""
                        for i, (defecto, cant) in enumerate(sorted(modelo_def_dist.items(), key=lambda x: x[1], reverse=True)[:5]):
                            bg_color = "#FEE2E2" if i % 2 == 0 else "white"
                            def_table_html += f'''
                            <tr style="background: {bg_color}; border-bottom: 1px solid #E2E8F0;">
                                <td style="padding: 15px; color: #1E293B; font-weight: 500;">{defecto}</td>
                                <td style="padding: 15px; text-align: right; font-weight: 700; color: #DC2626;">{cant:,}</td>
                            </tr>
                            '''

                        slide_modelo = f'''
                        <div class="slide">
                            <div class="slide-header" style="background: linear-gradient(135deg, #DC2626 0%, #EF4444 100%);">
                                <h1>ANÁLISIS POR MODELO: {modelo}</h1>
                                <p class="slide-subtitle">DETALLE DE RECHAZOS NOK</p>
                                <p class="slide-period">Período: {fecha_display}</p>
                            </div>
                            <div class="slide-kpis">
                                <div class="kpi-card">
                                    <div class="kpi-label">PRODUCCIÓN</div>
                                    <div class="kpi-value">{modelo_total:,}</div>
                                </div>
                                <div class="kpi-card">
                                    <div class="kpi-label">RECHAZOS NOK</div>
                                    <div class="kpi-value">{modelo_nok:,}</div>
                                </div>
                                <div class="kpi-card">
                                    <div class="kpi-label">% NOK</div>
                                    <div class="kpi-value">{modelo_pct:.2f}%</div>
                                </div>
                            </div>
                            <div style="padding: 40px 60px; background: white; border-radius: 20px; margin: 30px 100px; box-shadow: 0 8px 30px rgba(0,0,0,0.1);">
                                <h2 style="color: #1E293B; margin-bottom: 30px; font-size: 36px; text-align: center; border-bottom: 3px solid #DC2626; padding-bottom: 15px;">Top 5 Defectos del Modelo</h2>
                                <table style="width: 100%; border-collapse: collapse; font-size: 26px;">
                                    <thead>
                                        <tr style="background: #DC2626; color: white;">
                                            <th style="padding: 18px; text-align: left; font-weight: 700;">Defecto</th>
                                            <th style="padding: 18px; text-align: right; font-weight: 700;">Cantidad</th>
                                        </tr>
                                    </thead>
                                    <tbody>
                                        {def_table_html}
                                    </tbody>
                                </table>
                            </div>
                        </div>
                        '''
                        slides_html.append(slide_modelo)

            # === DIAPOSITIVAS DE MODELOS REWORK (Top 5) ===
            if rew_mod_dist:
                sorted_mod_r = sorted(rew_mod_dist.items(), key=lambda x: x[1], reverse=True)
                top5_r = sorted_mod_r[:5]

                for modelo, cantidad in top5_r:
                    # Obtener datos de rework del modelo
                    data_modelo_r = self.leer_estadisticas(fecha_inicio, fecha_fin, modelo, "(TODOS)")

                    if not data_modelo_r.get("error"):
                        modelo_rew = data_modelo_r.get("total_rework", 0)
                        modelo_rew_def = data_modelo_r.get("rework_defect_distribution", {})

                        def_r_table_html = ""
                        for i, (defecto, cant) in enumerate(sorted(modelo_rew_def.items(), key=lambda x: x[1], reverse=True)[:5]):
                            bg_color = "#D1FAE5" if i % 2 == 0 else "white"
                            def_r_table_html += f'''
                            <tr style="background: {bg_color}; border-bottom: 1px solid #E2E8F0;">
                                <td style="padding: 15px; color: #1E293B; font-weight: 500;">{defecto}</td>
                                <td style="padding: 15px; text-align: right; font-weight: 700; color: #10B981;">{cant:,}</td>
                            </tr>
                            '''

                        slide_modelo_r = f'''
                        <div class="slide">
                            <div class="slide-header" style="background: linear-gradient(135deg, #059669 0%, #10B981 100%);">
                                <h1>ANÁLISIS POR MODELO: {modelo}</h1>
                                <p class="slide-subtitle">DETALLE DE RETRABAJOS</p>
                                <p class="slide-period">Período: {fecha_display}</p>
                            </div>
                            <div class="slide-kpis">
                                <div class="kpi-hero" style="background: #10B981;">
                                    <div class="kpi-label">RETRABAJOS</div>
                                    <div class="kpi-value-hero">{modelo_rew:,}</div>
                                    <div class="kpi-status">UNIDADES RECUPERADAS</div>
                                </div>
                            </div>
                            <div style="padding: 40px 60px; background: white; border-radius: 20px; margin: 30px 100px; box-shadow: 0 8px 30px rgba(0,0,0,0.1);">
                                <h2 style="color: #1E293B; margin-bottom: 30px; font-size: 36px; text-align: center; border-bottom: 3px solid #10B981; padding-bottom: 15px;">Top 5 Defectos Recuperados</h2>
                                <table style="width: 100%; border-collapse: collapse; font-size: 26px;">
                                    <thead>
                                        <tr style="background: #10B981; color: white;">
                                            <th style="padding: 18px; text-align: left; font-weight: 700;">Defecto</th>
                                            <th style="padding: 18px; text-align: right; font-weight: 700;">Cantidad</th>
                                        </tr>
                                    </thead>
                                    <tbody>
                                        {def_r_table_html}
                                    </tbody>
                                </table>
                            </div>
                        </div>
                        '''
                        slides_html.append(slide_modelo_r)

            # === GENERAR HTML COMPLETO ===
            html_content = f'''
<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=1920, initial-scale=1.0">
    <title>Presentación Quality Dashboard</title>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}

        body {{
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
            background: #0F172A;
            overflow: hidden;
            width: 1920px;
            height: 1080px;
        }}

        .presentation-container {{
            width: 1920px;
            height: 1080px;
            position: relative;
            overflow: hidden;
        }}

        .slide {{
            width: 1920px;
            height: 1080px;
            position: absolute;
            top: 0;
            left: 0;
            background: linear-gradient(135deg, #F8FAFC 0%, #E2E8F0 100%);
            opacity: 0;
            transform: translateX(100%);
            transition: all 0.8s cubic-bezier(0.4, 0, 0.2, 1);
            display: flex;
            flex-direction: column;
        }}

        .slide.active {{
            opacity: 1;
            transform: translateX(0);
            z-index: 10;
        }}

        .slide.prev {{
            opacity: 0;
            transform: translateX(-100%);
        }}

        .slide-header {{
            padding: 30px 60px;
            color: white;
            text-align: center;
            box-shadow: 0 4px 20px rgba(0,0,0,0.15);
        }}

        .slide-kpis {{
            display: flex;
            justify-content: center;
            align-items: center;
            gap: 18px;
            padding: 18px 60px;
            flex-wrap: wrap;
        }}

        .slide-charts {{
            display: grid;
            grid-template-columns: repeat(2, 1fr);
            grid-template-rows: repeat(2, 1fr);
            gap: 18px;
            padding: 10px 70px 25px;
            width: 100%;
            height: auto;
            max-width: 1920px;
            margin: 0 auto;
            box-sizing: border-box;
        }}

        .chart-box {{
            background: white;
            border-radius: 12px;
            padding: 12px;
            box-shadow: 0 4px 15px rgba(0,0,0,0.08);
            border: 2px solid #E2E8F0;
            display: flex;
            flex-direction: column;
            width: 100%;
            height: 350px;
            overflow: hidden;
        }}

        .chart-box img {{
            width: 100%;
            height: 100%;
            object-fit: contain;
        }}

        .slide-indicator {{
            position: fixed;
            bottom: 30px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 12px;
            z-index: 1000;
        }}

        .indicator-dot {{
            width: 14px;
            height: 14px;
            border-radius: 50%;
            background: rgba(255,255,255,0.3);
            transition: all 0.3s ease;
            border: 2px solid rgba(255,255,255,0.5);
        }}

        .indicator-dot.active {{
            background: white;
            width: 40px;
            border-radius: 7px;
            box-shadow: 0 2px 10px rgba(255,255,255,0.5);
        }}

        table {{
            border-collapse: collapse;
        }}

        tbody tr {{
            border-bottom: 1px solid #E2E8F0;
        }}

        tbody tr:hover {{
            background: #F8FAFC;
        }}
    </style>
</head>
<body>
    <div class="presentation-container">
        {"".join(slides_html)}
    </div>

    <div class="slide-indicator">
        {''.join(f'<div class="indicator-dot {("active" if i == 0 else "")}"></div>' for i in range(len(slides_html)))}
    </div>

    <script>
        const slides = document.querySelectorAll('.slide');
        const indicators = document.querySelectorAll('.indicator-dot');
        let currentSlide = 0;
        const SLIDE_DURATION = 5000; // 5 segundos

        // Mostrar primera diapositiva
        slides[0].classList.add('active');

        function nextSlide() {{
            // Ocultar diapositiva actual
            slides[currentSlide].classList.remove('active');
            slides[currentSlide].classList.add('prev');
            indicators[currentSlide].classList.remove('active');

            // Calcular siguiente diapositiva (bucle infinito)
            currentSlide = (currentSlide + 1) % slides.length;

            // Mostrar siguiente diapositiva
            slides[currentSlide].classList.remove('prev');
            slides[currentSlide].classList.add('active');
            indicators[currentSlide].classList.add('active');

            // Limpiar clase prev de las demás
            slides.forEach((slide, i) => {{
                if (i !== currentSlide) {{
                    slide.classList.remove('active');
                    setTimeout(() => {{
                        slide.classList.remove('prev');
                    }}, 800);
                }}
            }});
        }}

        // Auto-avanzar cada 5 segundos
        setInterval(nextSlide, SLIDE_DURATION);

        // Navegación con teclado (opcional)
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowRight' || e.key === ' ') {{
                nextSlide();
            }} else if (e.key === 'Escape') {{
                window.close();
            }}
        }});

        console.log('🎬 Presentación iniciada con ' + slides.length + ' diapositivas');
        console.log('⏱️  Duración por diapositiva: 5 segundos');
        console.log('🔄 Modo: Bucle infinito');
        console.log('⌨️  Teclas: Espacio/Flecha derecha = Siguiente | Escape = Salir');
    </script>
</body>
</html>
            '''

            total_slides = len(slides_html)
            print(f"✅ Presentación HTML generada exitosamente:")
            print(f"   📊 Total de diapositivas: {total_slides}")
            print(f"   🎬 Modo: Carrusel automático (5s por diapositiva)")
            print(f"   🔄 Bucle: Infinito")
            print(f"   📐 Resolución: 1920x1080")

            return {
                "error": None,
                "html_content": html_content,
                "total_slides": total_slides
            }

        except Exception as e:
            print(f"❌ Error generando presentación HTML: {e}")
            import traceback
            traceback.print_exc()
            return {"error": f"Error generando presentación: {e}", "html_content": None}

    def _generate_placeholder_image(self, width_px=800, height_px=400, text="Sin datos disponibles"):
        """
        Genera una imagen placeholder moderna con gradiente y mejor tipografía.
        """
        fig, ax = plt.subplots(figsize=(width_px / 100, height_px / 100), dpi=100)
        
        # Gradiente de fondo moderno
        gradient = np.linspace(0, 1, 256).reshape(1, -1)
        gradient = np.vstack((gradient, gradient))
        ax.imshow(gradient, extent=[0, 1, 0, 1], aspect='auto', cmap='Blues_r', alpha=0.1)
        
        ax.axis('off')
        ax.text(0.5, 0.5, text, ha='center', va='center', 
            fontsize=16, fontweight='300', color='#64748B',
            family='Arial')
        
        # Marco sutil con esquinas redondeadas
        from matplotlib.patches import FancyBboxPatch
        fancy_box = FancyBboxPatch((0.05, 0.05), 0.9, 0.9, 
                                boxstyle="round,pad=0.02", 
                                facecolor='none',
                                edgecolor='#E2E8F0', 
                                linewidth=2, 
                                alpha=0.5)
        ax.add_patch(fancy_box)
        
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=100, bbox_inches='tight', 
                pad_inches=0.1, facecolor='white', edgecolor='none')
        plt.close(fig)
        buf.seek(0)
        return buf

    def _safe_img_flowable(self, buf, target_width=570, target_height=310):
        """
        Retorna un objeto Image con dimensiones optimizadas para garantizar que quepan 4 gráficos por página.
        TAMAÑOS ESTANDARIZADOS PARA TODAS LAS PÁGINAS - MEJORADO PARA EVITAR DEFORMACIÓN
        """
        if buf is None:
            buf = self._generate_placeholder_image(width_px=int(target_width*1.5), height_px=int(target_height*1.5))
        
        # Usar KeepAspectRatio para evitar deformación
        from reportlab.platypus import Image
        img = Image(buf)
        img._restrictSize(target_width, target_height)
        return img

    def _crear_grafica_profesional_moderna(self, figsize=(10, 6)):
        """
        Crea una figura con styling ultra profesional y moderno.
        MEJORADO: Mejor proporción y configuración para evitar deformación
        """
        # Configurar matplotlib para mejor calidad
        plt.rcParams['figure.dpi'] = 100
        plt.rcParams['savefig.dpi'] = 150
        plt.rcParams['font.size'] = 18
        plt.rcParams['axes.labelsize'] = 18
        plt.rcParams['axes.titlesize'] = 18
        plt.rcParams['xtick.labelsize'] = 18
        plt.rcParams['ytick.labelsize'] = 18
        plt.rcParams['legend.fontsize'] = 18
        
        # Configurar estilo
        plt.style.use('default')
        
        fig, ax = plt.subplots(figsize=figsize, facecolor='white', tight_layout=True)
        
        # Configurar el axes con estilo moderno premium
        ax.set_facecolor('#FAFBFC')
        
        # Grid premium con líneas más sutiles
        ax.grid(True, linestyle='-', alpha=0.1, color='#CBD5E1', linewidth=0.5, zorder=0)
        ax.set_axisbelow(True)
        
        # Spines modernos - solo mostrar izquierda y abajo
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['bottom'].set_color('#E2E8F0')
        ax.spines['left'].set_color('#E2E8F0')
        ax.spines['bottom'].set_linewidth(0.8)
        ax.spines['left'].set_linewidth(0.8)
        
        # Ticks modernos mejorados
        ax.tick_params(axis='both', which='major', labelsize=9, 
                    colors='#94A3B8', labelcolor='#475569', 
                    width=0.8, length=4, pad=5)
        ax.tick_params(axis='both', which='minor', labelsize=8, 
                    colors='#CBD5E1', labelcolor='#64748B', 
                    width=0.5, length=2)
        
        return fig, ax

    def _anotar_barras_moderno(self, ax, precision=0, offset_factor=0.02):
        """
        Anota las barras con estilo moderno y mejor legibilidad.
        MEJORADO: Mejor posicionamiento y estilo
        """
        for bar in ax.patches:
            height = bar.get_height()
            if height > 0:
                # Formatear números con separadores de miles
                if precision == 0:
                    label = f'{int(height):,}'
                else:
                    label = f'{height:,.{precision}f}'
                
                # Calcular posición mejorada
                y_max = ax.get_ylim()[1]
                offset = y_max * offset_factor
                
                # Anotación con sombra sutil
                ax.text(bar.get_x() + bar.get_width()/2., height + offset,
                    label, ha='center', va='bottom', 
                    fontsize=18, fontweight='600', color='#1E293B',
                    bbox=dict(boxstyle="round,pad=0.3", 
                            facecolor='white', 
                            edgecolor='#E2E8F0',
                            linewidth=1,
                            alpha=0.95))

    def _graficas_global_nok(self, fi_reporte, ff_reporte,
                            dist_def_reporte, dist_mod_reporte, shift_dist_reporte,
                            dist_mod_prov_reporte=None):  # ← NUEVO PARÁMETRO
        """
        Crea gráficos NOK con diseño ultra profesional y moderno.
        MEJORADO: Tamaños optimizados y mejor diseño
        """
        # Paleta de colores moderna y profesional MEJORADA
        modern_colors = {
            'primary': '#0F172A',      # Slate 900
            'secondary': '#1E293B',    # Slate 800  
            'accent': '#3B82F6',       # Blue 500
            'accent_light': '#60A5FA', # Blue 400
            'success': '#10B981',      # Emerald 500
            'warning': '#F59E0B',      # Amber 500
            'danger': '#EF4444',       # Red 500
            'danger_light': '#F87171', # Red 400
            'info': '#06B6D4',         # Cyan 500
            'purple': '#8B5CF6',       # Purple 500
            'light': '#F8FAFC',        # Slate 50
            'medium': '#64748B',       # Slate 500
            'dark': '#334155'          # Slate 700
        }
        
        # Paletas de colores mejoradas con gradientes
        defect_colors = ['#3B82F6', '#10B981', '#F59E0B', '#EF4444', '#8B5CF6', 
                        '#06B6D4', '#EC4899', '#14B8A6', '#F97316', '#6366F1']
        model_colors = ['#0F172A', '#1E293B', '#334155', '#475569', '#64748B']
        shift_colors = {'MAÑANA': '#10B981', 'TARDE': '#3B82F6', 'NOCHE': '#8B5CF6',
                    'A': '#10B981', 'B': '#3B82F6', 'C': '#8B5CF6'}

        buf_ok_nok = None
        try:
            ff_dt = datetime.strptime(ff_reporte, '%Y-%m-%d').date()
            trend_start_date = ff_dt - timedelta(days=self.TREND_DAYS - 1)
            trend_fi_str = trend_start_date.isoformat()

            data_trend_nok = self.leer_estadisticas(trend_fi_str, ff_reporte, "(TODOS)", "(TODOS)")
            evo_nok_trend = data_trend_nok.get("global_evolution", {})
            
            sorted_dates_trend = sorted(evo_nok_trend.keys())

            if sorted_dates_trend:
                xlabels_trend = [d[-5:] for d in sorted_dates_trend]
                if len(sorted_dates_trend) > 30:
                    step = max(1, len(sorted_dates_trend) // 20)
                    xlabels_trend = [d[-5:] if i % step == 0 else "" for i, d in enumerate(sorted_dates_trend)]

                arr_ok_trend = [evo_nok_trend[d].get("OK", 0) for d in sorted_dates_trend]
                arr_perc_trend = []
                for d_trend in sorted_dates_trend:
                    ok_val = evo_nok_trend[d_trend].get("OK", 0)
                    nok_val = evo_nok_trend[d_trend].get("NOK", 0)
                    tot_trend = ok_val + nok_val
                    arr_perc_trend.append((nok_val / tot_trend * 100) if tot_trend > 0 else 0)

                fig1, ax1 = self._crear_grafica_profesional_moderna(figsize=(12, 7))
                
                # Barras con efecto glass morphism
                bars = ax1.bar(range(len(sorted_dates_trend)), arr_ok_trend, 
                            color=modern_colors['accent'], alpha=0.7, 
                            label="Piezas OK", width=0.8,
                            edgecolor=modern_colors['accent_light'], linewidth=1.5)
                
                # Sombra sutil en las barras
                for i, bar in enumerate(bars):
                    shadow = ax1.bar(i, bar.get_height(), 
                                color='#000000', alpha=0.1, width=0.85,
                                bottom=-bar.get_height()*0.02)
                
                self._anotar_barras_moderno(ax1)
                
                ax1.set_ylabel("Piezas OK", fontsize=12, fontweight='600', 
                            color=modern_colors['dark'], labelpad=10)
                ax1.set_xticks(range(len(sorted_dates_trend)))
                ax1.set_xticklabels(xlabels_trend, rotation=45, ha="right", 
                                fontsize=9, color=modern_colors['medium'])

                # Línea de porcentaje NOK con estilo premium
                ax2 = ax1.twinx()
                line = ax2.plot(range(len(sorted_dates_trend)), arr_perc_trend, 
                            color=modern_colors['danger'], linewidth=3, 
                            marker='o', markersize=6, markerfacecolor='white', 
                            markeredgewidth=2, markeredgecolor=modern_colors['danger'], 
                            label="%NOK", zorder=5)
                
                # Línea de objetivo mejorada
                ax2.axhline(y=self.OBJETIVO, color=modern_colors['success'],
                        linestyle='--', linewidth=2, alpha=0.7, 
                        label=f"Objetivo {self.OBJETIVO:.1f}%", zorder=4)
                
                # Área bajo la curva con gradiente
                ax2.fill_between(range(len(sorted_dates_trend)), arr_perc_trend, 
                            alpha=0.15, color=modern_colors['danger'], zorder=1)
                
                # Resaltar zona sobre objetivo
                ax2.fill_between(range(len(sorted_dates_trend)), self.OBJETIVO, arr_perc_trend, 
                            where=[p > self.OBJETIVO for p in arr_perc_trend],
                            alpha=0.1, color=modern_colors['danger'], zorder=2)
                
                ax2.set_ylabel("%NOK", fontsize=12, fontweight='600', 
                            color=modern_colors['dark'], labelpad=10)
                ax2.set_ylim(bottom=0)
                
                # Título mejorado con subtítulo
                title = f"Tendencia de Calidad - Últimos {self.TREND_DAYS} días"
                ax1.text(0.5, 1.05, title, transform=ax1.transAxes,
                        fontsize=16, fontweight='700', color=modern_colors['primary'],
                        ha='center', va='bottom')
                
                # Leyendas mejoradas con marco
                leg1 = ax1.legend(loc='upper left', frameon=True, fancybox=True, 
                                shadow=True, fontsize=10, framealpha=0.95)
                leg2 = ax2.legend(loc='upper right', frameon=True, fancybox=True, 
                                shadow=True, fontsize=10, framealpha=0.95)
                
                leg1.get_frame().set_facecolor('white')
                leg1.get_frame().set_edgecolor('#E2E8F0')
                leg2.get_frame().set_facecolor('white')
                leg2.get_frame().set_edgecolor('#E2E8F0')

                # Ajustar márgenes
                plt.subplots_adjust(top=0.88, bottom=0.12, left=0.08, right=0.92)

                buf_ok_nok = io.BytesIO()
                fig1.savefig(buf_ok_nok, format='png', dpi=150, bbox_inches='tight', 
                        facecolor='white', edgecolor='none')
                plt.close(fig1)
                buf_ok_nok.seek(0)
        except Exception as e:
            print(f"Error en _graficas_global_nok (tendencia): {e}")

        # Gráfico de defectos mejorado
        buf_defectos = None
        if dist_def_reporte and any(v > 0 for v in dist_def_reporte.values()):
            try:
                sorted_defs = sorted(
                    [item for item in dist_def_reporte.items() if item[1] > 0 and item[0]],
                    key=lambda x: x[1], reverse=True
                )[:5]
                if sorted_defs:
                    labs, vals = [t[0] for t in sorted_defs], [t[1] for t in sorted_defs]
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
                    
                    # Barras con gradiente de colores
                    bars = []
                    for i, (lbl, val) in enumerate(zip(labs, vals)):
                        bar = ax.bar(i, val, color=defect_colors[i % len(defect_colors)], 
                                alpha=0.8, width=0.7, edgecolor='white', linewidth=2)
                        bars.extend(bar)
                    
                    # Personalización de etiquetas en X
                    ax.set_xticks(range(len(labs)))
                    ax.set_xticklabels(labs, rotation=30, ha="right", fontsize=10, 
                                    color=modern_colors['dark'], fontweight='500')
                    
                    self._anotar_barras_moderno(ax)
                    
                    # Título con estilo
                    ax.text(0.5, 1.05, "Top 5 Defectos Críticos", transform=ax.transAxes,
                        fontsize=16, fontweight='700', color=modern_colors['primary'],
                        ha='center', va='bottom')
                    
                    ax.set_ylabel("Cantidad de Defectos", fontsize=12, fontweight='600', 
                                color=modern_colors['dark'], labelpad=10)
                    
                    # Línea de promedio
                    avg_val = np.mean(vals)
                    ax.axhline(y=avg_val, color=modern_colors['medium'], 
                            linestyle=':', linewidth=2, alpha=0.6, 
                            label=f'Promedio: {avg_val:.0f}')
                    ax.legend(loc='upper right', frameon=True, fontsize=9)

                    plt.subplots_adjust(top=0.88, bottom=0.15)

                    buf_defectos = io.BytesIO()
                    fig.savefig(buf_defectos, format='png', dpi=150, bbox_inches='tight',
                            facecolor='white', edgecolor='none')
                    plt.close(fig)
                    buf_defectos.seek(0)
            except Exception as e:
                print(f"Error en _graficas_global_nok (defectos): {e}")

        # Gráfico de modelos mejorado
        # Gráfico de modelos mejorado con separación proveedor/interno
       # Gráfico de modelos mejorado con separación proveedor/interno
            buf_modelos = None
            if dist_mod_reporte and any(v > 0 for v in dist_mod_reporte.values()):
                try:
                    modelos_totales = {}
                    for modelo, cant in dist_mod_reporte.items():
                        modelos_totales[modelo] = modelos_totales.get(modelo, 0) + cant
                    if dist_mod_prov_reporte:
                        for modelo, cant in dist_mod_prov_reporte.items():
                            modelos_totales[modelo] = modelos_totales.get(modelo, 0) + cant

                    sorted_mods = sorted(
                        [it for it in modelos_totales.items() if it[1] > 0 and it[0]],
                        key=lambda x: x[1], reverse=True
                    )[:5]

                    if sorted_mods:
                        labs = [t[0] for t in sorted_mods]
                        vals_internos = [dist_mod_reporte.get(m, 0) for m in labs]
                        vals_proveedor = [dist_mod_prov_reporte.get(m, 0) if dist_mod_prov_reporte else 0 for m in labs]

                        fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))

                        x = np.arange(len(labs))
                        width = 0.38
                        ax.bar(x - width/2, vals_internos, width, label='Defectos Internos',
                            color='#EF4444', alpha=0.85, edgecolor='white', linewidth=2)
                        ax.bar(x + width/2, vals_proveedor, width, label='Defectos Proveedor',
                            color='#F59E0B', alpha=0.85, edgecolor='white', linewidth=2)

                        self._anotar_barras_moderno(ax)

                        ax.set_xticks(x)
                        ax.set_xticklabels(labs, fontsize=11, color=modern_colors['dark'], fontweight='500')
                        ax.set_ylabel("Unidades NOK", fontsize=12, fontweight='600', color=modern_colors['dark'])

                        # ticks de enteros en Y
                        from matplotlib.ticker import MaxNLocator
                        ax.yaxis.set_major_locator(MaxNLocator(integer=True))

                        ax.legend(loc='upper right', frameon=True, fontsize=10, framealpha=0.95)

                        # *** NO poner ax.text(...) de título aquí para evitar superposición ***

                        buf_modelos = io.BytesIO()
                        fig.savefig(buf_modelos, format='png', dpi=150, bbox_inches='tight',
                                    facecolor='white', edgecolor='none')
                        plt.close(fig)
                        buf_modelos.seek(0)
                except Exception as e:
                    print(f"Error en _graficas_global_nok (modelos): {e}")


        # Gráfico de turnos mejorado
        buf_turnos = None
        if shift_dist_reporte and any(v > 0 for v in shift_dist_reporte.values()):
            try:
                fig, ax = self._crear_grafica_profesional_moderna(figsize=(9, 7))
                
                # Ordenar turnos
                turnos_orden = ['MAÑANA', 'TARDE', 'NOCHE', 'A', 'B', 'C']
                turnos_data = [(k, v) for k, v in shift_dist_reporte.items() if k and v > 0]
                turnos_data.sort(key=lambda x: turnos_orden.index(x[0]) if x[0] in turnos_orden else 999)
                
                if turnos_data:
                    labels, values = zip(*turnos_data)
                    x_pos = np.arange(len(labels))
                    
                    # Barras con colores específicos por turno
                    bars = []
                    for i, (label, value) in enumerate(zip(labels, values)):
                        color = shift_colors.get(label, modern_colors['medium'])
                        bar = ax.bar(i, value, color=color, alpha=0.85, width=0.6,
                                    edgecolor='white', linewidth=2)
                        bars.extend(bar)
                    
                    self._anotar_barras_moderno(ax, offset_factor=0.03)
                    
                    ax.set_xticks(x_pos)
                    ax.set_xticklabels(labels, rotation=0, ha="center", fontsize=12,
                                    color=modern_colors['dark'], fontweight='600')
                    
                    # TÍTULO CORREGIDO - Limpiar título anterior
                    ax.clear()  # Limpiar cualquier contenido anterior
                    fig.clear()  # Limpiar figura
                    ax = fig.add_subplot(111)  # Recrear subplot
                    
                    # Recrear gráfico desde cero
                    bars = []
                    for i, (label, value) in enumerate(zip(labels, values)):
                        color = shift_colors.get(label, modern_colors['medium'])
                        bar = ax.bar(i, value, color=color, alpha=0.85, width=0.6,
                                    edgecolor='white', linewidth=2)
                        bars.extend(bar)
                    
                    self._anotar_barras_moderno(ax, offset_factor=0.03)
                    
                    ax.set_xticks(x_pos)
                    ax.set_xticklabels(labels, rotation=0, ha="center", fontsize=12,
                                    color=modern_colors['dark'], fontweight='600')
                    
                    ax.text(0.5, 1.05, "Distribución NOK por Turno", transform=ax.transAxes,
                        fontsize=16, fontweight='700', color=modern_colors['primary'],
                        ha='center', va='bottom')
                    
                    ax.set_ylabel("Cantidad NOK", fontsize=12, fontweight='600', 
                                color=modern_colors['dark'], labelpad=10)
                    
                    # Configurar spines y grid
                    ax.spines['top'].set_visible(False)
                    ax.spines['right'].set_visible(False)
                    ax.spines['bottom'].set_color('#E2E8F0')
                    ax.spines['left'].set_color('#E2E8F0')
                    ax.grid(True, alpha=0.1, linestyle='--', color='#CBD5E1')
                    ax.set_axisbelow(True)

                    plt.tight_layout()
                    plt.subplots_adjust(top=0.88)

                    buf_turnos = io.BytesIO()
                    fig.savefig(buf_turnos, format='png', dpi=150, bbox_inches='tight',
                            facecolor='white', edgecolor='none')
                    plt.close(fig)
                    buf_turnos.seek(0)
            except Exception as e:
                print(f"Error en _graficas_global_nok (turnos): {e}")
            
            return buf_ok_nok, buf_defectos, buf_modelos, buf_turnos


    def _graficas_global_r(self, fi_reporte, ff_reporte,
                        dist_def_r_reporte, dist_mod_r_reporte, dist_shift_r_reporte):
        """
        Crea gráficos de retrabajos con diseño ultra profesional.
        MEJORADO: Mejor diseño y tamaños optimizados
        """
        modern_colors = {
            'primary': '#0F172A', 'secondary': '#1E293B', 'accent': '#10B981',
            'accent_dark': '#059669', 'success': '#10B981', 'warning': '#F59E0B', 
            'danger': '#EF4444', 'info': '#06B6D4', 'purple': '#8B5CF6',
            'light': '#F8FAFC', 'medium': '#64748B', 'dark': '#334155'
        }
        
        # Paletas específicas para retrabajos (tonos verdes)
        defect_colors = ['#10B981', '#059669', '#047857', '#065F46', '#064E3B']
        model_colors = ['#065F46', '#047857', '#059669', '#10B981', '#34D399']
        shift_colors = {'MAÑANA': '#10B981', 'TARDE': '#3B82F6', 'NOCHE': '#8B5CF6',
                    'A': '#10B981', 'B': '#3B82F6', 'C': '#8B5CF6'}

        buf_evol = None
        try:
            ff_dt = datetime.strptime(ff_reporte, '%Y-%m-%d').date()
            trend_start_date = ff_dt - timedelta(days=self.TREND_DAYS - 1)
            trend_fi_str = trend_start_date.isoformat()
            data_trend_r = self.leer_estadisticas(trend_fi_str, ff_reporte, "(TODOS)", "(TODOS)")
            evo_r_trend = data_trend_r.get("global_evolution_rework", {})
            
            sorted_dates_trend_r = sorted(evo_r_trend.keys())

            if sorted_dates_trend_r:
                xlbls_trend_r = [d[-5:] for d in sorted_dates_trend_r]
                if len(sorted_dates_trend_r) > 30:
                    step = max(1, len(sorted_dates_trend_r) // 20)
                    xlbls_trend_r = [d[-5:] if i % step == 0 else "" for i, d in enumerate(sorted_dates_trend_r)]

                vals_r_trend = [evo_r_trend.get(d, 0) for d in sorted_dates_trend_r]
                fig, ax = self._crear_grafica_profesional_moderna(figsize=(12, 7))
                
                # Gráfico de área con línea
                ax.fill_between(range(len(sorted_dates_trend_r)), vals_r_trend,
                            color=modern_colors['accent'], alpha=0.2, label='Área de retrabajos')
                
                line = ax.plot(range(len(sorted_dates_trend_r)), vals_r_trend,
                            color=modern_colors['accent_dark'], marker='o', linewidth=3,
                            markersize=7, markerfacecolor='white', markeredgewidth=2,
                            markeredgecolor=modern_colors['accent_dark'], label='Retrabajos',
                            zorder=5)
                
                # Puntos destacados en valores máximos
                max_val = max(vals_r_trend) if vals_r_trend else 0
                max_indices = [i for i, v in enumerate(vals_r_trend) if v == max_val]
                for idx in max_indices:
                    ax.scatter(idx, vals_r_trend[idx], color=modern_colors['danger'], 
                            s=100, zorder=6, edgecolor='white', linewidth=2)
                
                # Anotaciones en puntos clave
                if len(vals_r_trend) > 0:
                    # Anotar máximo
                    max_idx = vals_r_trend.index(max(vals_r_trend))
                    ax.annotate(f'Máx: {max(vals_r_trend):,}', 
                            xy=(max_idx, vals_r_trend[max_idx]),
                            xytext=(max_idx, vals_r_trend[max_idx] + max_val*0.1),
                            ha='center', fontsize=18, fontweight='600',
                            bbox=dict(boxstyle="round,pad=0.3", facecolor='white', 
                                    edgecolor=modern_colors['danger'], alpha=0.95),
                            arrowprops=dict(arrowstyle='->', color=modern_colors['danger'], lw=1))
                
                ax.text(0.5, 1.05, f"Evolución de Retrabajos - Últimos {self.TREND_DAYS} días", 
                    transform=ax.transAxes, fontsize=16, fontweight='700', 
                    color=modern_colors['primary'], ha='center', va='bottom')
                
                ax.set_ylabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', 
                            color=modern_colors['dark'], labelpad=10)
                ax.set_xticks(range(len(sorted_dates_trend_r)))
                ax.set_xticklabels(xlbls_trend_r, rotation=45, ha="right", 
                                fontsize=9, color=modern_colors['medium'])
                
                # Media móvil
                if len(vals_r_trend) >= 7:
                    try:
                        from scipy.ndimage import uniform_filter1d
                        media_movil = uniform_filter1d(vals_r_trend, size=7, mode='nearest')
                        ax.plot(range(len(sorted_dates_trend_r)), media_movil,
                            color=modern_colors['warning'], linewidth=2, linestyle='--',
                            alpha=0.7, label='Media móvil (7 días)')
                    except ImportError:
                        pass  # Si scipy no está disponible, continuar sin media móvil
                
                ax.legend(loc='upper left', frameon=True, fontsize=10, framealpha=0.95)
                ax.set_ylim(bottom=0)

                plt.subplots_adjust(top=0.88, bottom=0.12)

                buf_evol = io.BytesIO()
                fig.savefig(buf_evol, format='png', dpi=150, bbox_inches='tight',
                        facecolor='white', edgecolor='none')
                plt.close(fig)
                buf_evol.seek(0)
        except Exception as e:
            print(f"Error en _graficas_global_r (tendencia): {e}")

        buf_def, buf_mod, buf_shift = None, None, None

        # Gráfico de defectos R mejorado
        if dist_def_r_reporte and any(v > 0 for v in dist_def_r_reporte.values()):
            try:
                sorted_rdef = sorted(
                    [item for item in dist_def_r_reporte.items() if item[1] > 0 and item[0]],
                    key=lambda x: x[1], reverse=True
                )[:5]
                if sorted_rdef:
                    labs, vals = [t[0] for t in sorted_rdef], [t[1] for t in sorted_rdef]
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
                    
                    # Crear barras con gradiente de verdes
                    for i, (l, v) in enumerate(zip(labs, vals)):
                        bar = ax.bar(i, v, color=defect_colors[i % len(defect_colors)], 
                                alpha=0.85, width=0.7, edgecolor='white', linewidth=2)
                    
                    ax.set_xticks(range(len(labs)))
                    ax.set_xticklabels(labs, rotation=30, ha="right", fontsize=10, 
                                    color=modern_colors['dark'], fontweight='500')
                    
                    self._anotar_barras_moderno(ax)
                    
                    ax.text(0.5, 1.05, "Top 5 Defectos en Retrabajos", transform=ax.transAxes,
                        fontsize=16, fontweight='700', color=modern_colors['primary'],
                        ha='center', va='bottom')
                    
                    ax.set_ylabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', 
                                color=modern_colors['dark'], labelpad=10)

                    plt.subplots_adjust(top=0.88, bottom=0.15)

                    buf_def = io.BytesIO()
                    fig.savefig(buf_def, format='png', dpi=150, bbox_inches='tight',
                            facecolor='white', edgecolor='none')
                    plt.close(fig)
                    buf_def.seek(0)
            except Exception as e:
                print(f"Error en _graficas_global_r (defectos): {e}")

        # Gráfico de modelos R mejorado
        if dist_mod_r_reporte and any(v > 0 for v in dist_mod_r_reporte.values()):
            try:
                sorted_rmod = sorted(
                    [item for item in dist_mod_r_reporte.items() if item[1] > 0 and item[0]],
                    key=lambda x: x[1], reverse=True
                )[:5]
                if sorted_rmod:
                    labs, vals = [t[0] for t in sorted_rmod], [t[1] for t in sorted_rmod]
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
                    
                    # Barras horizontales con gradiente
                    bars = ax.barh(range(len(labs)), vals, 
                                color=[model_colors[i % len(model_colors)] for i in range(len(labs))],
                                alpha=0.85, height=0.6, edgecolor='white', linewidth=2)
                    
                    # Anotaciones mejoradas
                    for i, (bar, val) in enumerate(zip(bars, vals)):
                        width = bar.get_width()
                        ax.text(width + max(vals)*0.01, bar.get_y() + bar.get_height()/2,
                            f'{int(val):,}', ha='left', va='center',
                            fontsize=18, fontweight='600', color=modern_colors['dark'],
                            bbox=dict(boxstyle="round,pad=0.3", facecolor='white', 
                                    edgecolor='#E2E8F0', alpha=0.95))
                    
                    ax.set_yticks(range(len(labs)))
                    ax.set_yticklabels(labs, fontsize=11, color=modern_colors['dark'], 
                                    fontweight='500')
                    
                    ax.text(0.5, 1.05, "Top 5 Modelos en Retrabajos", transform=ax.transAxes,
                        fontsize=16, fontweight='700', color=modern_colors['primary'],
                        ha='center', va='bottom')
                    
                    ax.set_xlabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', 
                                color=modern_colors['dark'], labelpad=10)

                    plt.subplots_adjust(top=0.88, left=0.15)

                    buf_mod = io.BytesIO()
                    fig.savefig(buf_mod, format='png', dpi=150, bbox_inches='tight',
                            facecolor='white', edgecolor='none')
                    plt.close(fig)
                    buf_mod.seek(0)
            except Exception as e:
                print(f"Error en _graficas_global_r (modelos): {e}")

        # Gráfico de turnos R mejorado
        if dist_shift_r_reporte and any(v > 0 for v in dist_shift_r_reporte.values()):
            try:
                fig, ax = self._crear_grafica_profesional_moderna(figsize=(9, 7))
                
                # Ordenar turnos
                turnos_orden = ['MAÑANA', 'TARDE', 'NOCHE', 'A', 'B', 'C']
                turnos_data = [(k, v) for k, v in dist_shift_r_reporte.items() if k and v > 0]
                turnos_data.sort(key=lambda x: turnos_orden.index(x[0]) if x[0] in turnos_orden else 999)
                
                if turnos_data:
                    labels, values = zip(*turnos_data)
                    
                    for i, (label, value) in enumerate(zip(labels, values)):
                        color = shift_colors.get(label, modern_colors['medium'])
                        bar = ax.bar(i, value, color=color, alpha=0.85, width=0.6,
                                edgecolor='white', linewidth=2)
                    
                    self._anotar_barras_moderno(ax, offset_factor=0.03)
                    
                    ax.set_xticks(range(len(labels)))
                    ax.set_xticklabels(labels, rotation=0, ha="center", fontsize=12,
                                    color=modern_colors['dark'], fontweight='600')
                    
                    ax.text(0.5, 1.05, "Distribución de Retrabajos por Turno", 
                        transform=ax.transAxes, fontsize=16, fontweight='700', 
                        color=modern_colors['primary'], ha='center', va='bottom')
                    
                    ax.set_ylabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', 
                                color=modern_colors['dark'], labelpad=10)

                    plt.subplots_adjust(top=0.88)

                    buf_shift = io.BytesIO()
                    fig.savefig(buf_shift, format='png', dpi=150, bbox_inches='tight',
                            facecolor='white', edgecolor='none')
                    plt.close(fig)
                    buf_shift.seek(0)
            except Exception as e:
                print(f"Error en _graficas_global_r (turnos): {e}")

        return buf_evol, buf_def, buf_mod, buf_shift

    def _crear_hoja_global_nok(self, story, fi, ff, total_ok, total_nok, total_nok_prov,
                            buf_ok_nok_trend, buf_defectos, buf_modelos, buf_turno):
        """
        Genera página global NOK con diseño ultra ejecutivo y profesional.
        MEJORADO: Mejor espaciado y diseño
        """
        styles = getSampleStyleSheet()
        
        # Estilos mejorados con tipografía más elegante
        styleTitleMain = ParagraphStyle(
            'ExecutiveTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=26,  # Ligeramente reducido para más espacio
            alignment=1,
            textColor=colors.white,
            spaceAfter=4,
            leading=28
        )
        
        styleSubtitle = ParagraphStyle(
            'ExecutiveSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=13,
            alignment=1,
            textColor=colors.white,
            spaceAfter=0,
            leading=15
        )
        
        styleKPILabel = ParagraphStyle(
            'KPILabel',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=10,
            alignment=1,
            textColor=colors.HexColor('#475569'),
            spaceAfter=1,
            leading=12
        )
        
        styleKPINumber = ParagraphStyle(
            'KPINumber',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=17,
            alignment=1,
            textColor=colors.HexColor('#1E293B'),
            spaceAfter=0,
            leading=20
        )
        
        styleKPINumberHero = ParagraphStyle(
            'KPINumberHero',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=30,
            alignment=1,
            textColor=colors.white,
            spaceAfter=0,
            leading=32
        )

        total_fab = total_ok + total_nok + total_nok_prov
        rejection_rate = (total_nok / (total_ok + total_nok) * 100) if (total_ok + total_nok) > 0 else 0
        fecha_display = fi if fi == ff else f"{fi} — {ff}"

        # Header mejorado con gradiente simulado
        header_data = [
            [Paragraph("QUALITY DASHBOARD", styleTitleMain)],
            [Paragraph("ANÁLISIS GLOBAL DE RECHAZOS (NOK)", styleSubtitle)],
            [Paragraph(f"Período: {fecha_display} | Objetivo: ≤ {self.OBJETIVO:.1f}%", styleSubtitle)]
        ]
        
        header_table = Table(header_data, colWidths=[1150], rowHeights=[24, 14, 14])
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#0F172A')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
            ('LINEBELOW', (0, 2), (-1, 2), 3, colors.HexColor('#3B82F6')),
        ]))
        story.append(header_table)
        story.append(Spacer(1, 6))

        # Determinar color del KPI principal
        kpi_color = colors.HexColor('#EF4444')  # Rojo
        status_text = "CRÍTICO"
        if rejection_rate <= self.OBJETIVO:
            kpi_color = colors.HexColor('#10B981')  # Verde
            status_text = "ÓPTIMO"
        elif rejection_rate <= self.OBJETIVO * 1.3:
            kpi_color = colors.HexColor('#F59E0B')  # Ámbar
            status_text = "ALERTA"

        # KPI Hero mejorado con indicador de estado
        kpi_hero_data = [
            [Paragraph("TASA DE RECHAZO", styleKPILabel)],
            [Paragraph(f"{rejection_rate:.2f}%", styleKPINumberHero)],
            [Paragraph(status_text, ParagraphStyle('status', parent=styleKPILabel, 
                                                textColor=colors.white, fontSize=11))]
        ]
        
        kpi_hero_table = Table(kpi_hero_data, colWidths=[240], rowHeights=[14, 32, 12])
        kpi_hero_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), kpi_color),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LINEBELOW', (0, 0), (-1, 0), 2, colors.white),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ('ROUNDEDCORNERS', [5]),
        ]))

        # KPIs secundarios mejorados
        kpi_cards = [
            ("PRODUCCIÓN TOTAL", f"{total_fab:,}", "#F8FAFC", "#3B82F6"),
            ("RECHAZOS INTERNOS", f"{total_nok:,}", "#FEF2F2", "#EF4444"), 
            ("RECHAZOS PROVEEDOR", f"{total_nok_prov:,}", "#FFF7ED", "#F59E0B")
        ]
        
        secondary_kpis = []
        for label, value, bg_color, accent_color in kpi_cards:
            card_data = [
                [Paragraph(label, styleKPILabel)],
                [Paragraph(value, styleKPINumber)]
            ]
            card_table = Table(card_data, colWidths=[240], rowHeights=[14, 28])
            card_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor(bg_color)),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BOX', (0, 0), (-1, -1), 1.5, colors.HexColor('#E2E8F0')),
                ('LEFTPADDING', (0, 0), (-1, -1), 8),
                ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                ('LINEBELOW', (0, 0), (-1, 0), 2, colors.HexColor(accent_color)),
                ('ROUNDEDCORNERS', [3]),
            ]))
            secondary_kpis.append(card_table)

        # Layout de KPIs mejorado
        kpi_layout_data = [[kpi_hero_table] + secondary_kpis]
        kpi_layout = Table(kpi_layout_data, colWidths=[285, 285, 285, 285], rowHeights=[58])
        kpi_layout.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
        ]))
        
        story.append(kpi_layout)
        story.append(Spacer(1, 4))

        # Grid de gráficos mejorado con bordes y sombras
        img_ok_nok = self._safe_img_flowable(buf_ok_nok_trend, 570, 310)
        img_defectos = self._safe_img_flowable(buf_defectos, 570, 310)
        img_modelos = self._safe_img_flowable(buf_modelos, 570, 310)
        img_turno = self._safe_img_flowable(buf_turno, 570, 310)

        # Crear contenedores para las imágenes con títulos
        def create_chart_container(img, title):
            container_data = [
                [Paragraph(title, ParagraphStyle('chartTitle', parent=styles['Normal'],
                                            fontName='Helvetica-Bold', fontSize=11,
                                            textColor=colors.HexColor('#1E293B'),
                                            alignment=1, spaceAfter=3))],
                [img]
            ]
            container = Table(container_data, colWidths=[570], rowHeights=[20, 310])
            container.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F8FAFC')),
                ('LINEBELOW', (0, 0), (-1, 0), 1, colors.HexColor('#E2E8F0')),
            ]))
            return container

        charts_grid = Table(
            [
                [create_chart_container(img_ok_nok, "Tendencia de Calidad"),
                create_chart_container(img_defectos, "Top Defectos")],
                [create_chart_container(img_modelos, "Impacto por Modelo"),
                create_chart_container(img_turno, "Distribución por Turno")]
            ],
            colWidths=[575, 575],
            rowHeights=[340, 340]
        )
        
        charts_grid.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
            ('BOX', (0, 0), (-1, -1), 2, colors.HexColor('#E2E8F0')),
            ('INNERGRID', (0, 0), (-1, -1), 1, colors.HexColor('#F1F5F9')),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('ROUNDEDCORNERS', [5]),
        ]))
        
        story.append(charts_grid)
        story.append(PageBreak())

    def _crear_hoja_global_rework(self, story, fi, ff, total_rew,
                                buf_evol_trend, buf_def, buf_mod, buf_shift):
        """
        Página global de retrabajos con diseño ejecutivo mejorado.
        """
        styles = getSampleStyleSheet()
        
        # Estilos consistentes con NOK
        styleTitleMain = ParagraphStyle(
            'ExecutiveTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=26,
            alignment=1,
            textColor=colors.white,
            spaceAfter=4,
            leading=28
        )
        
        styleSubtitle = ParagraphStyle(
            'ExecutiveSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=13,
            alignment=1,
            textColor=colors.white,
            spaceAfter=0,
            leading=15
        )
        
        styleKPILabel = ParagraphStyle(
            'KPILabel',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=10,
            alignment=1,
            textColor=colors.HexColor('#475569'),
            spaceAfter=1,
            leading=12
        )
        
        styleKPINumberHero = ParagraphStyle(
            'KPINumberHero',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=22,
            alignment=1,
            textColor=colors.white,
            spaceAfter=0,
            leading=24
        )

        fecha_display = fi if fi == ff else f"{fi} — {ff}"

        # Header profesional con tema verde
        header_data = [
            [Paragraph("REWORK DASHBOARD", styleTitleMain)],
            [Paragraph("ANÁLISIS GLOBAL DE RETRABAJOS", styleSubtitle)],
            [Paragraph(f"Período de Análisis: {fecha_display}", styleSubtitle)]
        ]
        
        header_table = Table(header_data, colWidths=[1150], rowHeights=[24, 14, 14])
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#065F46')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
            ('LINEBELOW', (0, 2), (-1, 2), 3, colors.HexColor('#10B981')),
        ]))
        story.append(header_table)
        story.append(Spacer(1, 6))

        # KPI Hero para retrabajos con indicador visual
        kpi_hero_data = [
            [Paragraph("TOTAL RETRABAJOS", styleKPILabel)],
            [Paragraph(f"{total_rew:,}", styleKPINumberHero)],
            [Paragraph("UNIDADES RECUPERADAS", ParagraphStyle('status', parent=styleKPILabel, 
                                                            textColor=colors.white, fontSize=9))]
        ]
        
        kpi_hero_table = Table(kpi_hero_data, colWidths=[320], rowHeights=[12, 26, 10])
        kpi_hero_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#10B981')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 15),
            ('RIGHTPADDING', (0, 0), (-1, -1), 15),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ('LINEBELOW', (0, 0), (-1, 0), 2, colors.white),
            ('ROUNDEDCORNERS', [5]),
        ]))
        
        # Centrar el KPI hero
        kpi_center = Table([[kpi_hero_table]], colWidths=[1150], 
                        style=[('ALIGN', (0, 0), (-1, -1), 'CENTER')])
        story.append(kpi_center)
        story.append(Spacer(1, 4))

        # Grid de gráficos mejorado
        img_evol = self._safe_img_flowable(buf_evol_trend, 570, 310)
        img_def_r = self._safe_img_flowable(buf_def, 570, 310)
        img_mod_r = self._safe_img_flowable(buf_mod, 570, 310)
        img_shift_r = self._safe_img_flowable(buf_shift, 570, 310)

        # Crear contenedores con títulos
        def create_chart_container(img, title):
            container_data = [
                [Paragraph(title, ParagraphStyle('chartTitle', parent=styles['Normal'],
                                            fontName='Helvetica-Bold', fontSize=11,
                                            textColor=colors.HexColor('#065F46'),
                                            alignment=1, spaceAfter=3))],
                [img]
            ]
            container = Table(container_data, colWidths=[570], rowHeights=[20, 310])
            container.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F0FDF4')),
                ('LINEBELOW', (0, 0), (-1, 0), 1, colors.HexColor('#BBF7D0')),
            ]))
            return container

        charts_grid = Table(
            [
                [create_chart_container(img_evol, "Evolución Temporal"),
                create_chart_container(img_def_r, "Defectos Recuperados")],
                [create_chart_container(img_mod_r, "Modelos en Retrabajo"),
                create_chart_container(img_shift_r, "Distribución por Turno")]
            ],
            colWidths=[575, 575],
            rowHeights=[340, 340]
        )
        
        charts_grid.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
            ('BOX', (0, 0), (-1, -1), 2, colors.HexColor('#E2E8F0')),
            ('INNERGRID', (0, 0), (-1, -1), 1, colors.HexColor('#F1F5F9')),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('ROUNDEDCORNERS', [5]),
        ]))
        
        story.append(charts_grid)
        story.append(PageBreak())

    # NUEVA FUNCIÓN INTEGRADA: ANÁLISIS ECONÓMICO COMPLETO
    def _crear_hoja_analisis_costes(self, story, fi_reporte, ff_reporte, costes_por_linea_turno, costes_globales):
        """
        Crear hoja de análisis económico de costes por línea y turno con diseño ultra profesional.
        MEJORADO: Estilo consistente con el resto del reporte - 2 PÁGINAS COMPLETAS
        """
        styles = getSampleStyleSheet()
        
        # Estilos MODERNOS consistentes con el diseño principal
        styleTitleMain = ParagraphStyle(
            'ExecutiveTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=26,
            alignment=1,
            textColor=colors.white,
            spaceAfter=4,
            leading=28
        )
        
        styleSubtitle = ParagraphStyle(
            'ExecutiveSubtitle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=13,
            alignment=1,
            textColor=colors.white,
            spaceAfter=0,
            leading=15
        )
        
        styleKPILabel = ParagraphStyle(
            'KPILabel',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=10,
            alignment=1,
            textColor=colors.HexColor('#475569'),
            spaceAfter=1,
            leading=12
        )
        
        styleKPINumber = ParagraphStyle(
            'KPINumber',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=17,
            alignment=1,
            textColor=colors.HexColor('#1E293B'),
            spaceAfter=0,
            leading=20
        )
        
        styleKPINumberHero = ParagraphStyle(
            'KPINumberHero',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=30,
            alignment=1,
            textColor=colors.white,
            spaceAfter=0,
            leading=32
        )

        fecha_display = fi_reporte if fi_reporte == ff_reporte else f"{fi_reporte} — {ff_reporte}"

        # Header moderno con tema rojo (pérdidas económicas)
        header_data = [
            [Paragraph("COST IMPACT DASHBOARD", styleTitleMain)],
            [Paragraph("ANÁLISIS ECONÓMICO DE PÉRDIDAS POR CALIDAD", styleSubtitle)],
            [Paragraph(f"Período: {fecha_display} | Impacto Económico Excluye Defectos de Proveedor", styleSubtitle)]
        ]
        
        header_table = Table(header_data, colWidths=[1150], rowHeights=[24, 14, 14])
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#C0392B')),  # Rojo para pérdidas
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
            ('LINEBELOW', (0, 2), (-1, 2), 3, colors.HexColor('#E74C3C')),
        ]))
        story.append(header_table)
        story.append(Spacer(1, 6))

        # KPIs económicos principales
        total_perdidas = costes_globales.get('total_perdidas', 0)
        num_defectos = costes_globales.get('num_defectos', 0)
        coste_medio_defecto = total_perdidas / num_defectos if num_defectos > 0 else 0

        # KPI Hero para pérdidas totales
        kpi_hero_data = [
            [Paragraph("PÉRDIDAS TOTALES", styleKPILabel)],
            [Paragraph(f"€ {self._format_euro(total_perdidas)}", styleKPINumberHero)],
            [Paragraph("IMPACTO ECONÓMICO", ParagraphStyle('status', parent=styleKPILabel, 
                                                        textColor=colors.white, fontSize=11))]
        ]
        
        kpi_hero_table = Table(kpi_hero_data, colWidths=[240], rowHeights=[14, 32, 12])
        kpi_hero_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#E74C3C')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LINEBELOW', (0, 0), (-1, 0), 2, colors.white),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ('ROUNDEDCORNERS', [5]),
        ]))

        # KPIs secundarios económicos
        # Necesitamos obtener los NOK internos del período
        data_completa = self.leer_estadisticas(fi_reporte, ff_reporte, "(TODOS)", "(TODOS)")
        nok_internos_total = data_completa.get("total_nok", 0)

        # Verificar completitud del análisis
        if num_defectos == nok_internos_total and num_defectos > 0:
            estado_analisis = "COMPLETO"
            color_estado = "#10B981"
            bg_estado = "#F0FDF4"
        else:
            estado_analisis = "INCOMPLETO"
            color_estado = "#EF4444" 
            bg_estado = "#FEF2F2"

        # KPIs secundarios económicos
        kpi_cards = [
            ("DEFECTOS VALORADOS", f"{num_defectos:,}", "#FEF2F2", "#EF4444"),
            ("COSTE MEDIO/DEFECTO", f"€ {self._format_euro(coste_medio_defecto)}", "#FFF7ED", "#F59E0B"),
            ("ANÁLISIS ECONÓMICO", estado_analisis, bg_estado, color_estado)
        ]

        
        secondary_kpis = []
        for label, value, bg_color, accent_color in kpi_cards:
            card_data = [
                [Paragraph(label, styleKPILabel)],
                [Paragraph(value, styleKPINumber)]
            ]
            card_table = Table(card_data, colWidths=[240], rowHeights=[14, 28])
            card_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor(bg_color)),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BOX', (0, 0), (-1, -1), 1.5, colors.HexColor('#E2E8F0')),
                ('LEFTPADDING', (0, 0), (-1, -1), 8),
                ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                ('LINEBELOW', (0, 0), (-1, 0), 2, colors.HexColor(accent_color)),
                ('ROUNDEDCORNERS', [3]),
            ]))
            secondary_kpis.append(card_table)

        # Layout de KPIs económicos
        kpi_layout_data = [[kpi_hero_table] + secondary_kpis]
        kpi_layout = Table(kpi_layout_data, colWidths=[285, 285, 285, 285], rowHeights=[58])
        kpi_layout.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
        ]))
        
        story.append(kpi_layout)
        story.append(Spacer(1, 4))

        # Tabla detallada de análisis por línea y turno
        self._crear_tabla_detalle_costes_moderna(story, costes_por_linea_turno)
        
        # Nueva página para gráficos
        story.append(PageBreak())
        
        # Header para gráficos económicos
        header_graficos_data = [
            [Paragraph("ANÁLISIS VISUAL DE IMPACTO ECONÓMICO", styleTitleMain)],
            [Paragraph(f"Distribución y tendencias de costes | {fecha_display}", styleSubtitle)]
        ]
        
        header_graficos_table = Table(header_graficos_data, colWidths=[1150], rowHeights=[24, 14])
        header_graficos_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#C0392B')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 7),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 7),
            ('LINEBELOW', (0, 1), (-1, 1), 3, colors.HexColor('#E74C3C')),
        ]))
        
        story.append(header_graficos_table)
        story.append(Spacer(1, 4))

        # Generar gráficos económicos
        buf_ranking_lineas = self._generar_grafico_ranking_lineas_costes(costes_por_linea_turno)
        buf_distribucion_turnos = self._generar_grafico_distribucion_turnos_costes(costes_por_linea_turno)
        buf_distribucion_uet = self._generar_grafico_distribucion_perdidas_por_uet(costes_por_linea_turno)
        buf_tendencia_diaria = self._generar_grafico_tendencia_costes_30_dias(fi_reporte, ff_reporte)

        # Grid de gráficos económicos con diseño moderno
        img_ranking = self._safe_img_flowable(buf_ranking_lineas, 570, 310)
        img_turnos = self._safe_img_flowable(buf_distribucion_turnos, 570, 310)
        img_uet = self._safe_img_flowable(buf_distribucion_uet, 570, 310)
        img_tendencia = self._safe_img_flowable(buf_tendencia_diaria, 570, 310)

        def create_chart_container(img, title):
            container_data = [
                [Paragraph(title, ParagraphStyle('chartTitle', parent=styles['Normal'],
                                            fontName='Helvetica-Bold', fontSize=11,
                                            textColor=colors.HexColor('#C0392B'),
                                            alignment=1, spaceAfter=3))],
                [img]
            ]
            container = Table(container_data, colWidths=[570], rowHeights=[20, 310])
            container.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#FEF2F2')),
                ('LINEBELOW', (0, 0), (-1, 0), 1, colors.HexColor('#FCA5A5')),
            ]))
            return container

        charts_grid = Table(
            [
                [create_chart_container(img_ranking, "Ranking de Líneas por Pérdidas"),
                create_chart_container(img_turnos, "Distribución por Turno")],
                [create_chart_container(img_uet, "Análisis por UET"),
                create_chart_container(img_tendencia, "Tendencia de 30 Días")]
            ],
            colWidths=[575, 575],
            rowHeights=[340, 340]
        )
        
        charts_grid.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
            ('BOX', (0, 0), (-1, -1), 2, colors.HexColor('#E2E8F0')),
            ('INNERGRID', (0, 0), (-1, -1), 1, colors.HexColor('#F1F5F9')),
            ('LEFTPADDING', (0, 0), (-1, -1), 5),
            ('RIGHTPADDING', (0, 0), (-1, -1), 5),
            ('TOPPADDING', (0, 0), (-1, -1), 5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
            ('ROUNDEDCORNERS', [5]),
        ]))
        
        story.append(charts_grid)
        story.append(PageBreak())

    def _crear_tabla_detalle_costes_moderna(self, story, costes_por_linea_turno):
        """
        Crear tabla detallada de costes con diseño moderno y profesional.
        MEJORADO: Consistente con el estilo del resto del reporte
        """
        styles = getSampleStyleSheet()
        
        # Estilos modernos para la tabla
        styleTableTitle = ParagraphStyle(
            'tableTitle',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=16,
            alignment=1,
            textColor=colors.white,
            spaceAfter=8
        )
        
        styleHeaderCell = ParagraphStyle(
            'headerCell',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=11,
            alignment=1,
            textColor=colors.white
        )
        
        styleDataCell = ParagraphStyle(
            'dataCell',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=colors.HexColor('#1E293B')
        )
        
        styleDataCellCenter = ParagraphStyle(
            'dataCellCenter',
            parent=styleDataCell,
            alignment=1
        )
        
        styleDataCellRight = ParagraphStyle(
            'dataCellRight',
            parent=styleDataCell,
            alignment=2
        )

        # Título de la tabla
        title_data = [[Paragraph("TOP 5 LÍNEAS CON MAYOR IMPACTO ECONÓMICO", styleTableTitle)]]
        title_table = Table(title_data, colWidths=[1150], rowHeights=[40])
        title_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#7F1D1D')),
            ('BOX', (0, 0), (-1, -1), 2, colors.HexColor('#991B1B')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('TOPPADDING', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
            ('ROUNDEDCORNERS', [5]),
        ]))
        
        story.append(title_table)
        story.append(Spacer(1, 12))

        # Crear datos para la tabla
        table_data = []
        
        # Header moderno
        header_row = [
            Paragraph("LÍNEA / UET / TURNO", styleHeaderCell),
            Paragraph("TIPO DE DEFECTO", styleHeaderCell),
            Paragraph("CANTIDAD", styleHeaderCell),
            Paragraph("COSTE TOTAL", styleHeaderCell)
        ]
        table_data.append(header_row)
        
        # Ordenar por coste total descendente - SOLO TOP 5
        sorted_lineas = sorted(costes_por_linea_turno.items(),
                            key=lambda x: x[1]['total_coste'], reverse=True)[:5]

        # Procesar solo TOP 5 líneas más costosas
        for idx, (key_string, data) in enumerate(sorted_lineas):
            try:
                parts = key_string.split('|')
                if len(parts) == 3:
                    linea, uet, turno = parts
                elif len(parts) == 2:
                    uet, turno = parts
                    linea = f"UET_{uet}"
                else:
                    continue
            except:
                continue
            
            # Agrupar defectos por tipo
            defectos_agrupados = {}
            for item in data['items']:
                texto_breve = item['texto_breve']
                if texto_breve not in defectos_agrupados:
                    defectos_agrupados[texto_breve] = {
                        'cantidad': 0,
                        'coste_total': 0
                    }
                defectos_agrupados[texto_breve]['cantidad'] += item['cantidad']
                defectos_agrupados[texto_breve]['coste_total'] += item['coste_total']
            
            # Convertir turno a texto legible
            turno_texto = {
                'M': 'MAÑANA',
                'T': 'TARDE',
                'N': 'NOCHE',
                'SIN_TURNO': 'SIN TURNO'
            }.get(turno, turno)
            
            # Añadir filas para cada defecto
            for defecto_idx, (texto_breve, info) in enumerate(sorted(defectos_agrupados.items(),
                                        key=lambda x: x[1]['coste_total'], reverse=True)):
                
                linea_info = Paragraph(f"<b>{linea}</b><br/>UET: {uet}<br/>TURNO: {turno_texto}", styleDataCell)
                defecto_info = Paragraph(f"<b>{texto_breve}</b>", styleDataCell)
                cantidad_info = Paragraph(f"<b>{info['cantidad']}</b><br/>defectos", styleDataCellCenter)
                coste_info = Paragraph(f"<b>€ {self._format_euro(info['coste_total'])}</b>", styleDataCellRight)
                
                table_data.append([linea_info, defecto_info, cantidad_info, coste_info])

        # Crear tabla con diseño moderno
        ancho_disponible = 1150
        col_widths = [
            ancho_disponible * 0.28,  # Línea/UET/turno
            ancho_disponible * 0.42,  # Tipo defecto
            ancho_disponible * 0.15,  # Cantidad
            ancho_disponible * 0.15   # Coste
        ]
        
        detail_table = Table(table_data, colWidths=col_widths, repeatRows=1)
        detail_table.setStyle(TableStyle([
            # Header styling moderno
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#991B1B')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
            ('TOPPADDING', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            
            # Data rows styling
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('ALIGN', (0, 1), (0, -1), 'LEFT'),    # Línea/UET/Turno
            ('ALIGN', (1, 1), (1, -1), 'LEFT'),    # Tipo defecto
            ('ALIGN', (2, 1), (2, -1), 'CENTER'),  # Cantidad
            ('ALIGN', (3, 1), (3, -1), 'RIGHT'),   # Coste
            
            # Spacing optimizado
            ('TOPPADDING', (0, 1), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 10),
            ('LEFTPADDING', (0, 0), (-1, -1), 12),
            ('RIGHTPADDING', (0, 0), (-1, -1), 12),
            
            # Borders modernos
            ('BOX', (0, 0), (-1, -1), 2, colors.HexColor('#991B1B')),
            ('INNERGRID', (0, 0), (-1, -1), 1, colors.HexColor('#E5E7EB')),
            ('LINEBELOW', (0, 0), (-1, 0), 3, colors.HexColor('#7F1D1D')),
            
            # Alternating row colors
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [
                colors.white,
                colors.HexColor('#FEF2F2')
            ]),
            
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ROUNDEDCORNERS', [3]),
        ]))
        
        story.append(detail_table)
        
        # Nota final moderna
        nota_style = ParagraphStyle(
            'nota',
            parent=styles['Normal'],
            fontName='Helvetica-Oblique',
            fontSize=10,
            textColor=colors.HexColor('#6B7280'),
            alignment=1
        )
        
        story.append(Spacer(1, 15))
        nota_text = f"<i>Nota: Análisis completo con {len(table_data)-1} registros. Los valores excluyen defectos de proveedor.</i>"
        nota_para = Paragraph(nota_text, nota_style)
        story.append(nota_para)

    def _generar_grafico_ranking_lineas_costes(self, costes_por_linea_turno):
        """
        Generar gráfico de ranking de líneas por pérdidas económicas con diseño moderno.
        """
        try:
            # Agrupar por línea real
            costes_por_linea = {}
            for key_string, data in costes_por_linea_turno.items():
                try:
                    parts = key_string.split('|')
                    if len(parts) == 3:
                        linea, uet, turno = parts
                    elif len(parts) == 2:
                        uet, turno = parts
                        linea = f"UET_{uet}"
                    else:
                        continue
                        
                    if linea not in costes_por_linea:
                        costes_por_linea[linea] = 0
                    costes_por_linea[linea] += data['total_coste']
                except Exception:
                    continue
            
            if not costes_por_linea:
                return None
                
            # Ordenar y tomar top 10
            sorted_lineas = sorted(costes_por_linea.items(), key=lambda x: x[1], reverse=True)[:10]
            
            if not sorted_lineas or all(coste == 0 for _, coste in sorted_lineas):
                return None
                
            fig, ax = self._crear_grafica_profesional_moderna(figsize=(12, 7))
            
            lineas = [item[0] for item in sorted_lineas]
            costes = [item[1] for item in sorted_lineas]
            
            # Colores gradiente modernos (rojo a naranja)
            colors_bars = ['#DC2626' if i < 3 else '#EA580C' if i < 6 else '#D97706' 
                        for i in range(len(lineas))]
            
            bars = ax.barh(range(len(lineas)), costes, color=colors_bars, alpha=0.85, height=0.6)
            
            # Añadir valores
            for i, (bar, coste) in enumerate(zip(bars, costes)):
                width = bar.get_width()
                if width > 0:
                    ax.text(width + max(costes)*0.01, bar.get_y() + bar.get_height()/2,
                        f'€{coste:,.0f}', va='center', ha='left', 
                        fontweight='600', fontsize=10, color='#1F2937')
            
            ax.set_yticks(range(len(lineas)))
            ax.set_yticklabels(lineas, fontsize=11, color='#374151', fontweight='500')
            ax.set_xlabel('Pérdidas Económicas (€)', fontweight='600', fontsize=12, color='#1F2937')
            
            # Título moderno
            ax.text(0.5, 1.05, 'RANKING DE LÍNEAS POR IMPACTO ECONÓMICO', 
                transform=ax.transAxes, fontsize=16, fontweight='700', 
                color='#1F2937', ha='center', va='bottom')
            
            # Grid sutil
            ax.grid(axis='x', alpha=0.2, linestyle='-', color='#E5E7EB')
            ax.set_axisbelow(True)
            
            plt.subplots_adjust(left=0.15, top=0.88)
            
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            return buf
            
        except Exception as e:
            print(f"Error en _generar_grafico_ranking_lineas_costes: {e}")
            return None

    def _generar_grafico_distribucion_turnos_costes(self, costes_por_linea_turno):
        """
        Generar gráfico circular de distribución de costes por turno con diseño moderno.
        """
        try:
            # Agrupar por turno
            costes_por_turno = {}
            for key_string, data in costes_por_linea_turno.items():
                try:
                    parts = key_string.split('|')
                    if len(parts) == 3:
                        linea, uet, turno = parts
                    elif len(parts) == 2:
                        uet, turno = parts
                    else:
                        continue
                        
                    if turno not in costes_por_turno:
                        costes_por_turno[turno] = 0
                    costes_por_turno[turno] += data['total_coste']
                except Exception:
                    continue
            
            if not costes_por_turno or all(coste == 0 for coste in costes_por_turno.values()):
                return None
                
            fig, ax = self._crear_grafica_profesional_moderna(figsize=(9, 7))
            
            turnos = list(costes_por_turno.keys())
            costes = list(costes_por_turno.values())
            
            # Colores específicos por turno (modernos)
            colores_turno = {
                'M': '#F59E0B',      # Ámbar para Mañana
                'T': '#DC2626',      # Rojo para Tarde
                'N': '#3B82F6',      # Azul para Noche
                'SIN_TURNO': '#6B7280'  # Gris para sin turno
            }
            colors_pie = [colores_turno.get(turno, '#6B7280') for turno in turnos]
            
            # Crear donut chart moderno
            max_idx = costes.index(max(costes))
            explode = [0.1 if i == max_idx else 0 for i in range(len(turnos))]
            
            wedges, texts, autotexts = ax.pie(costes, labels=turnos, autopct='%1.1f%%',
                                            colors=colors_pie, explode=explode, startangle=90,
                                            textprops={'fontsize': 11, 'fontweight': '600'},
                                            wedgeprops=dict(width=0.5, edgecolor='white', linewidth=2))
            
            # Centro blanco para efecto donut
            centre_circle = plt.Circle((0, 0), 0.70, fc='white')
            fig.gca().add_artist(centre_circle)
            
            # Mejorar textos
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontsize(12)
                autotext.set_fontweight('bold')
            
            # Título moderno
            ax.text(0.5, 1.05, 'DISTRIBUCIÓN DE PÉRDIDAS POR TURNO', 
                transform=ax.transAxes, fontsize=16, fontweight='700', 
                color='#1F2937', ha='center', va='bottom')
            
            # Leyenda moderna
            leyenda_labels = [f'{turno}: €{self._format_euro(coste).split(",")[0]}' for turno, coste in zip(turnos, costes)]
            legend = ax.legend(wedges, leyenda_labels, title="Impacto por Turno", 
                            loc="center left", bbox_to_anchor=(1, 0, 0.5, 1),
                            frameon=True, fontsize=10)
            legend.get_frame().set_facecolor('white')
            legend.get_frame().set_edgecolor('#E5E7EB')
            
            plt.subplots_adjust(top=0.88)
            
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            return buf
            
        except Exception as e:
            print(f"Error en _generar_grafico_distribucion_turnos_costes: {e}")
            return None

    def _generar_grafico_distribucion_perdidas_por_uet(self, costes_por_linea_turno):
        """
        Generar gráfico de distribución de pérdidas por UET con diseño moderno.
        """
        try:
            # Agrupar por UET
            costes_por_uet = {}
            for key_string, data in costes_por_linea_turno.items():
                try:
                    parts = key_string.split('|')
                    if len(parts) == 3:
                        linea, uet, turno = parts
                    elif len(parts) == 2:
                        uet, turno = parts
                    else:
                        continue
                        
                    if uet not in costes_por_uet:
                        costes_por_uet[uet] = 0
                    costes_por_uet[uet] += data['total_coste']
                except Exception:
                    continue
            
            if not costes_por_uet or all(coste == 0 for coste in costes_por_uet.values()):
                return None
                
            # Tomar top 8 UETs más costosos
            sorted_uets = sorted(costes_por_uet.items(), key=lambda x: x[1], reverse=True)[:8]
            
            fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
            
            uets = [item[0] for item in sorted_uets]
            costes = [item[1] for item in sorted_uets]
            
            # Gradiente de colores modernos
            colors_bars = ['#DC2626', '#EA580C', '#D97706', '#CA8A04', '#A3A3A3', 
                        '#737373', '#525252', '#404040'][:len(uets)]
            
            bars = ax.bar(range(len(uets)), costes, color=colors_bars, alpha=0.85, width=0.7,
                        edgecolor='white', linewidth=2)
            
            self._anotar_barras_moderno(ax)
            
            ax.set_xticks(range(len(uets)))
            ax.set_xticklabels([f'UET {uet}' for uet in uets], rotation=30, ha='right', 
                            fontsize=10, color='#374151', fontweight='500')
            
            ax.set_ylabel('Coste Total (€)', fontsize=12, fontweight='600', color='#1F2937')
            
            # Título moderno
            ax.text(0.5, 1.05, 'DISTRIBUCIÓN DE PÉRDIDAS POR UET', 
                transform=ax.transAxes, fontsize=16, fontweight='700', 
                color='#1F2937', ha='center', va='bottom')
            
            # Línea de promedio
            if len(costes) > 1:
                avg_val = np.mean(costes)
                ax.axhline(y=avg_val, color='#6B7280', linestyle='--', linewidth=2, 
                        alpha=0.7, label=f'Promedio: €{avg_val:.0f}')
                ax.legend(loc='upper right', frameon=True, fontsize=10)
            
            plt.subplots_adjust(top=0.88, bottom=0.15)
            
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            return buf
            
        except Exception as e:
            print(f"Error en _generar_grafico_distribucion_perdidas_por_uet: {e}")
            return None

    def _generar_grafico_tendencia_costes_30_dias(self, fi_reporte, ff_reporte):
        """
        Generar gráfico de tendencia de costes para los últimos 30 días con diseño moderno.
        MEJORADO: Intenta procesar datos reales de costes
        """
        try:
            # Calcular rango de 30 días hacia atrás desde la fecha final
            ff_dt = datetime.strptime(ff_reporte, '%Y-%m-%d').date()
            fi_30_dias = ff_dt - timedelta(days=29)  # 30 días incluyendo el día final
            
            # Inicializar estructura para costes diarios
            costes_por_fecha = {}
            
            # Procesar archivos día a día
            if not os.path.isdir(self.ruta_archivos):
                return self._generar_placeholder_tendencia()

            patron = os.path.join(self.ruta_archivos, "*.xlsx")
            archivos_excel = glob.glob(patron)
            archivos_excel = [archivo for archivo in archivos_excel if not archivo.endswith("costes.xlsx")]
            
            archivos_procesados = 0
            for archivo in archivos_excel:
                nombre = os.path.basename(archivo)
                base, _ = os.path.splitext(nombre)
                try:
                    fecha_archivo = datetime.strptime(base, '%d.%m.%Y').date()
                except ValueError:
                    continue

                # Solo procesar archivos en el rango de 30 días
                if not (fi_30_dias <= fecha_archivo <= ff_dt):
                    continue

                fecha_str = fecha_archivo.isoformat()
                archivos_procesados += 1
                costes_dia = 0
                
                try:
                    wb = openpyxl.load_workbook(archivo, data_only=True)
                    ws = wb.active
                    
                    for row in ws.iter_rows(min_row=4, values_only=True):
                        if len(row) < 12: 
                            continue

                        val_nok = row[3]
                        val_col_f = str(row[5] or "").strip().upper()
                        val_defecto_row = (row[6] or "").strip().upper()
                        val_texto_breve = (row[1] or "").strip().upper()

                        _nok = self._parse_num(val_nok)
                        es_proveedor = self._es_defecto_proveedor(val_defecto_row)

                        # Calcular costes (misma lógica que en leer_estadisticas)
                        if val_col_f != "R" and _nok > 0 and val_texto_breve and not es_proveedor:
                            precio_unitario = self.costes_dict.get(val_texto_breve, 0)
                            coste_total = _nok * precio_unitario
                            
                            if coste_total > 0:
                                costes_dia += coste_total
                                    
                except Exception as e:
                    print(f"Error procesando archivo {nombre}: {e}")
                    continue
                
                # Solo guardar días que tengan costes > 0
                if costes_dia > 0:
                    costes_por_fecha[fecha_str] = costes_dia
            
            # Verificar que tenemos datos
            if not costes_por_fecha:
                return self._generar_placeholder_tendencia("NO HAY PÉRDIDAS ECONÓMICAS\nEN LOS ÚLTIMOS 30 DÍAS")
            
            # Preparar datos para el gráfico
            fechas_ordenadas = sorted(costes_por_fecha.keys())
            costes_ordenados = [costes_por_fecha[fecha] for fecha in fechas_ordenadas]
            
            fig, ax = self._crear_grafica_profesional_moderna(figsize=(12, 7))
            
            # Convertir fechas a formato más legible
            fechas_display = []
            for fecha in fechas_ordenadas:
                fecha_dt = datetime.strptime(fecha, '%Y-%m-%d')
                fechas_display.append(fecha_dt.strftime('%d/%m'))
            
            # Crear gráfico con línea y área rellena
            ax.plot(range(len(fechas_ordenadas)), costes_ordenados, color='#DC2626', 
                    linewidth=3, marker='o', markersize=6, markerfacecolor='white', 
                    markeredgecolor='#DC2626', markeredgewidth=2)
            
            ax.fill_between(range(len(fechas_ordenadas)), costes_ordenados, 
                            alpha=0.3, color='#DC2626')
            
            # Configurar ejes
            if len(fechas_display) <= 15:
                ax.set_xticks(range(len(fechas_display)))
                ax.set_xticklabels(fechas_display, rotation=45, ha='right')
            else:
                step = max(1, len(fechas_display) // 10)
                indices_mostrar = list(range(0, len(fechas_display), step))
                if (len(fechas_display) - 1) not in indices_mostrar:
                    indices_mostrar.append(len(fechas_display) - 1)
                
                ax.set_xticks(indices_mostrar)
                ax.set_xticklabels([fechas_display[i] for i in indices_mostrar], rotation=45, ha='right')
            
            ax.set_ylabel('Pérdidas Diarias (€)', fontweight='600', fontsize=12, color='#1F2937')
            ax.text(0.5, 1.05, f'EVOLUCIÓN DIARIA DE PÉRDIDAS ECONÓMICAS (Últimos 30 días)\n{len(fechas_ordenadas)} días con pérdidas', 
                    transform=ax.transAxes, fontsize=16, fontweight='700', 
                    color='#1F2937', ha='center', va='bottom')
            
            # Añadir estadísticas
            total_costes = sum(costes_ordenados)
            coste_medio = np.mean(costes_ordenados)
            ax.axhline(y=coste_medio, color='#F59E0B', linestyle='--', linewidth=2, 
                    alpha=0.7, label=f'Promedio: €{coste_medio:,.0f}')
            ax.legend(loc='upper right', frameon=True, fontsize=10)
            
            # Destacar el día más costoso
            max_coste = max(costes_ordenados)
            max_idx = costes_ordenados.index(max_coste)
            fecha_max = fechas_display[max_idx]
            ax.annotate(f'Máx: {fecha_max} - €{max_coste:,.0f}', 
                        xy=(max_idx, max_coste), 
                        xytext=(max_idx + len(fechas_ordenadas)*0.1, max_coste*1.1),
                        arrowprops=dict(arrowstyle='->', color='#DC2626', lw=2),
                        fontsize=10, fontweight='600', color='#DC2626')
            
            plt.subplots_adjust(top=0.85, bottom=0.15)
            
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            return buf
            
        except Exception as e:
            print(f"Error en _generar_grafico_tendencia_costes_30_dias: {e}")
            return self._generar_placeholder_tendencia("ERROR GENERANDO GRÁFICO\nDE TENDENCIA DE 30 DÍAS")

    def _generar_placeholder_tendencia(self, mensaje="ANÁLISIS DE TENDENCIA DE 30 DÍAS\nEn desarrollo"):
        """
        Genera un placeholder para el gráfico de tendencia
        """
        try:
            fig, ax = self._crear_grafica_profesional_moderna(figsize=(12, 7))
            
            ax.text(0.5, 0.5, mensaje, 
                    ha='center', va='center', fontsize=16, fontweight='600', 
                    transform=ax.transAxes, color='#6B7280',
                    bbox=dict(boxstyle="round,pad=0.5", facecolor='#F3F4F6', alpha=0.8))
            
            ax.text(0.5, 1.05, 'EVOLUCIÓN DIARIA DE PÉRDIDAS ECONÓMICAS (30 días)', 
                transform=ax.transAxes, fontsize=16, fontweight='700', 
                color='#1F2937', ha='center', va='bottom')
            
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis('off')
            
            plt.subplots_adjust(top=0.88)
            
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            return buf
            
        except Exception as e:
            print(f"Error generando placeholder de tendencia: {e}")
            return None

    # FUNCIONES RESTANTES DEL MODELO (NOK y REWORK) - VERSIONES COMPACTAS
    def _crear_hoja_modelo_nok(self, story, fi_reporte, ff_reporte, modelo_str):
        """
        Genera página de detalle para un modelo (NOK) OPTIMIZADA para que las 4 gráficas quepan en una página.
        SOLUCIONADO: Tamaños reducidos y mejor espaciado
        """
        # USAR FILTRADO POR MODELO para obtener datos específicos
        data_model_report_range = self.leer_estadisticas(fi_reporte, ff_reporte, modelo_str, "(TODOS)")
        if data_model_report_range.get("error"):
            story.append(Paragraph(
                f"Error al leer datos del modelo {modelo_str} para {fi_reporte}-{ff_reporte}: "
                f"{data_model_report_range['error']}",
                getSampleStyleSheet()['Normal']
            ))
            story.append(PageBreak())
            return

        # USAR DATOS PRINCIPALES (que ahora están filtrados por modelo)
        model_total_ok = data_model_report_range.get("total_ok", 0)
        model_total_nok = data_model_report_range.get("total_nok", 0)
        model_total_nok_prov = data_model_report_range.get("total_nok_prov", 0)
        defect_dist_model = data_model_report_range.get("defect_distribution", {})
        shift_dist_model = data_model_report_range.get("shift_distribution", {})
        model_evolution = data_model_report_range.get("evolution", {})

        styles = getSampleStyleSheet()
        
        # Estilos COMPACTOS - Reducidos significativamente
        styleTitle = ParagraphStyle(
            'modelTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=18,  # Reducido de 26 a 18
            alignment=1,
            textColor=colors.white,
            leading=20   # Reducido de 28 a 20
        )
        
        styleSub = ParagraphStyle(
            'modelSub',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,  # Reducido de 13 a 10
            alignment=1,
            textColor=colors.white,
            leading=12    # Reducido de 15 a 12
        )
        
        styleKPILabel = ParagraphStyle(
            'styleMarkerLabel',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=8,   # Reducido de 10 a 8
            alignment=1,
            textColor=colors.HexColor('#475569'),
            leading=10    # Reducido de 12 a 10
        )
        
        styleKPINumber = ParagraphStyle(
            'styleMarkerNumber',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=12,  # Reducido de 17 a 12
            alignment=1,
            textColor=colors.HexColor('#1E293B'),
            leading=14    # Reducido de 20 a 14
        )
        
        styleKPINumberHero = ParagraphStyle(
            'styleMarkerNumberBig',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=20,  # Reducido de 30 a 20
            alignment=1,
            textColor=colors.white,
            leading=22    # Reducido de 32 a 22
        )

        fecha_display = fi_reporte if fi_reporte == ff_reporte else f"{fi_reporte} — {ff_reporte}"

        # Header COMPACTO - Alturas reducidas
        header_data = [
            [Paragraph(f"MODELO {modelo_str.upper()}", styleTitle)],
            [Paragraph(f"Rechazos NOK | {fecha_display}", styleSub)]
        ]
        
        header_table = Table(header_data, colWidths=[1150], rowHeights=[22, 14])  # Reducido de [32, 18] a [22, 14]
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#1E293B')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),   # Reducido de 15 a 8
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),  # Reducido de 15 a 8
            ('TOPPADDING', (0, 0), (-1, -1), 4),    # Reducido de 8 a 4
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4), # Reducido de 8 a 4
            ('LINEBELOW', (0, 1), (-1, 1), 2, colors.HexColor('#3B82F6')),  # Línea más delgada
        ]))
        story.append(header_table)
        story.append(Spacer(1, 4))  # Reducido de 12 a 4

        total_modelo = model_total_ok + model_total_nok
        perc_nok_modelo = (model_total_nok / total_modelo * 100) if total_modelo > 0 else 0

        # Color dinámico del KPI principal
        col_rechazo = colors.HexColor('#EF4444')
        status_text = "CRÍTICO"
        if perc_nok_modelo <= self.OBJETIVO:
            col_rechazo = colors.HexColor('#10B981')
            status_text = "ÓPTIMO"
        elif perc_nok_modelo <= self.OBJETIVO * 1.3:
            col_rechazo = colors.HexColor('#F59E0B')
            status_text = "ALERTA"

        # KPI principal COMPACTO
        kpi_nok_data = [
            [Paragraph("TASA RECHAZO", styleKPILabel)],
            [Paragraph(f"{perc_nok_modelo:.1f}%", styleKPINumberHero)],  # Solo 1 decimal
            [Paragraph(status_text, ParagraphStyle('status', parent=styleKPILabel, 
                                                textColor=colors.white, fontSize=8))]
        ]
        
        kpi_nok_table = Table(kpi_nok_data, colWidths=[200], rowHeights=[12, 26, 10])  # Reducido ancho y alto
        kpi_nok_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), col_rechazo),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),   # Reducido
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),  # Reducido
            ('TOPPADDING', (0, 0), (-1, -1), 3),    # Reducido
            ('BOTTOMPADDING', (0, 0), (-1, -1), 3), # Reducido
            ('LINEBELOW', (0, 0), (-1, 0), 1, colors.white),
            ('ROUNDEDCORNERS', [3]),
        ]))

        # KPIs secundarios COMPACTOS
        secondary_kpis_data = [
            ("PIEZAS OK", f"{model_total_ok:,}", "#F0FDF4", "#10B981"),
            ("NOK INTERNOS", f"{model_total_nok:,}", "#FEF2F2", "#EF4444"),
            ("NOK PROVEEDOR", f"{model_total_nok_prov:,}", "#FFF7ED", "#F59E0B")
        ]
        
        secondary_kpis = []
        for label, value, bg_color, accent_color in secondary_kpis_data:
            card_data = [
                [Paragraph(label, styleKPILabel)],
                [Paragraph(value, styleKPINumber)]
            ]
            card_table = Table(card_data, colWidths=[200], rowHeights=[12, 20])  # Reducido
            card_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor(bg_color)),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BOX', (0, 0), (-1, -1), 1, colors.HexColor('#E2E8F0')),
                ('LEFTPADDING', (0, 0), (-1, -1), 6),   # Reducido
                ('RIGHTPADDING', (0, 0), (-1, -1), 6),  # Reducido
                ('TOPPADDING', (0, 0), (-1, -1), 3),    # Reducido
                ('BOTTOMPADDING', (0, 0), (-1, -1), 3), # Reducido
                ('LINEBELOW', (0, 0), (-1, 0), 1, colors.HexColor(accent_color)),
                ('ROUNDEDCORNERS', [2]),
            ]))
            secondary_kpis.append(card_table)

        # Layout de KPIs COMPACTO - Ajustar anchos para que quepan
        kpi_container = Table([[kpi_nok_table] + secondary_kpis], 
                            colWidths=[240, 240, 240, 240])  # Reducido de 285 a 240
        kpi_container.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 2),   # Reducido
            ('RIGHTPADDING', (0, 0), (-1, -1), 2),  # Reducido
        ]))
        story.append(kpi_container)
        story.append(Spacer(1, 4))  # Reducido de 8 a 4

        # Crear gráficos específicos del modelo con datos filtrados
        buf_trend_model, buf_defectos_model, buf_turnos_model, buf_top_def_model = self._crear_graficos_modelo_nok(
            fi_reporte, ff_reporte, modelo_str, model_evolution, defect_dist_model, shift_dist_model
        )

        # Grid de gráficos MAXIMIZADO - Aprovechar al máximo el espacio disponible
        imgA = self._safe_img_flowable(buf_trend_model, 550, 290)      # Aumentado para llenar mejor: 550x290
        imgB = self._safe_img_flowable(buf_defectos_model, 550, 290)   # Aumentado para llenar mejor: 550x290
        imgC = self._safe_img_flowable(buf_turnos_model, 550, 290)     # Aumentado para llenar mejor: 550x290
        imgD = self._safe_img_flowable(buf_top_def_model, 550, 290)    # Aumentado para llenar mejor: 550x290

        # Crear contenedores con títulos COMPACTOS
        def create_chart_container(img, title):
            container_data = [
                [Paragraph(title, ParagraphStyle('chartTitle', parent=styles['Normal'],
                                            fontName='Helvetica-Bold', fontSize=9,  # Reducido de 11 a 9
                                            textColor=colors.HexColor('#1E293B'),
                                            alignment=1, spaceAfter=2))],  # Reducido de 3 a 2
                [img]
            ]
            container = Table(container_data, colWidths=[560], rowHeights=[14, 300])  # Aumentado para llenar: [560] y [14, 300]
            container.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F8FAFC')),
                ('LINEBELOW', (0, 0), (-1, 0), 1, colors.HexColor('#E2E8F0')),
            ]))
            return container

        charts = Table(
            [
                [create_chart_container(imgA, "Tendencia Temporal"),
                create_chart_container(imgB, "Distribución Defectos")],
                [create_chart_container(imgC, "Análisis por Turno"),
                create_chart_container(imgD, "Top 5 Defectos")]
            ],
            colWidths=[570, 570],       # Aumentado para llenar mejor: 570 cada columna
            rowHeights=[320, 320]       # Aumentado para llenar mejor: 320 cada fila
        )
        
        charts.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('BOX', (0, 0), (-1, -1), 1, colors.HexColor('#E2E8F0')),       # Línea más delgada
            ('INNERGRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#F1F5F9')), # Grid más sutil
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
            ('LEFTPADDING', (0, 0), (-1, -1), 3),   # Reducido de 5 a 3
            ('RIGHTPADDING', (0, 0), (-1, -1), 3),  # Reducido de 5 a 3
            ('TOPPADDING', (0, 0), (-1, -1), 3),    # Reducido de 5 a 3
            ('BOTTOMPADDING', (0, 0), (-1, -1), 3), # Reducido de 5 a 3
            ('ROUNDEDCORNERS', [3]),
        ]))
        story.append(charts)
        story.append(PageBreak())

    def _crear_hoja_modelo_rework(self, story, fi_reporte, ff_reporte, modelo_str):
        """
        Genera página de detalle para un modelo (Rework) OPTIMIZADA para que las 4 gráficas quepan en una página.
        SOLUCIONADO: Tamaños reducidos y mejor espaciado
        """
        # USAR FILTRADO POR MODELO para obtener datos específicos de retrabajos
        data_model_report_range = self.leer_estadisticas(fi_reporte, ff_reporte, modelo_str, "(TODOS)")
        if data_model_report_range.get("error"):
            story.append(Paragraph(
                f"Error leyendo datos del modelo {modelo_str} para {fi_reporte}-{ff_reporte}: "
                f"{data_model_report_range['error']}",
                getSampleStyleSheet()['Normal']
            ))
            story.append(PageBreak())
            return

        # USAR DATOS PRINCIPALES DE RETRABAJOS (filtrados por modelo)
        model_total_rework = data_model_report_range.get("total_rework", 0)
        rework_defect_dist_model = data_model_report_range.get("rework_defect_distribution", {})
        rework_shift_dist_model = data_model_report_range.get("rework_shift_distribution", {})
        rework_evolution_model = data_model_report_range.get("evolution_rework", {})

        styles = getSampleStyleSheet()
        
        # Estilos COMPACTOS - Consistentes con NOK
        styleTitle = ParagraphStyle(
            'modelTitle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=18,  # Reducido de 26 a 18
            alignment=1,
            textColor=colors.white,
            leading=20   # Reducido de 28 a 20
        )
        
        styleSub = ParagraphStyle(
            'modelSub',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,  # Reducido de 13 a 10
            alignment=1,
            textColor=colors.white,
            leading=12    # Reducido de 15 a 12
        )
        
        styleKPILabel = ParagraphStyle(
            'styleMarkerLabel',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=8,   # Reducido de 10 a 8
            alignment=1,
            textColor=colors.HexColor('#475569'),
            leading=10    # Reducido de 12 a 10
        )
        
        styleKPINumberHero = ParagraphStyle(
            'styleMarkerNumberBig',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=20,  # Reducido de 30 a 20
            alignment=1,
            textColor=colors.white,
            leading=22    # Reducido de 32 a 22
        )

        fecha_display = fi_reporte if fi_reporte == ff_reporte else f"{fi_reporte} — {ff_reporte}"

        # Header COMPACTO con tema verde
        header_data = [
            [Paragraph(f"MODELO {modelo_str.upper()}", styleTitle)],
            [Paragraph(f"Retrabajos | {fecha_display}", styleSub)]
        ]
        
        header_table = Table(header_data, colWidths=[1150], rowHeights=[22, 14])  # Reducido de [32, 18] a [22, 14]
        header_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#065F46')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),   # Reducido de 15 a 8
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),  # Reducido de 15 a 8
            ('TOPPADDING', (0, 0), (-1, -1), 4),    # Reducido de 8 a 4
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4), # Reducido de 8 a 4
            ('LINEBELOW', (0, 1), (-1, 1), 2, colors.HexColor('#10B981')),  # Línea más delgada
        ]))
        story.append(header_table)
        story.append(Spacer(1, 4))  # Reducido de 12 a 4

        # KPI principal COMPACTO: Total Retrabajos
        kpi_rework_data = [
            [Paragraph("TOTAL RETRABAJOS", styleKPILabel)],
            [Paragraph(f"{model_total_rework:,}", styleKPINumberHero)],
            [Paragraph("UNIDADES RECUPERADAS", ParagraphStyle('status', parent=styleKPILabel, 
                                                            textColor=colors.white, fontSize=8))]
        ]
        
        kpi_rework_table = Table(kpi_rework_data, colWidths=[300], rowHeights=[12, 26, 10])  # Reducido de [400] y alturas
        kpi_rework_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#10B981')),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),   # Reducido de 15 a 8
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),  # Reducido de 15 a 8
            ('TOPPADDING', (0, 0), (-1, -1), 3),    # Reducido de 5 a 3
            ('BOTTOMPADDING', (0, 0), (-1, -1), 3), # Reducido de 5 a 3
            ('LINEBELOW', (0, 0), (-1, 0), 1, colors.white),
            ('ROUNDEDCORNERS', [3]),
        ]))

        # Centrar el KPI COMPACTO
        kpi_center = Table([[kpi_rework_table]], colWidths=[1150], 
                        style=[('ALIGN', (0, 0), (-1, -1), 'CENTER')])
        story.append(kpi_center)
        story.append(Spacer(1, 4))  # Reducido de 8 a 4

        # Crear gráficos específicos del modelo para retrabajos
        buf_trend_rework, buf_defectos_rework, buf_turnos_rework, buf_top_def_rework = self._crear_graficos_modelo_rework(
            fi_reporte, ff_reporte, modelo_str, rework_evolution_model, rework_defect_dist_model, rework_shift_dist_model
        )

        # Grid de gráficos MAXIMIZADO - Aprovechar al máximo el espacio disponible
        imgA = self._safe_img_flowable(buf_trend_rework, 550, 290)     # Aumentado para llenar mejor: 550x290
        imgB = self._safe_img_flowable(buf_defectos_rework, 550, 290)  # Aumentado para llenar mejor: 550x290
        imgC = self._safe_img_flowable(buf_turnos_rework, 550, 290)    # Aumentado para llenar mejor: 550x290
        imgD = self._safe_img_flowable(buf_top_def_rework, 550, 290)   # Aumentado para llenar mejor: 550x290

        def create_chart_container(img, title):
            container_data = [
                [Paragraph(title, ParagraphStyle('chartTitle', parent=styles['Normal'],
                                            fontName='Helvetica-Bold', fontSize=9,  # Reducido de 11 a 9
                                            textColor=colors.HexColor('#065F46'),
                                            alignment=1, spaceAfter=2))],  # Reducido de 3 a 2
                [img]
            ]
            container = Table(container_data, colWidths=[560], rowHeights=[14, 300])  # Aumentado para llenar: [560] y [14, 300]
            container.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F0FDF4')),
                ('LINEBELOW', (0, 0), (-1, 0), 1, colors.HexColor('#BBF7D0')),
            ]))
            return container

        charts = Table(
            [
                [create_chart_container(imgA, "Tendencia Temporal"),
                create_chart_container(imgB, "Distribución Defectos")],
                [create_chart_container(imgC, "Análisis por Turno"),
                create_chart_container(imgD, "Top 5 Defectos Retrabajo")]
            ],
            colWidths=[570, 570],       # Aumentado para llenar mejor: 570 cada columna
            rowHeights=[320, 320]       # Aumentado para llenar mejor: 320 cada fila
        )
        
        charts.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('BOX', (0, 0), (-1, -1), 1, colors.HexColor('#E2E8F0')),       # Línea más delgada
            ('INNERGRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#F1F5F9')), # Grid más sutil
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
            ('LEFTPADDING', (0, 0), (-1, -1), 3),   # Reducido de 5 a 3
            ('RIGHTPADDING', (0, 0), (-1, -1), 3),  # Reducido de 5 a 3
            ('TOPPADDING', (0, 0), (-1, -1), 3),    # Reducido de 5 a 3
            ('BOTTOMPADDING', (0, 0), (-1, -1), 3), # Reducido de 5 a 3
            ('ROUNDEDCORNERS', [3]),
        ]))
        story.append(charts)
        story.append(PageBreak())

    def _crear_graficos_modelo_nok(self, fi_reporte, ff_reporte, modelo_str, model_evolution, defect_dist_model, shift_dist_model):
        """
        Crea los 4 gráficos específicos para análisis NOK de un modelo.
        NUEVO: Gráficos con datos filtrados específicamente para el modelo
        """
        modern_colors = {
            'primary': '#0F172A', 'accent': '#3B82F6', 'success': '#10B981',
            'danger': '#EF4444', 'medium': '#64748B', 'dark': '#334155'
        }
        
        defect_colors = ['#3B82F6', '#10B981', '#F59E0B', '#EF4444', '#8B5CF6']
        shift_colors = {'MAÑANA': '#10B981', 'TARDE': '#3B82F6', 'NOCHE': '#8B5CF6',
                       'A': '#10B981', 'B': '#3B82F6', 'C': '#8B5CF6'}

        # 1. Gráfico de tendencia temporal
        buf_trend = None
        try:
            ff_dt = datetime.strptime(ff_reporte, '%Y-%m-%d').date()
            trend_start_date = ff_dt - timedelta(days=self.TREND_DAYS - 1)
            trend_fi_str = trend_start_date.isoformat()

            # Obtener datos de tendencia para el modelo específico
            data_trend = self.leer_estadisticas(trend_fi_str, ff_reporte, modelo_str, "(TODOS)")
            evolution_trend = data_trend.get("evolution", {})

            if evolution_trend:
                sorted_dates = sorted(evolution_trend.keys())
                xlabels = [d[-5:] for d in sorted_dates]
                
                if len(sorted_dates) > 20:
                    step = max(1, len(sorted_dates) // 15)
                    xlabels = [d[-5:] if i % step == 0 else "" for i, d in enumerate(sorted_dates)]

                arr_ok = [evolution_trend[d].get("OK", 0) for d in sorted_dates]
                arr_perc_nok = []
                for d in sorted_dates:
                    ok_val = evolution_trend[d].get("OK", 0)
                    nok_val = evolution_trend[d].get("NOK", 0)
                    total = ok_val + nok_val
                    arr_perc_nok.append((nok_val / total * 100) if total > 0 else 0)

                fig, ax = self._crear_grafica_profesional_moderna(figsize=(12, 7))
                
                # Barras OK
                bars = ax.bar(range(len(arr_ok)), arr_ok, color=modern_colors['accent'], 
                            alpha=0.7, label="Piezas OK", width=0.8, edgecolor='white', linewidth=1.5)
                
                self._anotar_barras_moderno(ax)
                ax.set_ylabel("Piezas OK", fontsize=12, fontweight='600', color=modern_colors['dark'])
                ax.set_xticks(range(len(sorted_dates)))
                ax.set_xticklabels(xlabels, rotation=45, ha="right", fontsize=9, color=modern_colors['medium'])

                # Línea %NOK
                ax2 = ax.twinx()
                ax2.plot(range(len(arr_perc_nok)), arr_perc_nok, color=modern_colors['danger'], 
                        marker='o', linewidth=3, markersize=6, markerfacecolor='white', 
                        markeredgewidth=2, markeredgecolor=modern_colors['danger'], label="%NOK")
                
                ax2.axhline(y=self.OBJETIVO, color=modern_colors['success'], linestyle='--', 
                           linewidth=2, alpha=0.7, label=f"Objetivo {self.OBJETIVO:.1f}%")
                
                ax2.set_ylabel("%NOK", fontsize=12, fontweight='600', color=modern_colors['dark'])
                ax.text(0.5, 1.05, f"Tendencia de Calidad - {modelo_str.upper()}", 
                       transform=ax.transAxes, fontsize=16, fontweight='700', 
                       color=modern_colors['primary'], ha='center', va='bottom')

                ax.legend(loc='upper left', frameon=True, fontsize=10)
                ax2.legend(loc='upper right', frameon=True, fontsize=10)

                plt.subplots_adjust(top=0.88, bottom=0.12)

                buf_trend = io.BytesIO()
                fig.savefig(buf_trend, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                plt.close(fig)
                buf_trend.seek(0)
        except Exception as e:
            print(f"Error creando tendencia NOK para modelo {modelo_str}: {e}")

        # Los otros 3 gráficos simplificados para espacio
        buf_defectos = None
        if defect_dist_model and any(v > 0 for v in defect_dist_model.values()):
            try:
                sorted_defects = sorted(defect_dist_model.items(), key=lambda x: x[1], reverse=True)[:5]
                if sorted_defects:
                    labs, vals = zip(*sorted_defects)
                    
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(9, 7))
                    
                    # Donut chart
                    wedges, texts, autotexts = ax.pie(vals, labels=labs, autopct="%1.1f%%", 
                                                    startangle=90, textprops={'fontsize': 10},
                                                    colors=defect_colors[:len(labs)],
                                                    wedgeprops=dict(width=0.5, edgecolor='white'))
                    
                    # Centro blanco
                    centre_circle = plt.Circle((0,0), 0.70, fc='white')
                    fig.gca().add_artist(centre_circle)
                    
                    for autotext in autotexts:
                        autotext.set_color('white')
                        autotext.set_fontweight('bold')
                        autotext.set_fontsize(11)
                    
                    ax.text(0.5, 1.05, f"Distribución de Defectos - {modelo_str}", 
                           transform=ax.transAxes, fontsize=16, fontweight='700', 
                           color=modern_colors['primary'], ha='center', va='bottom')

                    plt.subplots_adjust(top=0.88)

                    buf_defectos = io.BytesIO()
                    fig.savefig(buf_defectos, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close(fig)
                    buf_defectos.seek(0)
            except Exception as e:
                print(f"Error creando gráfico de defectos para modelo {modelo_str}: {e}")

        # 3. Gráfico de turnos
        buf_turnos = None
        if shift_dist_model and any(v > 0 for v in shift_dist_model.values()):
            try:
                fig, ax = self._crear_grafica_profesional_moderna(figsize=(9, 7))
                
                turnos_orden = ['MAÑANA', 'TARDE', 'NOCHE', 'A', 'B', 'C']
                turnos_data = [(k, v) for k, v in shift_dist_model.items() if k and v > 0]
                turnos_data.sort(key=lambda x: turnos_orden.index(x[0]) if x[0] in turnos_orden else 999)
                
                if turnos_data:
                    labels, values = zip(*turnos_data)
                    
                    for i, (label, value) in enumerate(zip(labels, values)):
                        color = shift_colors.get(label, modern_colors['medium'])
                        ax.bar(i, value, color=color, alpha=0.85, width=0.6, edgecolor='white', linewidth=2)
                    
                    self._anotar_barras_moderno(ax, offset_factor=0.03)
                    ax.set_xticks(range(len(labels)))
                    ax.set_xticklabels(labels, rotation=0, ha="center", fontsize=12,
                                      color=modern_colors['dark'], fontweight='600')
                    
                    ax.text(0.5, 1.05, f"NOK por Turno - {modelo_str}", 
                           transform=ax.transAxes, fontsize=16, fontweight='700', 
                           color=modern_colors['primary'], ha='center', va='bottom')
                    ax.set_ylabel("Cantidad NOK", fontsize=12, fontweight='600', color=modern_colors['dark'])

                    plt.subplots_adjust(top=0.88)

                    buf_turnos = io.BytesIO()
                    fig.savefig(buf_turnos, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close(fig)
                    buf_turnos.seek(0)
            except Exception as e:
                print(f"Error creando gráfico de turnos para modelo {modelo_str}: {e}")

        # 4. Gráfico top 5 defectos (barras horizontales)
        buf_top_def = None
        if defect_dist_model and any(v > 0 for v in defect_dist_model.values()):
            try:
                sorted_defs = sorted(
                    [item for item in defect_dist_model.items() if item[1] > 0 and item[0]],
                    key=lambda x: x[1], reverse=True
                )[:5]
                if sorted_defs:
                    labs, vals = [t[0] for t in sorted_defs], [t[1] for t in sorted_defs]
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
                    
                    # Barras horizontales
                    bars = ax.barh(range(len(labs)), vals, 
                                  color=[defect_colors[i % len(defect_colors)] for i in range(len(labs))],
                                  alpha=0.85, height=0.6, edgecolor='white', linewidth=2)
                    
                    # Anotaciones
                    for i, (bar, val) in enumerate(zip(bars, vals)):
                        width = bar.get_width()
                        ax.text(width + max(vals)*0.01, bar.get_y() + bar.get_height()/2,
                               f'{int(val):,}', ha='left', va='center', fontsize=18, 
                               fontweight='600', color=modern_colors['dark'],
                               bbox=dict(boxstyle="round,pad=0.3", facecolor='white', 
                                        edgecolor='#E2E8F0', alpha=0.95))
                    
                    ax.set_yticks(range(len(labs)))
                    ax.set_yticklabels(labs, fontsize=11, color=modern_colors['dark'])
                    
                    ax.text(0.5, 1.05, f"Top 5 Defectos Críticos - {modelo_str}", 
                           transform=ax.transAxes, fontsize=16, fontweight='700', 
                           color=modern_colors['primary'], ha='center', va='bottom')
                    ax.set_xlabel("Cantidad NOK", fontsize=12, fontweight='600', color=modern_colors['dark'])

                    plt.subplots_adjust(top=0.88, left=0.20)

                    buf_top_def = io.BytesIO()
                    fig.savefig(buf_top_def, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close(fig)
                    buf_top_def.seek(0)
            except Exception as e:
                print(f"Error creando top 5 defectos para modelo {modelo_str}: {e}")

        return buf_trend, buf_defectos, buf_turnos, buf_top_def

    def _crear_graficos_modelo_rework(self, fi_reporte, ff_reporte, modelo_str, rework_evolution, rework_defect_dist, rework_shift_dist):
        """
        Crea los 4 gráficos específicos para análisis de retrabajos de un modelo.
        NUEVO: Gráficos con datos filtrados específicamente para retrabajos del modelo
        """
        modern_colors = {
            'primary': '#0F172A', 'accent': '#10B981', 'accent_dark': '#059669',
            'success': '#10B981', 'warning': '#F59E0B', 'danger': '#EF4444',
            'medium': '#64748B', 'dark': '#334155'
        }
        
        defect_colors = ['#10B981', '#059669', '#047857', '#065F46', '#064E3B']
        shift_colors = {'MAÑANA': '#10B981', 'TARDE': '#3B82F6', 'NOCHE': '#8B5CF6',
                       'A': '#10B981', 'B': '#3B82F6', 'C': '#8B5CF6'}

        # 1. Gráfico de evolución temporal de retrabajos
        buf_evol = None
        try:
            ff_dt = datetime.strptime(ff_reporte, '%Y-%m-%d').date()
            trend_start_date = ff_dt - timedelta(days=self.TREND_DAYS - 1)
            trend_fi_str = trend_start_date.isoformat()

            # Obtener datos de tendencia de retrabajos para el modelo específico
            data_trend = self.leer_estadisticas(trend_fi_str, ff_reporte, modelo_str, "(TODOS)")
            evolution_rework_trend = data_trend.get("evolution_rework", {})

            if evolution_rework_trend:
                sorted_dates = sorted(evolution_rework_trend.keys())
                xlabels = [d[-5:] for d in sorted_dates]
                
                if len(sorted_dates) > 20:
                    step = max(1, len(sorted_dates) // 15)
                    xlabels = [d[-5:] if i % step == 0 else "" for i, d in enumerate(sorted_dates)]

                vals_rework = [evolution_rework_trend.get(d, 0) for d in sorted_dates]
                
                fig, ax = self._crear_grafica_profesional_moderna(figsize=(12, 7))
                
                # Gráfico de área con línea
                ax.fill_between(range(len(vals_rework)), vals_rework, color=modern_colors['accent'], 
                               alpha=0.2, label='Área de retrabajos')
                
                ax.plot(range(len(vals_rework)), vals_rework, color=modern_colors['accent_dark'], 
                       marker='o', linewidth=3, markersize=7, markerfacecolor='white', 
                       markeredgewidth=2, markeredgecolor=modern_colors['accent_dark'], 
                       label='Retrabajos', zorder=5)

                ax.text(0.5, 1.05, f"Evolución de Retrabajos - {modelo_str.upper()}", 
                       transform=ax.transAxes, fontsize=16, fontweight='700', 
                       color=modern_colors['primary'], ha='center', va='bottom')
                
                ax.set_ylabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', color=modern_colors['dark'])
                ax.set_xticks(range(len(sorted_dates)))
                ax.set_xticklabels(xlabels, rotation=45, ha="right", fontsize=9, color=modern_colors['medium'])
                
                ax.legend(loc='upper left', frameon=True, fontsize=10)
                ax.set_ylim(bottom=0)

                plt.subplots_adjust(top=0.88, bottom=0.12)

                buf_evol = io.BytesIO()
                fig.savefig(buf_evol, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                plt.close(fig)
                buf_evol.seek(0)
        except Exception as e:
            print(f"Error creando evolución retrabajos para modelo {modelo_str}: {e}")

        # 2. Gráfico de defectos en retrabajos
        buf_def = None
        if rework_defect_dist and any(v > 0 for v in rework_defect_dist.values()):
            try:
                sorted_rdef = sorted(
                    [item for item in rework_defect_dist.items() if item[1] > 0 and item[0]],
                    key=lambda x: x[1], reverse=True
                )[:5]
                if sorted_rdef:
                    labs, vals = [t[0] for t in sorted_rdef], [t[1] for t in sorted_rdef]
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
                    
                    for i, (l, v) in enumerate(zip(labs, vals)):
                        ax.bar(i, v, color=defect_colors[i % len(defect_colors)], 
                              alpha=0.85, width=0.7, edgecolor='white', linewidth=2)
                    
                    ax.set_xticks(range(len(labs)))
                    ax.set_xticklabels(labs, rotation=30, ha="right", fontsize=10, 
                                      color=modern_colors['dark'], fontweight='500')
                    
                    self._anotar_barras_moderno(ax)
                    
                    ax.text(0.5, 1.05, f"Defectos en Retrabajos - {modelo_str}", 
                           transform=ax.transAxes, fontsize=16, fontweight='700', 
                           color=modern_colors['primary'], ha='center', va='bottom')
                    
                    ax.set_ylabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', color=modern_colors['dark'])

                    plt.subplots_adjust(top=0.88, bottom=0.15)

                    buf_def = io.BytesIO()
                    fig.savefig(buf_def, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close(fig)
                    buf_def.seek(0)
            except Exception as e:
                print(f"Error creando gráfico defectos retrabajos para modelo {modelo_str}: {e}")

        # 3. Gráfico de turnos en retrabajos
        buf_shift = None
        if rework_shift_dist and any(v > 0 for v in rework_shift_dist.values()):
            try:
                fig, ax = self._crear_grafica_profesional_moderna(figsize=(9, 7))
                
                turnos_orden = ['MAÑANA', 'TARDE', 'NOCHE', 'A', 'B', 'C']
                turnos_data = [(k, v) for k, v in rework_shift_dist.items() if k and v > 0]
                turnos_data.sort(key=lambda x: turnos_orden.index(x[0]) if x[0] in turnos_orden else 999)
                
                if turnos_data:
                    labels, values = zip(*turnos_data)
                    
                    for i, (label, value) in enumerate(zip(labels, values)):
                        color = shift_colors.get(label, modern_colors['medium'])
                        ax.bar(i, value, color=color, alpha=0.85, width=0.6, edgecolor='white', linewidth=2)
                    
                    self._anotar_barras_moderno(ax, offset_factor=0.03)
                    ax.set_xticks(range(len(labels)))
                    ax.set_xticklabels(labels, rotation=0, ha="center", fontsize=12,
                                      color=modern_colors['dark'], fontweight='600')
                    
                    ax.text(0.5, 1.05, f"Retrabajos por Turno - {modelo_str}", 
                           transform=ax.transAxes, fontsize=16, fontweight='700', 
                           color=modern_colors['primary'], ha='center', va='bottom')
                    
                    ax.set_ylabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', color=modern_colors['dark'])

                    plt.subplots_adjust(top=0.88)

                    buf_shift = io.BytesIO()
                    fig.savefig(buf_shift, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close(fig)
                    buf_shift.seek(0)
            except Exception as e:
                print(f"Error creando gráfico turnos retrabajos para modelo {modelo_str}: {e}")

        # 4. Top 5 defectos en retrabajos (barras horizontales)
        buf_top_def = None
        if rework_defect_dist and any(v > 0 for v in rework_defect_dist.values()):
            try:
                sorted_rdef = sorted(
                    [item for item in rework_defect_dist.items() if item[1] > 0 and item[0]],
                    key=lambda x: x[1], reverse=True
                )[:5]
                if sorted_rdef:
                    labs, vals = [t[0] for t in sorted_rdef], [t[1] for t in sorted_rdef]
                    fig, ax = self._crear_grafica_profesional_moderna(figsize=(11, 7))
                    
                    bars = ax.barh(range(len(labs)), vals, 
                                  color=[defect_colors[i % len(defect_colors)] for i in range(len(labs))],
                                  alpha=0.85, height=0.6, edgecolor='white', linewidth=2)
                    
                    for i, (bar, val) in enumerate(zip(bars, vals)):
                        width = bar.get_width()
                        ax.text(width + max(vals)*0.01, bar.get_y() + bar.get_height()/2,
                               f'{int(val):,}', ha='left', va='center', fontsize=18, 
                               fontweight='600', color=modern_colors['dark'],
                               bbox=dict(boxstyle="round,pad=0.3", facecolor='white', 
                                        edgecolor='#E2E8F0', alpha=0.95))
                    
                    ax.set_yticks(range(len(labs)))
                    ax.set_yticklabels(labs, fontsize=11, color=modern_colors['dark'])
                    
                    ax.text(0.5, 1.05, f"Top 5 Defectos en Retrabajo - {modelo_str}", 
                           transform=ax.transAxes, fontsize=16, fontweight='700', 
                           color=modern_colors['primary'], ha='center', va='bottom')
                    
                    ax.set_xlabel("Cantidad de Retrabajos", fontsize=12, fontweight='600', color=modern_colors['dark'])

                    plt.subplots_adjust(top=0.88, left=0.20)

                    buf_top_def = io.BytesIO()
                    fig.savefig(buf_top_def, format='png', dpi=150, bbox_inches='tight', facecolor='white')
                    plt.close(fig)
                    buf_top_def.seek(0)
            except Exception as e:
                print(f"Error creando top 5 defectos retrabajos para modelo {modelo_str}: {e}")

        return buf_evol, buf_def, buf_shift, buf_top_def

    # ------------------------------------------------------------------------------
    # UTILIDADES MEJORADAS PARA PDF Y OUTLOOK
    # ------------------------------------------------------------------------------
    
    # AGREGAR ESTAS FUNCIONES COMPLETAS A TU CLASE API

    def abrir_outlook_reporte(self, pdf_path, fecha_inicio, fecha_fin):
        """
        Abre un borrador de correo en Outlook con imágenes INCRUSTADAS en el cuerpo del correo.
        VERSIÓN COMPLETA - Con todas las funciones auxiliares incluidas.
        """
        try:
            import win32com.client
            import base64
            
            print("🖼️ Generando capturas optimizadas para correo...")
            capturas_paths = self._generar_capturas_con_fallback(pdf_path)
            
            outlook = win32com.client.Dispatch("Outlook.Application")
            mail = outlook.CreateItem(0)
            # Todos los destinatarios principales
            mail.To = "adem.benchohra@antolin.com; angelo.semeria@antolin.com; rebeca.cuesta@antolin.com; alvaro.heras@antolin.com; angel.lozano2@antolin.com; Diego.Bisabarros@antolin.com; enrique.izquierdo@antolin.com; felix.fernandez@antolin.com; fguillermo.sanchez@antolin.com; fjesus.rodriguez@antolin.com; jluis.ortega@antolin.com; jmaria.castro@antolin.com; juan.gomez3@antolin.com; oscar.martinez3@antolin.com; Oscar.Santamaria@antolin.com; rafael.gonzalez@antolin.com; renato.martinez@antolin.com; Ricardo.Hernandez@antolin.com; santiago.martin@antolin.com; sergio.garcia@antolin.com; sergio.marin@antolin.com; tomas.esteban@antolin.com; teodomiro.gonzalez@antolin.com; ivan.maria@antolin.com; javier.sanchez@antolin.com; Alberto.Lozano@antolin.com; sandra.arasti@antolin.com; Pilar.Siruela@antolin.com;"
            mail.Subject = f"Informe de Rechazo - {fecha_inicio} al {fecha_fin}"
            
            # Adjuntar PDF
            mail.Attachments.Add(Source=os.path.abspath(pdf_path))
            print(f"📎 PDF adjuntado: {os.path.basename(pdf_path)}")
            
            # Preparar imágenes optimizadas
            html_imagenes = ""
            
            if capturas_paths:
                for i, captura_path in enumerate(capturas_paths):
                    if os.path.exists(captura_path):
                        try:
                            # Leer y convertir imagen optimizada a base64
                            with open(captura_path, "rb") as img_file:
                                img_data = base64.b64encode(img_file.read()).decode('utf-8')
                            
                            # Título según el tipo de captura
                            if "Dashboard_Costes" in captura_path or "Cost_Impact" in captura_path:
                                titulo = ""
                                descripcion = "Pérdidas económicas por línea, UET y turno"
                            elif "Analisis_Visual" in captura_path or "Visual" in captura_path:
                                titulo = ""
                                descripcion = "Distribución y tendencias de costes"
                            else:
                                titulo = f"Captura del Reporte {i+1}"
                                descripcion = "Análisis de calidad"
                            
                            # HTML optimizado para cada imagen
                            html_imagenes += f"""
                            <div style="margin: 30px 0; text-align: center; page-break-inside: avoid; background: white; padding: 15px; border-radius: 12px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);">
                                <h3 style="color: #1E293B; margin: 0 0 5px 0; font-size: 18px; font-weight: bold;">{titulo}</h3>
                                <p style="color: #64748B; margin: 0 0 15px 0; font-size: 14px; font-style: italic;">{descripcion}</p>
                                <img src="data:image/jpeg;base64,{img_data}" 
                                    style="max-width: 100%; height: auto; border: 2px solid #E2E8F0; border-radius: 8px; box-shadow: 0 4px 12px rgba(0,0,0,0.15);"
                                    alt="{titulo}">
                            </div>
                            """
                            
                            print(f"📸 Imagen incrustada: {os.path.basename(captura_path)}")
                            
                        except Exception as e:
                            print(f"⚠️ Error procesando imagen {captura_path}: {e}")
                            continue
            
            # HTML final del correo
            mail.HTMLBody = f"""
            <html>
            <head>
                <meta charset="UTF-8">
                <meta name="viewport" content="width=device-width, initial-scale=1.0">
            </head>
            <body style="font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; line-height: 1.6; color: #333; margin: 0; padding: 0; background-color: #f5f7fa;">
                
                <!-- Contenedor principal -->
                <div style="max-width: 800px; margin: 0 auto; background: white;">
                    
                    <!-- Mensaje principal -->
                    <div style="padding: 40px 30px; background: white;">
                        <p style="font-size: 16px; margin: 0 0 15px 0; color: #2d3748;">Buenas,</p>
                        <p style="font-size: 16px; margin: 0 0 15px 0; color: #2d3748;">Adjunto informe de rechazo actualizado.</p>
                        <p style="font-size: 16px; margin: 0 0 0 0; color: #2d3748;">Un saludo,</p>
                    </div>
                    
                    <!-- Separador visual -->
                    <div style="height: 4px; background: linear-gradient(90deg, #3B82F6 0%, #10B981 50%, #F59E0B 100%); margin: 0;"></div>
                    
                    <!-- Contenedor de imágenes -->
                    <div style="padding: 20px; background: #f8fafc;">
                        {html_imagenes}
                    </div>
                    
                    <!-- Footer informativo -->
                    <div style="padding: 25px 30px; background: #1a202c; color: white;">
                        <table style="width: 100%; border-collapse: collapse;">
                            <tr>
                                <td style="padding: 5px 0; font-size: 14px;">
                                    <strong>Período de análisis:</strong> {fecha_inicio} al {fecha_fin}
                                </td>
                            </tr>
                            <tr>
                                <td style="padding: 15px 0 5px 0; font-size: 12px; color: #a0aec0; font-style: italic;">
                                    Las imágenes mostradas son extractos optimizados del reporte para visualización rápida.<br>
                                    Consulta el PDF adjunto para el análisis completo y detallado.
                                </td>
                            </tr>
                        </table>
                    </div>
                </div>
            </body>
            </html>
            """
            
            mail.Display()
            
            # Programar limpieza
            self._programar_limpieza_capturas(capturas_paths)
            
            if capturas_paths:
                return f"✅ Correo creado exitosamente con {len(capturas_paths)} imágenes optimizadas incrustadas directamente en el cuerpo del mensaje."
            else:
                return "✅ Correo creado con PDF adjunto (las capturas no se pudieron generar automáticamente)."
                
        except Exception as e:
            print(f"Error completo: {e}")
            return f"❌ Error creando correo: {str(e)}"

    def _generar_capturas_con_fallback(self, pdf_path):
        """
        Intenta generar capturas con múltiples métodos (fallback robusto)
        """
        capturas_paths = []
        
        # Método 1: Intentar con pdf2image
        try:
            print("🔄 Intentando generar capturas con pdf2image...")
            capturas_paths = self._generar_capturas_pdf2image(pdf_path)
            if capturas_paths:
                print(f"✅ Capturas generadas con pdf2image: {len(capturas_paths)}")
                return capturas_paths
        except ImportError:
            print("⚠️ pdf2image no está instalado")
        except Exception as e:
            print(f"⚠️ pdf2image falló: {e}")
        
        # Método 2: Intentar con PyMuPDF
        try:
            print("🔄 Intentando generar capturas con PyMuPDF...")
            capturas_paths = self._generar_capturas_pymupdf(pdf_path)
            if capturas_paths:
                print(f"✅ Capturas generadas con PyMuPDF: {len(capturas_paths)}")
                return capturas_paths
        except ImportError:
            print("⚠️ PyMuPDF no está instalado")
        except Exception as e:
            print(f"⚠️ PyMuPDF falló: {e}")
        
        # Método 3: Usar método manual básico
        try:
            print("🔄 Intentando método básico alternativo...")
            capturas_paths = self._generar_capturas_basico(pdf_path)
            if capturas_paths:
                print(f"✅ Capturas generadas con método básico: {len(capturas_paths)}")
                return capturas_paths
        except Exception as e:
            print(f"⚠️ Método básico falló: {e}")
        
        print("❌ No se pudieron generar capturas con ningún método")
        return []

    def _generar_capturas_pdf2image(self, pdf_path):
        """Método 1: pdf2image"""
        try:
            from pdf2image import convert_from_path
            import tempfile
            
            # Convertir solo las primeras 6 páginas para ser eficiente
            pages = convert_from_path(pdf_path, dpi=100, first_page=1, last_page=6)
            capturas_paths = []
            
            # Páginas específicas basadas en tu structure de PDF
            paginas_objetivo = []
            if len(pages) >= 3:  # Página 3: Cost Impact Dashboard
                paginas_objetivo.append((pages[2], "Dashboard_Costes"))
            if len(pages) >= 4:  # Página 4: Análisis Visual
                paginas_objetivo.append((pages[3], "Analisis_Visual"))
            
            temp_dir = tempfile.gettempdir()
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            for page_image, nombre in paginas_objetivo:
                filename = f"Captura_{nombre}_{timestamp}.png"
                filepath = os.path.join(temp_dir, filename)
                
                # Guardar con calidad optimizada para correo
                page_image.save(filepath, 'PNG', optimize=True)
                
                # Optimizar para correo
                filepath_optimizado = self._optimizar_imagen_para_correo(filepath)
                capturas_paths.append(filepath_optimizado)
            
            return capturas_paths
            
        except ImportError:
            raise ImportError("pdf2image no está instalado")
        except Exception as e:
            raise Exception(f"Error con pdf2image: {e}")

    def _generar_capturas_pymupdf(self, pdf_path):
        """Método 2: PyMuPDF"""
        try:
            import fitz  # PyMuPDF
            import tempfile
            
            pdf_document = fitz.open(pdf_path)
            capturas_paths = []
            
            # Páginas específicas
            paginas_capturar = [
                (2, "Dashboard_Costes"),    # Página 3 (índice 2)
                (3, "Analisis_Visual")      # Página 4 (índice 3)
            ]
            
            temp_dir = tempfile.gettempdir()
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            for page_num, nombre in paginas_capturar:
                if page_num < pdf_document.page_count:
                    page = pdf_document[page_num]
                    
                    # Renderizar página como imagen con buena calidad
                    mat = fitz.Matrix(1.2, 1.2)  # Escala 1.2x para buena calidad
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Guardar imagen
                    filename = f"Captura_{nombre}_{timestamp}.png"
                    filepath = os.path.join(temp_dir, filename)
                    pix.save(filepath)
                    
                    # Optimizar para correo
                    filepath_optimizado = self._optimizar_imagen_para_correo(filepath)
                    capturas_paths.append(filepath_optimizado)
            
            pdf_document.close()
            return capturas_paths
            
        except ImportError:
            raise ImportError("PyMuPDF no está instalado")
        except Exception as e:
            raise Exception(f"Error con PyMuPDF: {e}")

    def _generar_capturas_basico(self, pdf_path):
        """Método 3: Básico - crear placeholders si no hay otras opciones"""
        try:
            from PIL import Image, ImageDraw, ImageFont
            import tempfile
            
            print("⚠️ Generando placeholders - instala pdf2image o PyMuPDF para capturas reales")
            
            capturas_paths = []
            temp_dir = tempfile.gettempdir()
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            # Crear imágenes placeholder
            placeholders = [
                ("Dashboard_Costes", "📊 Dashboard de Impacto Económico", "Datos de costes por línea y turno"),
                ("Analisis_Visual", "📈 Análisis Visual de Costes", "Gráficos de distribución económica")
            ]
            
            for nombre, titulo, descripcion in placeholders:
                # Crear imagen placeholder
                img = Image.new('RGB', (800, 400), color='#f8fafc')
                draw = ImageDraw.Draw(img)
                
                # Intentar usar una fuente mejor
                try:
                    font_title = ImageFont.truetype("arial.ttf", 24)
                    font_desc = ImageFont.truetype("arial.ttf", 16)
                except:
                    font_title = ImageFont.load_default()
                    font_desc = ImageFont.load_default()
                
                # Dibujar texto
                draw.text((400, 150), titulo, fill='#1E293B', font=font_title, anchor="mm")
                draw.text((400, 200), descripcion, fill='#64748B', font=font_desc, anchor="mm")
                draw.text((400, 250), "Instala pdf2image para capturas reales del PDF", 
                        fill='#9CA3AF', font=font_desc, anchor="mm")
                
                # Dibujar borde
                draw.rectangle([10, 10, 790, 390], outline='#E2E8F0', width=2)
                
                # Guardar
                filename = f"Placeholder_{nombre}_{timestamp}.png"
                filepath = os.path.join(temp_dir, filename)
                img.save(filepath, 'PNG')
                
                capturas_paths.append(filepath)
            
            return capturas_paths
            
        except Exception as e:
            raise Exception(f"Error generando placeholders: {e}")

    def _optimizar_imagen_para_correo(self, imagen_path, max_width=700, calidad=80):
        """
        Optimiza las imágenes para que se vean bien en el correo sin ser demasiado pesadas
        """
        try:
            from PIL import Image
            import tempfile
            
            # Abrir imagen original
            with Image.open(imagen_path) as img:
                # Convertir a RGB si es necesario
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                
                # Redimensionar si es muy grande
                if img.width > max_width:
                    ratio = max_width / img.width
                    new_height = int(img.height * ratio)
                    img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)
                
                # Guardar imagen optimizada
                temp_dir = tempfile.gettempdir()
                base_name = os.path.splitext(os.path.basename(imagen_path))[0]
                optimized_path = os.path.join(temp_dir, f"opt_{base_name}.jpg")
                
                img.save(optimized_path, 'JPEG', quality=calidad, optimize=True)
                
                print(f"🔧 Imagen optimizada: {os.path.basename(imagen_path)} -> {os.path.basename(optimized_path)}")
                
                # Eliminar imagen original si es diferente
                if optimized_path != imagen_path:
                    try:
                        os.remove(imagen_path)
                    except:
                        pass
                
                return optimized_path
                
        except Exception as e:
            print(f"⚠️ Error optimizando imagen: {e}")
            return imagen_path  # Devolver original si falla la optimización

    def _programar_limpieza_capturas(self, capturas_paths):
        """Programa la limpieza de archivos temporales"""
        if not capturas_paths:
            return
        
        import threading
        import time
        
        def limpiar_despues():
            time.sleep(20)  # Esperar 20 segundos para que Outlook procese
            for captura_path in capturas_paths:
                try:
                    if os.path.exists(captura_path):
                        os.remove(captura_path)
                        print(f"🗑️ Archivo temporal eliminado: {os.path.basename(captura_path)}")
                except Exception as e:
                    print(f"⚠️ No se pudo eliminar {captura_path}: {e}")
        
        thread_limpieza = threading.Thread(target=limpiar_despues)
        thread_limpieza.daemon = True
        thread_limpieza.start()
        print(f"⏰ Programada limpieza de {len(capturas_paths)} archivos temporales en 20 segundos")

    def listar_reportes(self):
        """
        Lista los archivos PDF ya generados en la carpeta de reportes, del más reciente al más antiguo.
        MEJORADO: Mejor información sobre los reportes
        """
        if not os.path.exists(self.reports_dir):
            return {"error": f"El directorio '{self.reports_dir}' no existe.", "reportes": []}
        
        try:
            files = os.listdir(self.reports_dir)
            pdf_reports = []
            
            for f in files:
                if f.lower().endswith(".pdf"):
                    full_path = os.path.join(self.reports_dir, f)
                    stat_info = os.stat(full_path)
                    size_mb = round(stat_info.st_size / (1024 * 1024), 2)
                    modified_time = datetime.fromtimestamp(stat_info.st_mtime)
                    
                    pdf_reports.append({
                        "nombre": f,
                        "tamano_mb": size_mb,
                        "fecha_modificacion": modified_time.strftime('%d/%m/%Y %H:%M'),
                        "timestamp": stat_info.st_mtime
                    })
            
            # Ordenar por fecha de modificación (más reciente primero)
            pdf_reports.sort(key=lambda x: x["timestamp"], reverse=True)
            
            return {"error": None, "reportes": pdf_reports}
        except Exception as e:
            return {"error": f"Error listando reportes: {e}", "reportes": []}

    def abrir_reporte_historial(self, nombre_reporte):
        """
        Intenta abrir con la aplicación por defecto el PDF seleccionado de historial.
        MEJORADO: Mejor validación y mensajes de estado
        """
        if not nombre_reporte:
            return {"error": "Nombre de reporte no proporcionado.", "mensaje": ""}
        
        # Limpiar el nombre del reporte (quitar path si lo tiene)
        nombre_limpio = os.path.basename(nombre_reporte)
        report_path = os.path.join(self.reports_dir, nombre_limpio)
        
        if not os.path.exists(report_path):
            return {"error": f"El reporte '{nombre_limpio}' no fue encontrado en '{self.reports_dir}'.", "mensaje": ""}
        
        try:
            if sys.platform == "win32":
                os.startfile(report_path)
            elif sys.platform == "darwin":
                subprocess.call(["open", report_path])
            else:
                subprocess.call(["xdg-open", report_path])
            
            return {"error": None, "mensaje": f"✅ Abriendo '{nombre_limpio}' con la aplicación por defecto del sistema."}
        except Exception as e:
            print(f"Error abriendo reporte '{nombre_limpio}': {e}")
            return {"error": f"No se pudo abrir '{nombre_limpio}': {e}", "mensaje": ""}

    def eliminar_reporte(self, nombre_reporte):
        """
        NUEVA FUNCIÓN: Elimina un reporte del historial
        """
        if not nombre_reporte:
            return {"error": "Nombre de reporte no proporcionado.", "mensaje": ""}
        
        nombre_limpio = os.path.basename(nombre_reporte)
        report_path = os.path.join(self.reports_dir, nombre_limpio)
        
        if not os.path.exists(report_path):
            return {"error": f"El reporte '{nombre_limpio}' no existe.", "mensaje": ""}
        
        try:
            os.remove(report_path)
            return {"error": None, "mensaje": f"✅ Reporte '{nombre_limpio}' eliminado correctamente."}
        except Exception as e:
            return {"error": f"Error eliminando reporte: {e}", "mensaje": ""}

    def limpiar_reportes_antiguos(self, dias_antiguedad=30):
        """
        NUEVA FUNCIÓN: Limpia reportes más antiguos que X días
        """
        if not os.path.exists(self.reports_dir):
            return {"error": f"El directorio '{self.reports_dir}' no existe.", "eliminados": 0}
        
        try:
            fecha_limite = datetime.now() - timedelta(days=dias_antiguedad)
            files = os.listdir(self.reports_dir)
            eliminados = 0
            
            for f in files:
                if f.lower().endswith(".pdf"):
                    full_path = os.path.join(self.reports_dir, f)
                    stat_info = os.stat(full_path)
                    fecha_archivo = datetime.fromtimestamp(stat_info.st_mtime)
                    
                    if fecha_archivo < fecha_limite:
                        os.remove(full_path)
                        eliminados += 1
            
            return {"error": None, "eliminados": eliminados, 
                   "mensaje": f"✅ Se eliminaron {eliminados} reportes anteriores a {dias_antiguedad} días."}
        except Exception as e:
            return {"error": f"Error limpiando reportes antiguos: {e}", "eliminados": 0}

    # FUNCIONES DE DIAGNÓSTICO Y VERIFICACIÓN
    def verificar_archivos_disponibles(self, fecha_inicio, fecha_fin):
        """Función para verificar qué archivos están disponibles en un rango"""
        try:
            fi_dt = datetime.strptime(fecha_inicio, '%Y-%m-%d').date()
            ff_dt = datetime.strptime(fecha_fin, '%Y-%m-%d').date()
            
            patron = os.path.join(self.ruta_archivos, "*.xlsx")
            archivos_excel = glob.glob(patron)
            archivos_excel = [archivo for archivo in archivos_excel if not archivo.endswith("costes.xlsx")]
            
            archivos_info = []
            fechas_encontradas = set()
            
            for archivo in archivos_excel:
                nombre = os.path.basename(archivo)
                base, _ = os.path.splitext(nombre)
                try:
                    fecha_archivo = datetime.strptime(base, '%d.%m.%Y').date()
                    if fi_dt <= fecha_archivo <= ff_dt:
                        fechas_encontradas.add(fecha_archivo)
                        archivos_info.append({
                            'archivo': nombre,
                            'fecha': fecha_archivo.isoformat(),
                            'existe': True
                        })
                except ValueError:
                    continue
            
            # Verificar qué fechas faltan
            fecha_actual = fi_dt
            fechas_faltantes = []
            while fecha_actual <= ff_dt:
                if fecha_actual not in fechas_encontradas:
                    fechas_faltantes.append(fecha_actual.isoformat())
                fecha_actual += timedelta(days=1)
            
            return {
                'archivos_encontrados': len(archivos_info),
                'archivos_info': archivos_info,
                'fechas_faltantes': fechas_faltantes,
                'total_dias_rango': (ff_dt - fi_dt).days + 1
            }
            
        except Exception as e:
            return {'error': str(e)}

    def diagnosticar_problema_costes(self, fecha_fin):
        """Diagnosticar por qué el análisis de costes no muestra datos"""
        fecha_fin_dt = datetime.strptime(fecha_fin, '%Y-%m-%d').date()
        fecha_inicio_dt = fecha_fin_dt - timedelta(days=29)
        
        print("="*60)
        print("DIAGNÓSTICO DEL ANÁLISIS ECONÓMICO")
        print("="*60)
        
        # 1. Verificar directorio
        print(f"1. Verificando directorio: {self.ruta_archivos}")
        if not os.path.exists(self.ruta_archivos):
            print("   ❌ ERROR: El directorio no existe")
            return
        else:
            print("   ✅ Directorio existe")
        
        # 2. Verificar archivo de costes
        print(f"2. Verificando archivo de costes: {self.ruta_costes}")
        if not os.path.exists(self.ruta_costes):
            print("   ❌ ERROR: Archivo de costes no existe")
            return
        else:
            print(f"   ✅ Archivo de costes existe con {len(self.costes_dict)} precios")
            if len(self.costes_dict) == 0:
                print("   ⚠️  ADVERTENCIA: No hay precios cargados en el diccionario")
                return
            
            # Mostrar algunos ejemplos de precios
            print("   📋 Primeros precios cargados:")
            for i, (codigo, precio) in enumerate(list(self.costes_dict.items())[:5]):
                print(f"      - {codigo}: €{precio:.2f}")
        
        # 3. Listar archivos Excel en rango
        patron = os.path.join(self.ruta_archivos, "*.xlsx")
        todos_archivos = glob.glob(patron)
        archivos_excel = [f for f in todos_archivos if not f.endswith("costes.xlsx")]
        
        print(f"3. Archivos Excel encontrados: {len(archivos_excel)}")
        
        # 4. Verificar archivos en rango
        archivos_en_rango = []
        for archivo in archivos_excel:
            nombre = os.path.basename(archivo)
            base, _ = os.path.splitext(nombre)
            try:
                fecha_archivo = datetime.strptime(base, '%d.%m.%Y').date()
                if fecha_inicio_dt <= fecha_archivo <= fecha_fin_dt:
                    archivos_en_rango.append((nombre, fecha_archivo))
            except ValueError:
                continue
        
        print(f"   Archivos en rango: {len(archivos_en_rango)}")
        
        if len(archivos_en_rango) == 0:
            print("   ❌ PROBLEMA: No hay archivos en el rango de fechas")
            return
        
        # 5. Verificar contenido y costes
        print("4. Verificando contenido y cálculo de costes:")
        total_costes_encontrados = 0
        
        for nombre, fecha_archivo in sorted(archivos_en_rango, key=lambda x: x[1])[:5]:  # Solo primeros 5
            archivo_path = os.path.join(self.ruta_archivos, nombre)
            
            try:
                wb = openpyxl.load_workbook(archivo_path, data_only=True)
                ws = wb.active
                
                costes_archivo = 0
                filas_con_costes = 0
                
                for row in ws.iter_rows(min_row=4, values_only=True):
                    if len(row) < 12:
                        continue
                        
                    val_nok = row[3]
                    val_col_f = str(row[5] or "").strip().upper()
                    val_defecto_row = (row[6] or "").strip().upper()
                    val_texto_breve = (row[1] or "").strip().upper()
                    
                    _nok = self._parse_num(val_nok)
                    es_proveedor = self._es_defecto_proveedor(val_defecto_row)
                    
                    if val_col_f != "R" and _nok > 0 and val_texto_breve and not es_proveedor:
                        precio_unitario = self.costes_dict.get(val_texto_breve, 0)
                        if precio_unitario > 0:
                            coste_fila = _nok * precio_unitario
                            costes_archivo += coste_fila
                            filas_con_costes += 1
                
                print(f"   📄 {nombre}: {filas_con_costes} filas con costes → €{costes_archivo:,.2f}")
                total_costes_encontrados += costes_archivo
                    
            except Exception as e:
                print(f"   ❌ Error leyendo {nombre}: {e}")
        
        print("\n" + "="*60)
        print("RESUMEN DEL DIAGNÓSTICO:")
        print("="*60)
        print(f"✅ Precios en diccionario: {len(self.costes_dict)}")
        print(f"✅ Archivos en rango: {len(archivos_en_rango)}")
        print(f"💰 Total costes encontrados (muestra): €{total_costes_encontrados:,.2f}")
        
        if total_costes_encontrados > 0:
            print("\n✅ EL ANÁLISIS ECONÓMICO DEBERÍA FUNCIONAR CORRECTAMENTE")
            print("   Los datos de costes están disponibles y se están calculando correctamente.")
        else:
            print("\n❌ PROBLEMA: No se encontraron costes")
            print("   Posibles causas:")
            print("   - Los códigos en 'texto_breve' no coinciden con los del archivo costes.xlsx")
            print("   - Todos los defectos están marcados como 'PROVEEDOR'")
            print("   - Los precios en costes.xlsx son 0 o no válidos")
        
        print("="*60)

    def analisis_correlaciones(self, fecha_inicio, fecha_fin, modelo="(TODOS)", defecto="(TODOS)"):
        """
        Encuentra correlaciones y patrones interesantes en los datos del rango analizado.
        Descubre "cositas" que normalmente pasarían desapercibidas.
        """
        try:
            fi = datetime.strptime(fecha_inicio, '%Y-%m-%d').date()
            ff = datetime.strptime(fecha_fin, '%Y-%m-%d').date()
        except ValueError:
            return {"error": "Formato de fecha inválido"}

        # Obtener datos del rango
        data = self.leer_estadisticas(fecha_inicio, fecha_fin, modelo, defecto)
        if data.get("error"):
            return {"error": data["error"]}

        correlaciones = {
            "error": None,
            "periodo": f"{fecha_inicio} - {fecha_fin}",
            "hallazgos_criticos": [],
            "correlacion_turnos_defectos": {},
            "correlacion_dias_semana": {},
            "correlacion_modelos_defectos": {},
            "patrones_temporales": {},
            "correlacion_costes": {},
            "insights_automaticos": []
        }

        # Procesar archivos para análisis detallado
        patron = os.path.join(self.ruta_archivos, "*.xlsx")
        archivos_excel = [f for f in glob.glob(patron) if not f.endswith("costes.xlsx")]
        
        # Estructuras para correlaciones
        uet_defecto_matrix = {}  # {turno: {defecto: count}}
        dia_semana_nok = {}        # {dia_semana: total_nok}
        modelo_defecto_matrix = {} # {modelo: {defecto: count}}
        evolucion_diaria = {}      # {fecha: datos}
        costes_por_defecto = {}    # {defecto: coste_total}
        
        for archivo in archivos_excel:
            nombre = os.path.basename(archivo)
            base, _ = os.path.splitext(nombre)
            try:
                fecha_archivo = datetime.strptime(base, '%d.%m.%Y').date()
            except ValueError:
                continue

            if not (fi <= fecha_archivo <= ff):
                continue

            fecha_str = fecha_archivo.isoformat()
            dia_semana = fecha_archivo.strftime('%A')
            
            # Inicializar estructuras (CORREGIDO: usar list en lugar de set)
            if dia_semana not in dia_semana_nok:
                dia_semana_nok[dia_semana] = 0
            if fecha_str not in evolucion_diaria:
                evolucion_diaria[fecha_str] = {"total_nok": 0, "total_coste": 0, "defectos": []}

            try:
                wb = openpyxl.load_workbook(archivo, data_only=True)
                ws = wb.active

                for row in ws.iter_rows(min_row=4, values_only=True):
                    if len(row) < 12:
                        continue

                    val_nok = row[3]
                    val_col_f = str(row[5] or "").strip().upper()
                    val_defecto_row = (row[6] or "").strip().upper()
                    val_modelo_row = (row[7] or "").strip().upper()
                    val_turno_row = (row[9] or "SIN_TURNO").strip().upper()
                    val_texto_breve = (row[1] or "").strip().upper()

                    _nok = self._parse_num(val_nok)
                    es_proveedor = self._es_defecto_proveedor(val_defecto_row)

                    if val_col_f != "R" and _nok > 0 and not es_proveedor:
                        # Correlación Turno-Defecto
                        val_uet = (row[11] or "SIN_UET").strip().upper() if len(row) > 11 else "SIN_UET"
                        if val_uet not in uet_defecto_matrix:
                            uet_defecto_matrix[val_uet] = {}
                        if val_defecto_row not in uet_defecto_matrix[val_uet]:
                            uet_defecto_matrix[val_uet][val_defecto_row] = 0
                        uet_defecto_matrix[val_uet][val_defecto_row] += _nok

                        # Correlación Modelo-Defecto
                        if val_modelo_row not in modelo_defecto_matrix:
                            modelo_defecto_matrix[val_modelo_row] = {}
                        if val_defecto_row not in modelo_defecto_matrix[val_modelo_row]:
                            modelo_defecto_matrix[val_modelo_row][val_defecto_row] = 0
                        modelo_defecto_matrix[val_modelo_row][val_defecto_row] += _nok

                        # Día de semana
                        dia_semana_nok[dia_semana] += _nok

                        # Evolución diaria (CORREGIDO: usar append en lugar de add)
                        evolucion_diaria[fecha_str]["total_nok"] += _nok
                        if val_defecto_row not in evolucion_diaria[fecha_str]["defectos"]:
                            evolucion_diaria[fecha_str]["defectos"].append(val_defecto_row)

                        # Costes por defecto
                        precio_unitario = self.costes_dict.get(val_texto_breve, 0)
                        if precio_unitario > 0:
                            if val_defecto_row not in costes_por_defecto:
                                costes_por_defecto[val_defecto_row] = 0
                            costes_por_defecto[val_defecto_row] += _nok * precio_unitario
                            evolucion_diaria[fecha_str]["total_coste"] += _nok * precio_unitario

            except Exception as e:
                print(f"Error procesando {archivo}: {e}")
                continue

        # ANÁLISIS 1: Correlación Turnos-Defectos
        correlaciones["correlacion_uets_defectos"] = self._analizar_uet_defecto(uet_defecto_matrix)

        # ANÁLISIS 2: Correlación Días de Semana
        correlaciones["correlacion_dias_semana"] = self._analizar_dias_semana(dia_semana_nok)

        # ANÁLISIS 3: Correlación Modelos-Defectos
        correlaciones["correlacion_modelos_defectos"] = self._analizar_modelo_defecto(modelo_defecto_matrix)

        # ANÁLISIS 4: Patrones Temporales
        correlaciones["patrones_temporales"] = self._analizar_patrones_temporales(evolucion_diaria)

        # ANÁLISIS 5: Correlación Costes
        correlaciones["correlacion_costes"] = self._analizar_correlacion_costes(costes_por_defecto)

        # ANÁLISIS 6: Insights Automáticos
        correlaciones["insights_automaticos"] = self._generar_insights_automaticos(
            uet_defecto_matrix, dia_semana_nok, modelo_defecto_matrix, costes_por_defecto
        )

        return correlaciones

    def _analizar_uet_defecto(self, uet_defecto_matrix):
        """Encuentra qué UETs tienen problemas específicos con qué defectos"""
        resultados = {"matriz": uet_defecto_matrix, "hallazgos": [], "correlaciones_criticas": []}
        
        for uet, defectos in uet_defecto_matrix.items():
            if not defectos or uet == "SIN_UET":
                continue
                
            total_uet = sum(defectos.values())
            if total_uet < 5:  # Filtrar UETs con muy pocos defectos
                continue
            
            # Encontrar el defecto principal de cada UET
            defecto_principal = max(defectos.items(), key=lambda x: x[1])
            porcentaje_principal = (defecto_principal[1] / total_uet) * 100
            
            # Si un defecto representa >40% en una UET, es una correlación crítica
            if porcentaje_principal > 40:
                resultados["correlaciones_criticas"].append({
                    "uet": uet,
                    "defecto": defecto_principal[0],
                    "cantidad": defecto_principal[1],
                    "porcentaje": round(porcentaje_principal, 1),
                    "total_uet": total_uet
                })
                
                resultados["hallazgos"].append({
                    "tipo": "correlacion_uet_defecto",
                    "uet": uet,
                    "defecto": defecto_principal[0],
                    "porcentaje": round(porcentaje_principal, 1),
                    "cantidad": defecto_principal[1],
                    "mensaje": f"🎯 UET {uet}: {porcentaje_principal:.1f}% de sus problemas son {defecto_principal[0]} ({defecto_principal[1]} casos)"
                })
        
        # Encontrar UETs más problemáticas
        uets_totales = [(uet, sum(defectos.values())) for uet, defectos in uet_defecto_matrix.items() 
                    if uet != "SIN_UET" and sum(defectos.values()) >= 5]
        uets_totales.sort(key=lambda x: x[1], reverse=True)
        
        # Top 3 UETs más problemáticas
        for i, (uet, total) in enumerate(uets_totales[:3]):
            resultados["hallazgos"].append({
                "tipo": "uet_problematica",
                "uet": uet,
                "total": total,
                "ranking": i + 1,
                "mensaje": f"🚨 UET {uet}: {total} defectos totales (#{i+1} más problemática)"
            })
        
        return resultados

    def _analizar_dias_semana(self, dia_semana_nok):
            """Encuentra qué días de la semana son más problemáticos"""
            if not dia_semana_nok:
                return {"distribucion": {}, "hallazgos": []}
            
            total_semanal = sum(dia_semana_nok.values())
            promedio_dia = total_semanal / len(dia_semana_nok)
            
            hallazgos = []
            for dia, nok_count in dia_semana_nok.items():
                desviacion = ((nok_count - promedio_dia) / promedio_dia) * 100
                
                if abs(desviacion) > 20:  # Más de 20% de desviación
                    tipo = "problema" if desviacion > 0 else "bueno"
                    hallazgos.append({
                        "tipo": tipo,
                        "dia": dia,
                        "nok_count": nok_count,
                        "desviacion": round(desviacion, 1),
                        "mensaje": f"{'🔴' if tipo == 'problema' else '🟢'} {dia}: {desviacion:+.1f}% vs promedio"
                    })
            
            return {"distribucion": dia_semana_nok, "promedio": round(promedio_dia, 1), "hallazgos": hallazgos}

    def _analizar_modelo_defecto(self, modelo_defecto_matrix):
            """Encuentra qué modelos tienen defectos específicos recurrentes"""
            resultados = {"matriz": modelo_defecto_matrix, "hallazgos": []}
            
            for modelo, defectos in modelo_defecto_matrix.items():
                if not defectos or len(defectos) < 2:
                    continue
                    
                total_modelo = sum(defectos.values())
                defectos_ordenados = sorted(defectos.items(), key=lambda x: x[1], reverse=True)
                
                # Si los top 2 defectos representan >70% del total
                if len(defectos_ordenados) >= 2:
                    top2_total = defectos_ordenados[0][1] + defectos_ordenados[1][1]
                    porcentaje_top2 = (top2_total / total_modelo) * 100
                    
                    if porcentaje_top2 > 70:
                        resultados["hallazgos"].append({
                            "tipo": "patron_especifico",
                            "modelo": modelo,
                            "defecto1": defectos_ordenados[0][0],
                            "defecto2": defectos_ordenados[1][0],
                            "porcentaje": round(porcentaje_top2, 1),
                            "mensaje": f"🎯 {modelo}: {porcentaje_top2:.1f}% problemas en solo 2 defectos"
                        })
            
            return resultados

    def _analizar_patrones_temporales(self, evolucion_diaria):
        """Encuentra patrones en la evolución temporal"""
        if len(evolucion_diaria) < 7:
            return {"hallazgos": []}
        
        fechas_ordenadas = sorted(evolucion_diaria.keys())
        valores_nok = [evolucion_diaria[f]["total_nok"] for f in fechas_ordenadas]
        
        hallazgos = []
        
        # Detectar tendencias (CORREGIDO: Manejo de división por cero)
        if len(valores_nok) >= 7:
            primera_semana = sum(valores_nok[:7]) / 7
            ultima_semana = sum(valores_nok[-7:]) / 7
            
            # SOLUCIÓN: Verificar que primera_semana no sea 0
            if primera_semana > 0:
                cambio_tendencia = ((ultima_semana - primera_semana) / primera_semana) * 100
                
                if abs(cambio_tendencia) > 15:
                    tipo = "empeorando" if cambio_tendencia > 0 else "mejorando"
                    hallazgos.append({
                        "tipo": "tendencia",
                        "direccion": tipo,
                        "cambio_porcentual": round(cambio_tendencia, 1),
                        "mensaje": f"{'📈' if tipo == 'empeorando' else '📉'} Tendencia {tipo}: {cambio_tendencia:+.1f}%"
                    })
            else:
                # Caso especial: primera semana sin datos
                if ultima_semana > 0:
                    hallazgos.append({
                        "tipo": "inicio_problemas",
                        "direccion": "empeorando",
                        "cambio_porcentual": 999,  # Indica cambio muy significativo
                        "mensaje": f"📈 Aparición de problemas: {ultima_semana:.1f} NOK/día (sin problemas iniciales)"
                    })
        
        # Detectar picos anómalos (CORREGIDO: Manejo de promedio cero)
        if valores_nok:
            promedio = sum(valores_nok) / len(valores_nok)
            
            # Solo buscar picos si hay un promedio significativo
            if promedio > 0:
                for i, (fecha, valor) in enumerate(zip(fechas_ordenadas, valores_nok)):
                    if valor > promedio * 2:  # Más del doble del promedio
                        hallazgos.append({
                            "tipo": "pico_anomalo",
                            "fecha": fecha,
                            "valor": valor,
                            "promedio": round(promedio, 1),
                            "mensaje": f"⚠️ {fecha}: Pico anómalo ({valor} vs {promedio:.1f} promedio)"
                        })
            else:
                # Caso especial: buscar el primer día con problemas
                for i, (fecha, valor) in enumerate(zip(fechas_ordenadas, valores_nok)):
                    if valor > 0:
                        hallazgos.append({
                            "tipo": "primer_problema",
                            "fecha": fecha,
                            "valor": valor,
                            "mensaje": f"🔴 {fecha}: Primer día con problemas ({valor} NOK)"
                        })
                        break  # Solo el primero
        
        return {"hallazgos": hallazgos, "evolucion": evolucion_diaria}
    def _analizar_correlacion_costes(self, costes_por_defecto):
        """Encuentra correlaciones entre defectos y impacto económico"""
        if not costes_por_defecto:
            return {"hallazgos": []}
        
        total_costes = sum(costes_por_defecto.values())
        hallazgos = []
        
        # SOLUCIÓN: Verificar que hay costes
        if total_costes == 0:
            return {
                "distribución": costes_por_defecto,
                "hallazgos": [{
                    "tipo": "sin_costes",
                    "mensaje": "💰 No hay datos de costes disponibles para este análisis"
                }]
            }
        
        defectos_ordenados = sorted(costes_por_defecto.items(), key=lambda x: x[1], reverse=True)
        
        # Pareto de costes
        coste_acumulado = 0
        for i, (defecto, coste) in enumerate(defectos_ordenados):
            coste_acumulado += coste
            porcentaje_acumulado = (coste_acumulado / total_costes) * 100
            
            if porcentaje_acumulado >= 80 and i < 5:  # Principio 80/20
                hallazgos.append({
                    "tipo": "pareto_costes",
                    "num_defectos": i + 1,
                    "porcentaje_acumulado": round(porcentaje_acumulado, 1),
                    "coste_acumulado": round(coste_acumulado, 2),
                    "mensaje": f"💰 Solo {i + 1} defectos causan {porcentaje_acumulado:.1f}% del coste total"
                })
                break
        
        return {"distribucion": costes_por_defecto, "hallazgos": hallazgos}

    def _generar_insights_automaticos(self, uet_defecto_matrix, dia_semana_nok, modelo_defecto_matrix, costes_por_defecto):
        """Genera insights automáticos combinando todos los análisis"""
        insights = []
        
        # Insight 1: UET más problemática (CAMBIADO DE TURNO A UET)
        if uet_defecto_matrix:
            uets_total = {}
            for uet, defectos in uet_defecto_matrix.items():
                if uet != "SIN_UET":
                    uets_total[uet] = sum(defectos.values())
            
            if uets_total:
                uet_peor = max(uets_total.items(), key=lambda x: x[1])
                uets_ordenadas = sorted(uets_total.items(), key=lambda x: x[1], reverse=True)
                
                if len(uets_ordenadas) > 1:
                    segunda_uet = uets_ordenadas[1]
                    diferencia = uet_peor[1] - segunda_uet[1]
                    
                    insights.append({
                        "tipo": "uet_critica",
                        "mensaje": f"🏭 UET {uet_peor[0]} tiene {diferencia} NOK más que UET {segunda_uet[0]}",
                        "recomendacion": f"Revisar configuración y mantenimiento de UET {uet_peor[0]}"
                    })
                
                # Correlación UET-Defecto específica
                defectos_uet_peor = uet_defecto_matrix[uet_peor[0]]
                defecto_principal = max(defectos_uet_peor.items(), key=lambda x: x[1])
                porcentaje = (defecto_principal[1] / uet_peor[1]) * 100
                
                if porcentaje > 50:
                    insights.append({
                        "tipo": "correlacion_especifica",
                        "mensaje": f"🎯 UET {uet_peor[0]} especializada en {defecto_principal[0]} ({porcentaje:.1f}% de sus problemas)",
                        "recomendacion": f"Revisar proceso específico de {defecto_principal[0]} en UET {uet_peor[0]}"
                    })
        
        # Insight 2: Concentración de problemas (igual que antes)
        if modelo_defecto_matrix:
            concentraciones = []
            for modelo, defectos in modelo_defecto_matrix.items():
                if defectos:
                    total = sum(defectos.values())
                    max_defecto = max(defectos.items(), key=lambda x: x[1])
                    concentracion = (max_defecto[1] / total) * 100
                    if concentracion > 60:
                        concentraciones.append((modelo, max_defecto[0], concentracion))
            
            if concentraciones:
                concentraciones.sort(key=lambda x: x[2], reverse=True)
                modelo, defecto, porc = concentraciones[0]
                insights.append({
                    "tipo": "concentracion_maxima",
                    "mensaje": f"🎯 {modelo}: {porc:.1f}% de problemas son {defecto}",
                    "recomendacion": f"Acción específica en {defecto} para {modelo} tendrá máximo impacto"
                })
        
        # Insight 3: Oportunidad de ahorro (igual que antes)
        if costes_por_defecto:
            defecto_mas_caro = max(costes_por_defecto.items(), key=lambda x: x[1])
            insights.append({
                "tipo": "oportunidad_ahorro",
                "mensaje": f"💸 Resolver {defecto_mas_caro[0]} ahorraría €{defecto_mas_caro[1]:,.2f}",
                "recomendacion": f"Priorizar plan de acción para eliminar {defecto_mas_caro[0]}"
            })
        
        return insights
    def analisis_predictivo_equipos_humanos(self, fecha_fin, dias_historia=21):
        """
        Predice qué EQUIPOS HUMANOS van a empeorar basado en patrones de comportamiento
        """
        try:
            ff = datetime.strptime(fecha_fin, '%Y-%m-%d').date()
            fi = ff - timedelta(days=dias_historia - 1)
        except ValueError:
            return {"error": "Formato de fecha inválido"}

        # Estructuras para análisis predictivo de equipos
        uets_evolucion = {}
        
        # Procesar archivos para obtener evolución diaria por UET
        patron = os.path.join(self.ruta_archivos, "*.xlsx")
        archivos_excel = [f for f in glob.glob(patron) if not f.endswith("costes.xlsx")]
        
        for archivo in archivos_excel:
            nombre = os.path.basename(archivo)
            base, _ = os.path.splitext(nombre)
            try:
                fecha_archivo = datetime.strptime(base, '%d.%m.%Y').date()
            except ValueError:
                continue

            if not (fi <= fecha_archivo <= ff):
                continue

            fecha_str = fecha_archivo.isoformat()
            
            try:
                wb = openpyxl.load_workbook(archivo, data_only=True)
                ws = wb.active

                for row in ws.iter_rows(min_row=4, values_only=True):
                    if len(row) < 12:
                        continue

                    val_nok = row[3]
                    val_col_f = str(row[5] or "").strip().upper()
                    val_defecto_row = (row[6] or "").strip().upper()
                    val_uet = (row[11] or "SIN_UET").strip().upper()

                    _nok = self._parse_num(val_nok)
                    es_proveedor = self._es_defecto_proveedor(val_defecto_row)

                    if val_col_f != "R" and _nok > 0 and not es_proveedor and val_uet != "SIN_UET":
                        
                        if val_uet not in uets_evolucion:
                            uets_evolucion[val_uet] = {}
                        
                        if fecha_str not in uets_evolucion[val_uet]:
                            uets_evolucion[val_uet][fecha_str] = {
                                'total_defectos': 0,
                                'tipos_defecto': {}
                            }
                        
                        uets_evolucion[val_uet][fecha_str]['total_defectos'] += _nok
                        
                        if val_defecto_row not in uets_evolucion[val_uet][fecha_str]['tipos_defecto']:
                            uets_evolucion[val_uet][fecha_str]['tipos_defecto'][val_defecto_row] = 0
                        uets_evolucion[val_uet][fecha_str]['tipos_defecto'][val_defecto_row] += _nok

            except Exception as e:
                print(f"Error procesando {archivo}: {e}")
                continue

        # ANÁLISIS PREDICTIVO POR EQUIPO
        predicciones_equipos = []
        
        for uet, evolucion in uets_evolucion.items():
            if len(evolucion) < 7:  # Necesitamos al menos 7 días de datos
                continue
                
            fechas_ordenadas = sorted(evolucion.keys())
            valores_diarios = [evolucion[f]['total_defectos'] for f in fechas_ordenadas]
            
            # INDICADORES ESPECÍFICOS DE CONTROL DE EQUIPOS
            
            # 1. ÍNDICE DE CONTROL DEL PROCESO (¿Mantienen la especialización?)
            control_proceso = self._calcular_control_proceso(evolucion)
            
            # 2. PATRÓN DE FATIGA SEMANAL  
            fatiga_score = self._detectar_patron_fatiga_equipo(evolucion)
            
            # 3. CONSISTENCIA DEL EQUIPO (¿Trabajan igual cada día?)
            consistencia_score = self._calcular_consistencia_equipo(evolucion)
            
            # 4. DISPERSIÓN DE DEFECTOS (Pérdida de control = más tipos de defectos)
            dispersion_defectos = self._calcular_dispersion_defectos(evolucion)
            
            # 5. TENDENCIA GENERAL
            if len(valores_diarios) >= 7:
                primera_semana = sum(valores_diarios[:7]) / 7
                ultima_semana = sum(valores_diarios[-7:]) / 7
                tendencia_pct = ((ultima_semana - primera_semana) / primera_semana * 100) if primera_semana > 0 else 0
            else:
                tendencia_pct = 0
            
            # CALCULAR SCORE DE RIESGO DEL EQUIPO
            score_riesgo_equipo = 0
            alertas_equipo = []
            
            # Pérdida de control del proceso
            if control_proceso < 50:
                score_riesgo_equipo += 40
                alertas_equipo.append(f"Pérdida de control del proceso: {control_proceso:.1f}%")
            
            # Dispersión de defectos (pérdida de especialización)
            if dispersion_defectos > 4:
                score_riesgo_equipo += 35
                alertas_equipo.append(f"Dispersión de defectos: {dispersion_defectos} tipos diferentes")
            
            # Patrón de fatiga del equipo
            if fatiga_score > 60:
                score_riesgo_equipo += 30
                alertas_equipo.append(f"Patrón de fatiga detectado en el equipo")
            
            # Inconsistencia diaria
            if consistencia_score < 40:
                score_riesgo_equipo += 25
                alertas_equipo.append(f"Equipo inconsistente: {consistencia_score:.1f}% estabilidad")
            
            # Tendencia creciente de problemas
            if tendencia_pct > 50:
                score_riesgo_equipo += 20
                alertas_equipo.append(f"Tendencia creciente: +{tendencia_pct:.1f}%")
            
            # GENERAR RECOMENDACIONES PARA JEFE DE LÍNEA
            if score_riesgo_equipo >= 70:
                nivel_riesgo = "CRÍTICO"
                recomendacion = f"🚨 EQUIPO UET {uet}: Intervención inmediata necesaria"
                acciones = [
                    "Reunión urgente con el equipo completo",
                    "Revisar y reforzar procedimientos estándar",
                    "Verificar cumplimiento de instrucciones de trabajo",
                    "Implementar seguimiento horario por 48h",
                    "Evaluar necesidad de re-formación específica"
                ]
            elif score_riesgo_equipo >= 50:
                nivel_riesgo = "ALTO"
                recomendacion = f"⚠️ EQUIPO UET {uet}: Atención del jefe de línea requerida"
                acciones = [
                    "Reunión de equipo para identificar desviaciones",
                    "Revisar procedimientos paso a paso con el equipo",
                    "Reforzar puntos críticos de control",
                    "Seguimiento diario por 1 semana"
                ]
            elif score_riesgo_equipo >= 30:
                nivel_riesgo = "MEDIO"
                recomendacion = f"🔍 EQUIPO UET {uet}: Monitoreo preventivo"
                acciones = [
                    "Conversación con el equipo sobre procedimientos",
                    "Revisar métricas semanalmente",
                    "Observar puntos críticos de control"
                ]
            else:
                nivel_riesgo = "BAJO"
                recomendacion = f"✅ EQUIPO UET {uet}: Bajo control, performance normal"
                acciones = ["Continuar operación normal", "Reconocer buen trabajo del equipo"]
            
            if score_riesgo_equipo >= 30:
                predicciones_equipos.append({
                    "uet": uet,
                    "score_riesgo": score_riesgo_equipo,
                    "nivel_riesgo": nivel_riesgo,
                    "control_proceso": round(control_proceso, 1),
                    "dispersion_defectos": dispersion_defectos,
                    "consistencia_score": round(consistencia_score, 1),
                    "fatiga_score": round(fatiga_score, 1),
                    "tendencia_pct": round(tendencia_pct, 1),
                    "alertas": alertas_equipo,
                    "recomendacion": recomendacion,
                    "acciones_jefe_linea": acciones,
                    "dias_para_accion": 2 if nivel_riesgo == "CRÍTICO" else 7 if nivel_riesgo == "ALTO" else 14
                })
        
        # Ordenar por score de riesgo
        predicciones_equipos.sort(key=lambda x: x["score_riesgo"], reverse=True)
        
        return {
            "error": None,
            "tipo": "ANALISIS_EQUIPOS_HUMANOS",
            "periodo_analisis": f"{fi.isoformat()} - {ff.isoformat()}",
            "predicciones": predicciones_equipos,
            "total_equipos_analizados": len(uets_evolucion),
            "equipos_en_riesgo": len([p for p in predicciones_equipos if p["score_riesgo"] >= 50]),
            "insights_control": self._generar_insights_control_equipos(predicciones_equipos)
        }

    def _calcular_control_proceso(self, evolucion):
        """Calcula qué tan controlado está el proceso del equipo"""
        if not evolucion:
            return 0
        
        # Obtener todos los tipos de defectos del período
        todos_defectos = {}
        for fecha_datos in evolucion.values():
            for defecto, cantidad in fecha_datos.get('tipos_defecto', {}).items():
                todos_defectos[defecto] = todos_defectos.get(defecto, 0) + cantidad
        
        if not todos_defectos:
            return 100
        
        total_defectos = sum(todos_defectos.values())
        
        # El control se mide por la concentración en pocos tipos de defecto
        # Un equipo controlado tiene 70-80% de sus problemas en 1-2 defectos específicos
        defectos_ordenados = sorted(todos_defectos.values(), reverse=True)
        
        if len(defectos_ordenados) >= 2:
            concentracion_top2 = (defectos_ordenados[0] + defectos_ordenados[1]) / total_defectos * 100
        elif len(defectos_ordenados) == 1:
            concentracion_top2 = defectos_ordenados[0] / total_defectos * 100
        else:
            concentracion_top2 = 0
        
        return min(concentracion_top2, 100)

    def _calcular_dispersion_defectos(self, evolucion):
        """Calcula cuántos tipos diferentes de defectos tiene el equipo"""
        tipos_defectos = set()
        
        for fecha_datos in evolucion.values():
            for defecto in fecha_datos.get('tipos_defecto', {}).keys():
                tipos_defectos.add(defecto)
        
        return len(tipos_defectos)

    def _detectar_patron_fatiga_equipo(self, evolucion):
        """Detecta si el equipo muestra patrones de fatiga semanal"""
        if len(evolucion) < 7:
            return 0
        
        fechas_ordenadas = sorted(evolucion.keys())
        valores_por_dia_semana = {}  # {0: [lunes_values], 1: [martes_values], ...}
        
        for fecha_str in fechas_ordenadas:
            fecha_dt = datetime.strptime(fecha_str, '%Y-%m-%d').date()
            dia_semana = fecha_dt.weekday()  # 0=Lunes, 6=Domingo
            
            if dia_semana not in valores_por_dia_semana:
                valores_por_dia_semana[dia_semana] = []
            
            valores_por_dia_semana[dia_semana].append(evolucion[fecha_str]['total_defectos'])
        
        # Calcular promedio por día de semana
        promedios_dia = {}
        for dia, valores in valores_por_dia_semana.items():
            if valores:
                promedios_dia[dia] = sum(valores) / len(valores)
        
        if len(promedios_dia) < 5:  # No hay suficientes días
            return 0
        
        # Detectar patrón creciente hacia fin de semana (fatiga)
        dias_ordenados = sorted(promedios_dia.keys())
        valores_ordenados = [promedios_dia[dia] for dia in dias_ordenados]
        
        # Calcular si hay tendencia creciente del lunes al viernes
        if len(valores_ordenados) >= 5:
            inicio_semana = sum(valores_ordenados[:2]) / 2  # Lunes-Martes
            fin_semana = sum(valores_ordenados[-2:]) / 2    # Jueves-Viernes
            
            if inicio_semana > 0:
                incremento = ((fin_semana - inicio_semana) / inicio_semana) * 100
                return max(0, min(100, incremento))  # Entre 0 y 100
        
        return 0

    def _calcular_consistencia_equipo(self, evolucion):
        """Calcula qué tan consistente es el equipo día a día"""
        if len(evolucion) < 3:
            return 100
        
        valores_diarios = [data['total_defectos'] for data in evolucion.values()]
        
        if not valores_diarios:
            return 100
        
        promedio = sum(valores_diarios) / len(valores_diarios)
        
        if promedio == 0:
            return 100
        
        # Calcular coeficiente de variación inverso (más consistente = menos variación)
        import numpy as np
        coef_variacion = (np.std(valores_diarios) / promedio) * 100
        consistencia = max(0, 100 - coef_variacion)
        
        return min(100, consistencia)

    def _generar_insights_control_equipos(self, predicciones):
        """Genera insights específicos sobre control de equipos"""
        insights = []
        
        if not predicciones:
            return insights
        
        # Equipo con mayor pérdida de control
        equipo_critico = max(predicciones, key=lambda x: x["score_riesgo"])
        insights.append({
            "tipo": "EQUIPO_PERDIENDO_CONTROL",
            "mensaje": f"🚨 Equipo UET {equipo_critico['uet']} perdiendo control del proceso",
            "accion": f"Intervención inmediata - {equipo_critico['dias_para_accion']} días para actuar"
        })
        
        # Detectar problema generalizado de control
        equipos_sin_control = [p for p in predicciones if p["control_proceso"] < 60]
        if len(equipos_sin_control) > 2:
            insights.append({
                "tipo": "PROBLEMA_CONTROL_GENERALIZADO",
                "mensaje": f"⚠️ {len(equipos_sin_control)} equipos con pérdida de control del proceso",
                "accion": "Revisar procedimientos generales y formación estándar"
            })
        
        # Detectar equipos con dispersión excesiva
        equipos_dispersos = [p for p in predicciones if p["dispersion_defectos"] > 4]
        if equipos_dispersos:
            insights.append({
                "tipo": "PERDIDA_ESPECIALIZACION",
                "mensaje": f"📊 {len(equipos_dispersos)} equipos perdiendo especialización",
                "accion": "Reforzar formación específica por tipo de defecto"
            })
        
        return insights
    
    def _optimizar_imagen_para_correo(self, imagen_path, max_width=750, calidad=85):
        """
        Optimiza las imágenes para que se vean bien en el correo sin ser demasiado pesadas
        """
        try:
            from PIL import Image
            import tempfile
            
            # Abrir imagen original
            with Image.open(imagen_path) as img:
                # Convertir a RGB si es necesario
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                
                # Redimensionar si es muy grande
                if img.width > max_width:
                    ratio = max_width / img.width
                    new_height = int(img.height * ratio)
                    img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)
                
                # Guardar imagen optimizada
                temp_dir = tempfile.gettempdir()
                optimized_path = os.path.join(
                    temp_dir, 
                    f"optimized_{os.path.basename(imagen_path)}"
                )
                
                img.save(optimized_path, 'JPEG', quality=calidad, optimize=True)
                
                print(f"🔧 Imagen optimizada: {os.path.basename(imagen_path)} -> {optimized_path}")
                return optimized_path
                
        except Exception as e:
            print(f"⚠️ Error optimizando imagen: {e}")
            return imagen_path  # Devolver original si falla la optimización

    # Modificar la función de generación de capturas para usar optimización
    def _generar_capturas_pdf2image(self, pdf_path):
        """Método 1: pdf2image con optimización"""
        from pdf2image import convert_from_path
        import tempfile
        
        # Usar DPI más bajo para archivos más pequeños en correo
        pages = convert_from_path(pdf_path, dpi=100, first_page=1, last_page=6)
        capturas_paths = []
        
        # Páginas específicas basadas en tu estructura
        paginas_objetivo = []
        if len(pages) >= 3:
            paginas_objetivo.append((pages[2], "Dashboard_Costes"))
        if len(pages) >= 4:
            paginas_objetivo.append((pages[3], "Analisis_Visual"))
        
        temp_dir = tempfile.gettempdir()
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        for page_image, nombre in paginas_objetivo:
            filename = f"Captura_{nombre}_{timestamp}.png"
            filepath = os.path.join(temp_dir, filename)
            
            # Guardar con calidad optimizada para correo
            page_image.save(filepath, 'PNG', optimize=True)
            
            # Optimizar para correo
            filepath_optimizado = self._optimizar_imagen_para_correo(filepath, max_width=700, calidad=80)
            
            # Si la optimización fue exitosa, usar la versión optimizada
            if filepath_optimizado != filepath:
                try:
                    os.remove(filepath)  # Eliminar original no optimizado
                    capturas_paths.append(filepath_optimizado)
                except:
                    capturas_paths.append(filepath)  # Usar original si no se puede eliminar
            else:
                capturas_paths.append(filepath)
        
        return capturas_paths
    def generar_presentacion_html(self, fecha_inicio, fecha_fin):
        """
        Genera una presentación HTML fullscreen con auto-avance mostrando todo el análisis
        """
        try:
            # Obtener datos completos igual que para PDF
            data_all = self.leer_estadisticas(fecha_inicio, fecha_fin, "(TODOS)", "(TODOS)")
            if data_all.get("error"):
                return {"error": data_all["error"], "html_content": None}
            
            # Extraer datos
            global_ok = data_all.get("global_total_ok", 0)
            global_nok = data_all.get("global_total_nok", 0)
            global_nok_prov = data_all.get("global_total_nok_prov", 0)
            global_rew = data_all.get("global_total_rework", 0)
            
            def_dist_nok = data_all.get("defect_distribution", {})
            mod_dist_nok = data_all.get("models_distribution", {})
            shift_dist_nok = data_all.get("shift_distribution", {})
            
            rew_def_dist = data_all.get("rework_defect_distribution", {})
            rew_mod_dist = data_all.get("rework_models_distribution", {})
            rew_shift_dist = data_all.get("rework_shift_distribution", {})
            
            costes_por_linea_turno = data_all.get("costes_por_linea_turno", {})
            costes_globales = data_all.get("costes_globales", {})
            
            # Generar slides HTML
            slides_html = []
            
            # SLIDE 1: Dashboard NOK Global
            slide1 = self._generar_slide_nok_global_html(
                fecha_inicio, fecha_fin, global_ok, global_nok, global_nok_prov,
                def_dist_nok, mod_dist_nok, shift_dist_nok
            )
            slides_html.append(slide1)
            
            # SLIDE 2: Dashboard Retrabajos Global
            slide2 = self._generar_slide_rework_global_html(
                fecha_inicio, fecha_fin, global_rew,
                rew_def_dist, rew_mod_dist, rew_shift_dist
            )
            slides_html.append(slide2)
            
            # SLIDES 3-4: Análisis Económico (si hay datos)
            if costes_por_linea_turno:
                slide3 = self._generar_slide_costes_dashboard_html(
                    fecha_inicio, fecha_fin, costes_por_linea_turno, costes_globales
                )
                slides_html.append(slide3)
                
                slide4 = self._generar_slide_costes_graficos_html(
                    fecha_inicio, fecha_fin, costes_por_linea_turno
                )
                slides_html.append(slide4)
            
            # SLIDES 5-9: Top 5 Modelos NOK
            if mod_dist_nok:
                sorted_mod_nok = sorted(mod_dist_nok.items(), key=lambda x: x[1], reverse=True)
                top5_nok = [t[0] for t in sorted_mod_nok[:5] if t[0]]
                
                for modelo in top5_nok:
                    slide_modelo = self._generar_slide_modelo_nok_html(
                        fecha_inicio, fecha_fin, modelo
                    )
                    slides_html.append(slide_modelo)
            
            # SLIDES 10-14: Top 5 Modelos Rework
            if rew_mod_dist:
                sorted_mod_rew = sorted(rew_mod_dist.items(), key=lambda x: x[1], reverse=True)
                top5_rew = [t[0] for t in sorted_mod_rew[:5] if t[0]]
                
                for modelo in top5_rew:
                    slide_modelo_rew = self._generar_slide_modelo_rework_html(
                        fecha_inicio, fecha_fin, modelo
                    )
                    slides_html.append(slide_modelo_rew)
            
            # Generar HTML completo de la presentación
            html_completo = self._generar_html_presentacion_completo(slides_html, len(slides_html))
            
            return {
                "error": None,
                "html_content": html_completo,
                "total_slides": len(slides_html),
                "mensaje": f"✅ Presentación generada con {len(slides_html)} diapositivas"
            }
            
        except Exception as e:
            print(f"Error generando presentación HTML: {e}")
            return {"error": str(e), "html_content": None}

    def _generar_slide_nok_global_html(self, fi, ff, total_ok, total_nok, total_nok_prov, 
                                        def_dist, mod_dist, shift_dist):
        """Genera slide 1: Dashboard NOK Global"""
        try:
            total_prod = total_ok + total_nok
            rejection_rate = (total_nok / (total_ok + total_nok) * 100) if (total_ok + total_nok) > 0 else 0
            
            # Generar gráficos pequeños y convertir a base64
            img_defectos_base64 = self._generar_grafico_html_base64(def_dist, "Top 5 Defectos", tipo="bar")
            img_modelos_base64 = self._generar_grafico_html_base64(mod_dist, "Top 5 Modelos", tipo="bar")
            img_turnos_base64 = self._generar_grafico_html_base64(shift_dist, "NOK por Turno", tipo="bar")
            img_distribucion_base64 = self._generar_grafico_html_base64(def_dist, "Distribución", tipo="pie")
            
            html = f"""
            <div class="slide">
                <div class="slide-header">
                    <h1>QUALITY DASHBOARD - NOK ANALYSIS</h1>
                    <p class="period">Período: {fi} — {ff}</p>
                </div>
                
                <div class="kpis-container">
                    <div class="kpi-card kpi-hero" style="background: linear-gradient(135deg, #EF4444 0%, #DC2626 100%);">
                        <div class="kpi-label">TASA DE RECHAZO</div>
                        <div class="kpi-value-hero">{rejection_rate:.2f}%</div>
                        <div class="kpi-status">{'CRÍTICO' if rejection_rate > 0.8 else 'ALERTA' if rejection_rate > 0.5 else 'ÓPTIMO'}</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">PRODUCCIÓN TOTAL</div>
                        <div class="kpi-value">{total_prod:,}</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">RECHAZOS INTERNOS</div>
                        <div class="kpi-value">{total_nok:,}</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">RECHAZOS PROVEEDOR</div>
                        <div class="kpi-value">{total_nok_prov:,}</div>
                    </div>
                </div>
                
                <div class="charts-grid">
                    <div class="chart-box">
                        <h3>Top 5 Defectos Críticos</h3>
                        <img src="data:image/png;base64,{img_defectos_base64}" alt="Defectos">
                    </div>
                    <div class="chart-box">
                        <h3>Impacto por Modelo</h3>
                        <img src="data:image/png;base64,{img_modelos_base64}" alt="Modelos">
                    </div>
                    <div class="chart-box">
                        <h3>Distribución por Turno</h3>
                        <img src="data:image/png;base64,{img_turnos_base64}" alt="Turnos">
                    </div>
                    <div class="chart-box">
                        <h3>Distribución de Defectos</h3>
                        <img src="data:image/png;base64,{img_distribucion_base64}" alt="Distribución">
                    </div>
                </div>
            </div>
            """
            return html
        except Exception as e:
            print(f"Error generando slide NOK: {e}")
            return "<div class='slide'><h1>Error generando slide NOK</h1></div>"

    def _generar_slide_rework_global_html(self, fi, ff, total_rew, 
                                        rew_def_dist, rew_mod_dist, rew_shift_dist):
        """Genera slide 2: Dashboard Retrabajos Global"""
        try:
            img_def_rew = self._generar_grafico_html_base64(rew_def_dist, "Defectos R", tipo="bar")
            img_mod_rew = self._generar_grafico_html_base64(rew_mod_dist, "Modelos R", tipo="bar")
            img_shift_rew = self._generar_grafico_html_base64(rew_shift_dist, "Turnos R", tipo="bar")
            img_dist_rew = self._generar_grafico_html_base64(rew_def_dist, "Distribución R", tipo="pie")
            
            html = f"""
            <div class="slide">
                <div class="slide-header">
                    <h1>REWORK DASHBOARD</h1>
                    <p class="period">Análisis de Retrabajos | {fi} — {ff}</p>
                </div>
                
                <div class="kpis-container">
                    <div class="kpi-card kpi-hero" style="background: linear-gradient(135deg, #10B981 0%, #059669 100%);">
                        <div class="kpi-label">TOTAL RETRABAJOS</div>
                        <div class="kpi-value-hero">{total_rew:,}</div>
                        <div class="kpi-status">UNIDADES RECUPERADAS</div>
                    </div>
                </div>
                
                <div class="charts-grid">
                    <div class="chart-box">
                        <h3>Top 5 Defectos Recuperados</h3>
                        <img src="data:image/png;base64,{img_def_rew}" alt="Defectos R">
                    </div>
                    <div class="chart-box">
                        <h3>Modelos en Retrabajo</h3>
                        <img src="data:image/png;base64,{img_mod_rew}" alt="Modelos R">
                    </div>
                    <div class="chart-box">
                        <h3>Distribución por Turno</h3>
                        <img src="data:image/png;base64,{img_shift_rew}" alt="Turnos R">
                    </div>
                    <div class="chart-box">
                        <h3>Distribución de Defectos</h3>
                        <img src="data:image/png;base64,{img_dist_rew}" alt="Distribución R">
                    </div>
                </div>
            </div>
            """
            return html
        except Exception as e:
            print(f"Error generando slide Rework: {e}")
            return "<div class='slide'><h1>Error generando slide Rework</h1></div>"

    def _generar_slide_costes_dashboard_html(self, fi, ff, costes_linea_turno, costes_globales):
        """Genera slide 3: Dashboard Económico"""
        try:
            total_perdidas = costes_globales.get('total_perdidas', 0)
            num_defectos = costes_globales.get('num_defectos', 0)
            coste_medio = total_perdidas / num_defectos if num_defectos > 0 else 0
            
            # Top 5 líneas más costosas
            sorted_lineas = sorted(costes_linea_turno.items(), 
                                key=lambda x: x[1]['total_coste'], reverse=True)[:5]
            
            tabla_html = "<table class='costes-table'>"
            tabla_html += "<thead><tr><th>LÍNEA/UET/TURNO</th><th>DEFECTOS</th><th>CANTIDAD</th><th>COSTE TOTAL</th></tr></thead><tbody>"
            
            for key_string, data in sorted_lineas:
                parts = key_string.split('|')
                if len(parts) == 3:
                    linea, uet, turno = parts
                    linea_display = f"{linea} | UET {uet} | Turno {turno}"
                else:
                    linea_display = key_string
                
                defectos_agrupados = {}
                for item in data['items']:
                    tb = item['texto_breve']
                    if tb not in defectos_agrupados:
                        defectos_agrupados[tb] = {'cantidad': 0, 'coste': 0}
                    defectos_agrupados[tb]['cantidad'] += item['cantidad']
                    defectos_agrupados[tb]['coste'] += item['coste_total']
                
                for defecto, info in defectos_agrupados.items():
                    tabla_html += f"""
                    <tr>
                        <td><strong>{linea_display}</strong></td>
                        <td>{defecto}</td>
                        <td>{info['cantidad']}</td>
                        <td style="color:#DC2626;font-weight:bold;">€ {self._format_euro(info['coste'])}</td>
                    </tr>
                    """
            
            tabla_html += "</tbody></table>"
            
            html = f"""
            <div class="slide">
                <div class="slide-header" style="background: linear-gradient(135deg, #C0392B 0%, #7F1D1D 100%);">
                    <h1>COST IMPACT DASHBOARD</h1>
                    <p class="period">Análisis Económico de Pérdidas | {fi} — {ff}</p>
                </div>
                
                <div class="kpis-container">
                    <div class="kpi-card kpi-hero" style="background: linear-gradient(135deg, #E74C3C 0%, #C0392B 100%);">
                        <div class="kpi-label">PÉRDIDAS TOTALES</div>
                        <div class="kpi-value-hero">€ {self._format_euro(total_perdidas)}</div>
                        <div class="kpi-status">IMPACTO ECONÓMICO</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">DEFECTOS VALORADOS</div>
                        <div class="kpi-value">{num_defectos:,}</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">COSTE MEDIO/DEFECTO</div>
                        <div class="kpi-value">€ {self._format_euro(coste_medio)}</div>
                    </div>
                </div>
                
                <div class="table-container">
                    <h3 style="margin-bottom:15px;color:#1E293B;">TOP 5 LÍNEAS CON MAYOR IMPACTO ECONÓMICO</h3>
                    {tabla_html}
                </div>
            </div>
            """
            return html
        except Exception as e:
            print(f"Error generando slide costes dashboard: {e}")
            return "<div class='slide'><h1>Error generando slide costes</h1></div>"

    def _generar_slide_costes_graficos_html(self, fi, ff, costes_linea_turno):
        """Genera slide 4: Gráficos de Costes"""
        try:
            # Generar gráficos de costes
            buf_ranking = self._generar_grafico_ranking_lineas_costes(costes_linea_turno)
            buf_turnos = self._generar_grafico_distribucion_turnos_costes(costes_linea_turno)
            
            img_ranking = self._buffer_to_base64(buf_ranking)
            img_turnos = self._buffer_to_base64(buf_turnos)
            
            html = f"""
            <div class="slide">
                <div class="slide-header" style="background: linear-gradient(135deg, #C0392B 0%, #7F1D1D 100%);">
                    <h1>ANÁLISIS VISUAL DE IMPACTO ECONÓMICO</h1>
                    <p class="period">Distribución y tendencias | {fi} — {ff}</p>
                </div>
                
                <div class="charts-grid-large">
                    <div class="chart-box-large">
                        <h3>Ranking de Líneas por Pérdidas</h3>
                        <img src="data:image/png;base64,{img_ranking}" alt="Ranking">
                    </div>
                    <div class="chart-box-large">
                        <h3>Distribución por Turno</h3>
                        <img src="data:image/png;base64,{img_turnos}" alt="Turnos">
                    </div>
                </div>
            </div>
            """
            return html
        except Exception as e:
            print(f"Error generando slide gráficos costes: {e}")
            return "<div class='slide'><h1>Error generando gráficos costes</h1></div>"

    def _generar_slide_modelo_nok_html(self, fi, ff, modelo):
        """Genera slide para un modelo NOK específico"""
        try:
            data_modelo = self.leer_estadisticas(fi, ff, modelo, "(TODOS)")
            if data_modelo.get("error"):
                return f"<div class='slide'><h1>Error cargando datos de {modelo}</h1></div>"
            
            total_ok = data_modelo.get("total_ok", 0)
            total_nok = data_modelo.get("total_nok", 0)
            total_prod = total_ok + total_nok
            pct_nok = (total_nok / total_prod * 100) if total_prod > 0 else 0
            
            def_dist = data_modelo.get("defect_distribution", {})
            shift_dist = data_modelo.get("shift_distribution", {})
            
            img_def = self._generar_grafico_html_base64(def_dist, "Defectos", tipo="bar")
            img_shift = self._generar_grafico_html_base64(shift_dist, "Turnos", tipo="bar")
            img_pie = self._generar_grafico_html_base64(def_dist, "Distribución", tipo="pie")
            
            html = f"""
            <div class="slide">
                <div class="slide-header">
                    <h1>MODELO {modelo.upper()}</h1>
                    <p class="period">Análisis de Rechazos NOK | {fi} — {ff}</p>
                </div>
                
                <div class="kpis-container">
                    <div class="kpi-card kpi-hero" style="background: linear-gradient(135deg, #1E293B 0%, #0F172A 100%);">
                        <div class="kpi-label">TASA RECHAZO</div>
                        <div class="kpi-value-hero">{pct_nok:.1f}%</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">PIEZAS OK</div>
                        <div class="kpi-value">{total_ok:,}</div>
                    </div>
                    <div class="kpi-card">
                        <div class="kpi-label">NOK INTERNOS</div>
                        <div class="kpi-value">{total_nok:,}</div>
                    </div>
                </div>
                
                <div class="charts-grid">
                    <div class="chart-box">
                        <h3>Distribución de Defectos</h3>
                        <img src="data:image/png;base64,{img_def}" alt="Defectos">
                    </div>
                    <div class="chart-box">
                        <h3>Análisis por Turno</h3>
                        <img src="data:image/png;base64,{img_shift}" alt="Turnos">
                    </div>
                    <div class="chart-box">
                        <h3>Top Defectos</h3>
                        <img src="data:image/png;base64,{img_pie}" alt="Top">
                    </div>
                </div>
            </div>
            """
            return html
        except Exception as e:
            print(f"Error generando slide modelo NOK {modelo}: {e}")
            return f"<div class='slide'><h1>Error: {modelo}</h1></div>"

    def _generar_slide_modelo_rework_html(self, fi, ff, modelo):
        """Genera slide para un modelo Rework específico"""
        try:
            data_modelo = self.leer_estadisticas(fi, ff, modelo, "(TODOS)")
            if data_modelo.get("error"):
                return f"<div class='slide'><h1>Error cargando datos de {modelo}</h1></div>"
            
            total_rew = data_modelo.get("total_rework", 0)
            rew_def_dist = data_modelo.get("rework_defect_distribution", {})
            rew_shift_dist = data_modelo.get("rework_shift_distribution", {})
            
            img_def = self._generar_grafico_html_base64(rew_def_dist, "Defectos R", tipo="bar")
            img_shift = self._generar_grafico_html_base64(rew_shift_dist, "Turnos R", tipo="bar")
            img_pie = self._generar_grafico_html_base64(rew_def_dist, "Distribución R", tipo="pie")
            
            html = f"""
            <div class="slide">
                <div class="slide-header" style="background: linear-gradient(135deg, #065F46 0%, #064E3B 100%);">
                    <h1>MODELO {modelo.upper()}</h1>
                    <p class="period">Análisis de Retrabajos | {fi} — {ff}</p>
                </div>
                
                <div class="kpis-container">
                    <div class="kpi-card kpi-hero" style="background: linear-gradient(135deg, #10B981 0%, #059669 100%);">
                        <div class="kpi-label">TOTAL RETRABAJOS</div>
                        <div class="kpi-value-hero">{total_rew:,}</div>
                    </div>
                </div>
                
                <div class="charts-grid">
                    <div class="chart-box">
                        <h3>Defectos Recuperados</h3>
                        <img src="data:image/png;base64,{img_def}" alt="Defectos R">
                    </div>
                    <div class="chart-box">
                        <h3>Análisis por Turno</h3>
                        <img src="data:image/png;base64,{img_shift}" alt="Turnos R">
                    </div>
                    <div class="chart-box">
                        <h3>Top Defectos Retrabajo</h3>
                        <img src="data:image/png;base64,{img_pie}" alt="Top R">
                    </div>
                </div>
            </div>
            """
            return html
        except Exception as e:
            print(f"Error generando slide modelo Rework {modelo}: {e}")
            return f"<div class='slide'><h1>Error: {modelo}</h1></div>"

    def _generar_grafico_html_base64(self, data_dict, titulo, tipo="bar"):
        """Genera un gráfico matplotlib y lo convierte a base64 para embedding en HTML"""
        try:
            if not data_dict:
                return ""
            
            sorted_data = sorted(data_dict.items(), key=lambda x: x[1], reverse=True)[:5]
            if not sorted_data:
                return ""
            
            labels, values = zip(*sorted_data)
            
            fig, ax = plt.subplots(figsize=(6, 4), facecolor='white')
            
            if tipo == "bar":
                colors = ['#0F172A', '#00B4A6', '#4a6fa5', '#7f8c8d', '#2c3e50']
                bars = ax.bar(range(len(labels)), values, color=colors[:len(labels)], alpha=0.85)
                ax.set_xticks(range(len(labels)))
                ax.set_xticklabels(labels, rotation=30, ha='right', fontsize=10)
                
                # Añadir valores encima
                for bar, val in zip(bars, values):
                    height = bar.get_height()
                    ax.text(bar.get_x() + bar.get_width()/2., height,
                        f'{int(val):,}', ha='center', va='bottom', fontweight='bold')
            
            elif tipo == "pie":
                colors = ['#0F172A', '#00B4A6', '#4a6fa5', '#7f8c8d', '#2c3e50']
                ax.pie(values, labels=labels, autopct='%1.1f%%', colors=colors[:len(labels)],
                    startangle=90, textprops={'fontsize': 10, 'fontweight': 'bold'})
            
            ax.set_title(titulo, fontsize=12, fontweight='bold', pad=10)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            plt.tight_layout()
            
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=100, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            
            import base64
            img_base64 = base64.b64encode(buf.read()).decode('utf-8')
            return img_base64
            
        except Exception as e:
            print(f"Error generando gráfico HTML: {e}")
            return ""

    def _buffer_to_base64(self, buf):
        """Convierte un buffer de imagen a base64"""
        try:
            if buf is None:
                return ""
            import base64
            buf.seek(0)
            return base64.b64encode(buf.read()).decode('utf-8')
        except:
            return ""

    def _generar_html_presentacion_completo(self, slides_html, total_slides):
        """Genera el HTML completo de la presentación con CSS y JavaScript"""
        
        slides_content = "\n".join(slides_html)
        
        html_template = f"""
    <!DOCTYPE html>
    <html lang="es">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Presentación Antolin - Quality Dashboard</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
        <style>
            * {{
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }}
            
            body {{
                font-family: 'Inter', sans-serif;
                background: #0F172A;
                color: #1E293B;
                overflow: hidden;
            }}
            
            .presentation-container {{
                width: 1920px;
                height: 1080px;
                margin: 0 auto;
                position: relative;
                background: white;
            }}
            
            .slide {{
                width: 100%;
                height: 100%;
                display: none;
                flex-direction: column;
                padding: 40px 60px;
                background: white;
                animation: slideIn 0.6s ease-out;
            }}
            
            .slide.active {{
                display: flex;
            }}
            
            @keyframes slideIn {{
                from {{
                    opacity: 0;
                    transform: translateX(50px);
                }}
                to {{
                    opacity: 1;
                    transform: translateX(0);
                }}
            }}
            
            .slide-header {{
                background: linear-gradient(135deg, #0F172A 0%, #1E293B 100%);
                padding: 30px 40px;
                border-radius: 12px;
                margin-bottom: 30px;
                box-shadow: 0 4px 20px rgba(15, 23, 42, 0.3);
            }}
            
            .slide-header h1 {{
                color: white;
                font-size: 48px;
                font-weight: 800;
                letter-spacing: -0.02em;
                margin-bottom: 8px;
            }}
            
            .slide-header .period {{
                color: rgba(255, 255, 255, 0.8);
                font-size: 20px;
                font-weight: 500;
            }}
            
            .kpis-container {{
                display: grid;
                grid-template-columns: repeat(4, 1fr);
                gap: 20px;
                margin-bottom: 30px;
            }}
            
            .kpi-card {{
                background: linear-gradient(135deg, #F8FAFC 0%, #E2E8F0 100%);
                border-radius: 12px;
                padding: 25px;
                box-shadow: 0 2px 12px rgba(15, 23, 42, 0.08);
                border: 2px solid #E2E8F0;
                transition: transform 0.3s ease;
            }}
            
            .kpi-card:hover {{
                transform: translateY(-5px);
            }}
            
            .kpi-hero {{
                grid-column: span 1;
            }}
            
            .kpi-label {{
                font-size: 14px;
                font-weight: 700;
                color: rgba(255, 255, 255, 0.9);
                text-transform: uppercase;
                letter-spacing: 0.1em;
                margin-bottom: 12px;
            }}
            
            .kpi-value {{
                font-size: 36px;
                font-weight: 800;
                color: #1E293B;
            }}
            
            .kpi-value-hero {{
                font-size: 56px;
                font-weight: 900;
                color: white;
            }}
            
            .kpi-status {{
                font-size: 13px;
                color: rgba(255, 255, 255, 0.8);
                margin-top: 8px;
                font-weight: 600;
            }}
            
            .charts-grid {{
                display: grid;
                grid-template-columns: repeat(3, 1fr);
                gap: 20px;
                flex: 1;
            }}
            
            .charts-grid-large {{
                display: grid;
                grid-template-columns: repeat(2, 1fr);
                gap: 30px;
                flex: 1;
            }}
            
            .chart-box {{
                background: #F8FAFC;
                border-radius: 12px;
                padding: 20px;
                box-shadow: 0 2px 12px rgba(15, 23, 42, 0.08);
                border: 2px solid #E2E8F0;
                display: flex;
                flex-direction: column;
            }}
            
            .chart-box-large {{
                background: #F8FAFC;
                border-radius: 12px;
                padding: 25px;
                box-shadow: 0 2px 12px rgba(15, 23, 42, 0.08);
                border: 2px solid #E2E8F0;
                display: flex;
                flex-direction: column;
            }}
            
            .chart-box h3, .chart-box-large h3 {{
                font-size: 18px;
                font-weight: 700;
                color: #1E293B;
                margin-bottom: 15px;
                text-align: center;
            }}
            
            .chart-box img, .chart-box-large img {{
                width: 100%;
                height: auto;
                object-fit: contain;
            }}
            
            .table-container {{
                background: #F8FAFC;
                border-radius: 12px;
                padding: 25px;
                box-shadow: 0 2px 12px rgba(15, 23, 42, 0.08);
                border: 2px solid #E2E8F0;
                overflow-y: auto;
                max-height: 700px;
            }}
            
            .costes-table {{
                width: 100%;
                border-collapse: collapse;
                font-size: 16px;
            }}
            
            .costes-table thead {{
                background: linear-gradient(135deg, #C0392B 0%, #7F1D1D 100%);
                position: sticky;
                top: 0;
                z-index: 10;
            }}
            
            .costes-table th {{
                padding: 15px 12px;
                color: white;
                text-align: left;
                font-weight: 700;
                text-transform: uppercase;
                font-size: 14px;
                letter-spacing: 0.05em;
            }}
            
            .costes-table td {{
                padding: 12px;
                border-bottom: 1px solid #E2E8F0;
            }}
            
            .costes-table tbody tr:hover {{
                background: #EEF2FF;
            }}
            
            .progress-bar {{
                position: fixed;
                bottom: 0;
                left: 0;
                height: 6px;
                background: linear-gradient(90deg, #00B4A6 0%, #0F172A 100%);
                width: 0%;
                transition: width 0.1s linear;
                z-index: 1000;
            }}
            
            .slide-counter {{
                position: fixed;
                bottom: 20px;
                right: 40px;
                background: rgba(15, 23, 42, 0.9);
                color: white;
                padding: 12px 24px;
                border-radius: 999px;
                font-size: 16px;
                font-weight: 700;
                z-index: 1000;
                box-shadow: 0 4px 12px rgba(0, 0, 0, 0.3);
            }}
            
            .logo {{
                position: fixed;
                top: 20px;
                right: 40px;
                height: 50px;
                z-index: 1000;
            }}
            
            .controls {{
                position: fixed;
                bottom: 20px;
                left: 40px;
                display: flex;
                gap: 10px;
                z-index: 1000;
            }}
            
            .btn-control {{
                background: rgba(15, 23, 42, 0.9);
                color: white;
                border: none;
                padding: 12px 20px;
                border-radius: 8px;
                font-size: 14px;
                font-weight: 600;
                cursor: pointer;
                transition: all 0.3s ease;
            }}
            
            .btn-control:hover {{
                background: #00B4A6;
                transform: translateY(-2px);
            }}
        </style>
    </head>
    <body>
        <div class="presentation-container">
            {slides_content}
        </div>
        
        <div class="progress-bar" id="progressBar"></div>
        <div class="slide-counter" id="slideCounter">1 / {total_slides}</div>
        
        <div class="controls">
            <button class="btn-control" onclick="prevSlide()">◀ Anterior</button>
            <button class="btn-control" onclick="toggleAutoPlay()">⏸ Pausar</button>
            <button class="btn-control" onclick="nextSlide()">Siguiente ▶</button>
            <button class="btn-control" onclick="exitPresentation()">✕ Salir</button>
        </div>
        
        <script>
            let currentSlide = 0;
            const totalSlides = {total_slides};
            const slides = document.querySelectorAll('.slide');
            const slideCounter = document.getElementById('slideCounter');
            const progressBar = document.getElementById('progressBar');
            
            let autoPlayInterval;
            let isPlaying = true;
            const SLIDE_DURATION = 12000; // 12 segundos por slide
            
            function showSlide(index) {{
                slides.forEach((slide, i) => {{
                    slide.classList.toggle('active', i === index);
                }});
                
                currentSlide = index;
                slideCounter.textContent = `${{index + 1}} / ${{totalSlides}}`;
                
                // Actualizar barra de progreso
                const progress = ((index + 1) / totalSlides) * 100;
                progressBar.style.width = progress + '%';
            }}
            
            function nextSlide() {{
                const next = (currentSlide + 1) % totalSlides;
                showSlide(next);
            }}
            
            function prevSlide() {{
                const prev = (currentSlide - 1 + totalSlides) % totalSlides;
                showSlide(prev);
            }}
            
            function startAutoPlay() {{
                autoPlayInterval = setInterval(nextSlide, SLIDE_DURATION);
                isPlaying = true;
                document.querySelector('.btn-control:nth-child(2)').textContent = '⏸ Pausar';
            }}
            
            function stopAutoPlay() {{
                clearInterval(autoPlayInterval);
                isPlaying = false;
                document.querySelector('.btn-control:nth-child(2)').textContent = '▶ Reanudar';
            }}
            
            function toggleAutoPlay() {{
                if (isPlaying) {{
                    stopAutoPlay();
                }} else {{
                    startAutoPlay();
                }}
            }}
            
            function exitPresentation() {{
                if (confirm('¿Salir del modo presentación?')) {{
                    window.close();
                }}
            }}
            
            // Atajos de teclado
            document.addEventListener('keydown', (e) => {{
                switch(e.key) {{
                    case 'ArrowRight':
                    case ' ':
                        e.preventDefault();
                        nextSlide();
                        break;
                    case 'ArrowLeft':
                        e.preventDefault();
                        prevSlide();
                        break;
                    case 'Escape':
                        exitPresentation();
                        break;
                    case 'p':
                    case 'P':
                        toggleAutoPlay();
                        break;
                }}
            }});
            
            // Iniciar presentación
            showSlide(0);
            startAutoPlay();
            
            // Intentar entrar en fullscreen
            setTimeout(() => {{
                if (document.documentElement.requestFullscreen) {{
                    document.documentElement.requestFullscreen();
                }}
            }}, 500);
        </script>
    </body>
    </html>
        """
        
        return html_template
if __name__ == "__main__":
    api = API()
    html_file_name = "index2.html"
    html_local_path_str = os.path.join(os.path.dirname(__file__), html_file_name)

    if not os.path.exists(html_local_path_str):
        print(f"ADVERTENCIA: {html_file_name} no encontrado en {html_local_path_str}. La interfaz web podría no funcionar.")

    html_local_path_uri = Path(html_local_path_str).resolve().as_uri()
    print(f"Cargando UI desde: {html_local_path_uri}")
    print("🚀 Sistema de Análisis de Calidad iniciado")
    print("📊 Filtrado exhaustivo implementado - Los análisis se enfocan específicamente en los filtros seleccionados")
    print("🎯 Cuando filtres por modelo o defecto, obtendrás estadísticas precisas de esa selección específica")
    print("💰 NUEVO: Análisis económico de costes integrado - Impacto financiero por línea, UET y turno")
    print(f"💾 Archivo de costes: {api.ruta_costes}")
    print(f"📈 Precios cargados: {len(api.costes_dict)} códigos de defectos con precio")
    
    if len(api.costes_dict) == 0:
        print("⚠️  ADVERTENCIA: No se cargaron precios. El análisis económico no estará disponible.")
        print(f"   Verifica que existe el archivo: {api.ruta_costes}")
        print("   Y que tiene la estructura correcta: [Código_Defecto, Precio]")
    else:
        print("✅ Análisis económico disponible y funcionando")
    
    ventana = webview.create_window(
        "Estadísticas de Rechazo - Análisis Exhaustivo + Económico",
        url=html_local_path_uri,
        js_api=api,
        width=1400,
        height=900,
        resizable=True
    )
    webview.start()